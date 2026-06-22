#!/usr/bin/env python3
"""
make_family_slides.py — build the advisor-meeting deck for the multi-copy gene family
definition (the read-conflict graph). Generates a genome-wide validation figure and a
PowerPoint deck (process + examples/counterexamples + the airtight case).

Numbers are taken verbatim from bench/family_definition_formal.md / family_def_genomewide.py
(this session). Run: /home/juanfra/miniforge3/bin/python bench/make_family_slides.py
Output: bench/family_definition_slides.pptx + bench/family_def_genomewide_fig.png
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
FIG_CONFLICT = os.path.join(HERE, "family_definition_figure.png")  # (A) graph + (B) ledger
FIG_GW = os.path.join(HERE, "family_def_genomewide_fig.png")
OUT = os.path.join(HERE, "family_definition_slides.pptx")

NAVY = RGBColor(0x1F, 0x2D, 0x5A)
TEAL = RGBColor(0x12, 0x7A, 0x6E)
GREY = RGBColor(0x44, 0x44, 0x44)
ORANGE = RGBColor(0xD9, 0x6B, 0x27)

# ----------------------------------------------------------------- genome-wide figure
def make_gw_figure():
    delta = [0.003, 0.004, 0.005, 0.006, 0.007, 0.009, 0.017]
    fam_genes = [388, 403, 416, 427, 436, 460, 538]
    fam_dn = [195, 202, 212, 220, 227, 243, 279]
    sizes = {2: 238, 3: 73, 4: 27, 5: 15, 6: 18, 7: 7, 8: 3, 9: 7, 10: 1,
             11: 7, 12: 4, "13+": 16}

    fig, ax = plt.subplots(1, 3, figsize=(15, 4.3))

    # A: Δ-sweep — stable valley
    ax[0].plot(delta, fam_genes, "o-", color=GREY.__str__() and "#444", label="annotated genes (34k vertices)")
    ax[0].plot(delta, fam_dn, "s-", color="#127A6E", label="de-novo loci (101k, production)")
    ax[0].axvspan(0.003, 0.006, color="#cfe8e4", alpha=0.5, zorder=0)
    ax[0].axvline(0.005, color="#D96B27", ls="--", lw=1.5)
    ax[0].text(0.0052, ax[0].get_ylim()[1] * 0.55, "operating\npoint Δ=0.005", color="#D96B27", fontsize=8)
    ax[0].set_xlabel("Δ (de-tie tolerance)"); ax[0].set_ylabel("# families")
    ax[0].set_title("A. Δ is FLAT genome-wide\n(±7% over the band), not a panel fit", fontsize=11, weight="bold")
    ax[0].legend(fontsize=7, loc="upper left"); ax[0].grid(alpha=0.3)

    # B: bridge collapse at production loci
    cats = ["annotated\ngenes (coarse)", "de-novo loci\n(production)"]
    bridge_pct = [14, 9]; coherent_pct = [86, 91]
    x = range(len(cats)); w = 0.38
    ax[1].bar([i - w / 2 for i in x], coherent_pct, w, color="#127A6E", label="coherent (co-located)")
    ax[1].bar([i + w / 2 for i in x], bridge_pct, w, color="#D96B27", label="cross-chrom over-merge (coarse vertex)")
    for i, (c, b) in enumerate(zip(coherent_pct, bridge_pct)):
        ax[1].text(i - w / 2, c + 1, f"{c}%", ha="center", fontsize=9)
        ax[1].text(i + w / 2, b + 1, f"{b}%", ha="center", fontsize=9, color="#D96B27")
    ax[1].annotate("bridges 59→20\n(worst mega-bridges vanish)", xy=(1, 9), xytext=(0.25, 45),
                   fontsize=8, color="#D96B27", arrowprops=dict(arrowstyle="->", color="#D96B27"))
    ax[1].set_xticks(list(x)); ax[1].set_xticklabels(cats); ax[1].set_ylabel("% of families")
    ax[1].set_ylim(0, 100)
    ax[1].set_title("B. Coarse-vertex over-merge COLLAPSES\nat production loci", fontsize=11, weight="bold")
    ax[1].legend(fontsize=7, loc="upper right")

    # C: family-size distribution (size-2 dominated)
    keys = [str(k) for k in sizes]; vals = list(sizes.values())
    ax[2].bar(keys, vals, color="#1F2D5A")
    ax[2].set_xlabel("family size (# loci)"); ax[2].set_ylabel("# families")
    ax[2].set_title("C. Families are size-2 dominated\n(recent-paralog pairs)", fontsize=11, weight="bold")
    ax[2].text(0.5, 200, "57% are size 2", fontsize=9, color="#1F2D5A")

    fig.suptitle("Genome-wide on 34,114 INDEPENDENT gene vertices (no circularity): "
                 "416 families · Δ flat to ±7% · 86–91% co-located",
                 fontsize=11.5, weight="bold", y=1.02)
    fig.tight_layout()
    fig.savefig(FIG_GW, dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- pptx helpers
def add_title_slide(prs, title, subtitle, foot):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bx = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.0))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(20); p2.font.color.rgb = TEAL
    fb = s.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5))
    fp = fb.text_frame.paragraphs[0]; fp.text = foot; fp.font.size = Pt(12); fp.font.color.rgb = GREY
    return s


def add_content_slide(prs, title, bullets, note=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tp = tb.text_frame.paragraphs[0]; tp.text = title
    tp.font.size = Pt(26); tp.font.bold = True; tp.font.color.rgb = NAVY
    body = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.4))
    tf = body.text_frame; tf.word_wrap = True
    for i, (txt, lvl, color) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " if lvl == 0 else "      –  ") + txt
        p.font.size = Pt(18 if lvl == 0 else 15)
        p.font.color.rgb = color or RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(7)
    if note:
        nb = s.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.5))
        np_ = nb.text_frame.paragraphs[0]; np_.text = note
        np_.font.size = Pt(11); np_.font.italic = True; np_.font.color.rgb = GREY
    return s


def add_figure_slide(prs, title, img, caption):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.8))
    tp = tb.text_frame.paragraphs[0]; tp.text = title
    tp.font.size = Pt(24); tp.font.bold = True; tp.font.color.rgb = NAVY
    from PIL import Image
    w, h = Image.open(img).size
    maxw, maxh = Inches(12.3), Inches(5.3)
    scale = min(maxw / w, maxh / h)
    iw, ih = int(w * scale), int(h * scale)
    left = int((Inches(13.33) - iw) / 2)
    s.shapes.add_picture(img, left, Inches(1.15), width=iw, height=ih)
    cb = s.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.6))
    cp = cb.text_frame.paragraphs[0]; cp.text = caption
    cp.font.size = Pt(12); cp.font.color.rgb = GREY; cp.alignment = PP_ALIGN.CENTER
    return s


def B(t, lvl=0, color=None):  # bullet helper
    return (t, lvl, color)


# ----------------------------------------------------------------- build deck
def build():
    make_gw_figure()
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Defining a multi-copy gene family at the RNA level",
        "A read-conflict graph:  recent paralogy  ∩  read-confusability",
        "GGO HiFi IsoSeq · de-tie criterion · TP=7 TN=10 FP=0 FN=0 · genome-wide validated")

    add_content_slide(prs, "The question", [
        B("What is a 'gene family' — for the copy-assignment problem?"),
        B("The biological definition (sequence similarity / phylogeny) needs the genome or protein —", 1),
        B("not derivable from RNA, and circular if used to DETECT families from reads.", 1),
        B("The copy-assignment problem operates on a different unit:", 0),
        B("the set of loci among which the READS are genuinely confused (multimapping).", 1),
        B("Goal: a formal, RNA-derived definition with an honest false-positive / false-negative ledger.", 0, TEAL),
    ])

    add_content_slide(prs, "Key idea: multimapping reads are shared evidence, not noise", [
        B("A read that maps equally well to two loci = a genuine alternative placement"),
        B("(the multimapping conflict — Canzar's conflict graph).", 1),
        B("Two loci are in the same family  ⇔  reads cross-map between them at TIED divergence.", 0, NAVY),
        B("A family = a connected component of the read-conflict graph.", 0, NAVY),
        B("This is exactly the unit on which read-to-copy assignment must be solved.", 0, TEAL),
    ])

    add_content_slide(prs, "The definition (formal)", [
        B("Vertices  V  =  de-novo EXPRESSED loci (per-transcript intervals, not annotated gene spans)."),
        B("Placement: each alignment RECORD → its single best-overlap locus.", 0),
        B("⇒ two cross-locus placements come from distinct records = genuine multimapping, BY CONSTRUCTION", 1),
        B("(a single alignment over nested genes yields one entry → no distance guard needed).", 1),
        B("Conflict predicate (read r on loci i,j):   |de_i − de_j| ≤ Δ   AND   max(de_i,de_j) ≤ DE_max", 0, NAVY),
        B("tied at the HiFi error floor: minimap2 cannot decide which copy the read came from.", 1),
        B("Edge: ≥ MIN_READS = 3 conflicting reads.    Family: connected component, |C| ≥ 2.", 0),
    ], note="Δ=0.005, DE_max=0.05, MIN_READS=3 — derived, not tuned (see later slide).")

    add_content_slide(prs, "Worked example — one read becomes one edge", [
        B("Take a single HiFi molecule r that minimap2 places twice:", 0),
        B("primary on locus A  (de = 0.0010),   secondary on locus B  (de = 0.0035).", 1, NAVY),
        B("Test the conflict predicate:  |0.0010 − 0.0035| = 0.0025 ≤ Δ=0.005,  and both ≤ DE_max.", 0),
        B("⇒ r is TIED: the aligner cannot decide A vs B → r is a conflict read for the pair (A,B).", 1, TEAL),
        B("Repeat over all reads. 47 of them tie on (A,B) ≥ MIN_READS=3 → draw edge A—B.", 0),
        B("A and B land in one connected component ⇒ {A,B} is a family   (this IS RABL2A~RABL2B).", 0, NAVY),
        B("Contrast: a read with de_A=0.001, de_B=0.018 → |diff|=0.017 > Δ → decidable → NO edge", 0, ORANGE),
        B("(the EEF1A1 retrocopy: distinguishable, correctly never a family).", 1, ORANGE),
    ])

    add_content_slide(prs, "Why raw divergence (de), not the aligner's score (AS)", [
        B("AS (alignment score) folds in read LENGTH →"),
        B("a long retrocopy alignment scores like the real copy → AS TIES a false positive.", 1, ORANGE),
        B("de (gap-compressed per-base divergence) sees the actual sequence gap.", 0),
        B("de-tie  ⊊  AS-tie   — a shipped regression invariant.", 0, NAVY),
        B("The principled choice removes exactly the false positives:", 0),
        B("EEF1A1 retrocopy: AS ties 3,347 reads → de = 0 (|de diff| 0.016 > Δ, decidable).", 1, ORANGE),
    ])

    add_figure_slide(
        prs, "Examples & counterexamples — 17-candidate panel (TP=7, TN=10, FP=0, FN=0)",
        FIG_CONFLICT,
        "(A) green = de-tie families (components); red dashed = AS false-positive edges avoided.   "
        "(B) de fires on the 7 families; AS over-links the decoys (EEF1A1-retro: AS 3,347 vs de 0).")

    add_content_slide(prs, "Examples — the 7 families found (reads ARE confused)", [
        B("RABL2A ~ RABL2B — recent paralog, separate chromosomes (47 de-conflict reads)."),
        B("AK6 ~ LOC,  CCDC196 ~ LOC — high-identity cross-chromosome copies (24 each)."),
        B("MAGEA de-novo arrays ×4 — co-located tandem copies (75–303 conflict reads)."),
        B("Every edge verified read-by-read from raw GGO.bam (samtools/pysam).", 0, TEAL),
        B("All three regimes covered: recent paralog · cross-chrom copy · co-located array.", 0),
    ])

    add_content_slide(prs, "Counterexamples — three ways it correctly says 'NOT a family'", [
        B("THE CRUX — APOBEC3:  a CONFIRMED Ensembl-Compara recent paralog  →  STILL 0 conflict.", 0, NAVY),
        B("recent paralogy is NOT sufficient: APOBEC3's copies are read-resolvable, so they are excluded.", 1, NAVY),
        B("⇒ the definition is  recent-paralogy  ∩  read-confusability  — a strict subset, not a sequence test.", 1, TEAL),
        B("Domain-sharers  (CREB1~METTL21A, 192 shared reads)  →  0 conflict.", 0),
        B("single-valued best-overlap selects the SAME physical record for both → 0 edges (by construction).", 1),
        B("Retrocopies  (EEF1A1, AS ties 3,347)  →  0 conflict.", 0),
        B("reads decidable: |de diff| 0.016 > Δ → no tie.", 1),
    ])

    add_content_slide(prs, "Three structural robustnesses — by construction, not by threshold", [
        B("Domain-sharers: single-valued best-overlap → 0 cross-locus edges (verified 0/192, 0/429)."),
        B("Over-split fragments: co-positioned isoforms share no conflicting read.", 0),
        B("a sequence-SIMILARITY definition made 42% false families here; the conflict definition needs no guard.", 1),
        B("Retrocopies / decidable paralogs: decisive reads don't tie → excluded.", 0),
        B("None of these depend on Δ — they hold for any threshold, GIVEN tight per-locus vertices.", 0, TEAL),
        B("The one precondition: vertices must be per-transcript loci (the production unit).", 1),
        B("coarse gene-span vertices reintroduce over-merge — and that bridge is exactly what collapses 59→20", 1),
        B("when we move to de-novo loci (next slide). So the only knob is vertex resolution, and it is pinned.", 1),
    ])

    add_content_slide(prs, "The thresholds are DERIVED, not tuned", [
        B("Δ = 0.005 = the single-read divergence-discrimination RESOLUTION at HiFi error.", 0, NAVY),
        B("per-read de has SE √(ε/L) ≈ 0.0009; the TIE statistic |de_i − de_j| has SE √(2ε/L) ≈ 0.0013.", 1),
        B("⇒ two copies differing < ~4σ ≈ 0.005 are indistinguishable by one read — a measurement constant.", 1),
        B("the panel Δ-sweep plateau (correct down to 0.005, first error at 0.007) CONFIRMS it empirically.", 1),
        B("DE_max = 0.05: loose copy-vs-distinct-gene ceiling.   MIN_READS = 3: minimum-evidence quorum.", 0),
        B("Independent Ensembl Compara: RABL2 (Homininae), MAGEA (Catarrhini) confirmed recent paralogs;", 0),
        B("APOBEC3 confirms the strict-subset (recent paralog, yet correctly excluded).", 1),
    ])

    add_figure_slide(
        prs, "Genome-wide — stable, sensible, NOT overfit to the panel",
        FIG_GW,
        "34k independent gene vertices → 416 families, 57% size-2, 86% co-located; Δ flat to ±7%; "
        "production de-novo loci collapse the coarse over-merge bridges 59→20 (worst mega-bridges vanish).")

    add_content_slide(prs, "Honest limits (disclosed up front)", [
        B("FN=0 holds for ≤6-copy families: minimap2's default secondary cap (N=5) fragments LARGER arrays.", 0, ORANGE),
        B("MEASURED fix: uncapped re-align (−N 50 −p 0.1) heals a ≥12-copy chr23 array (core 5→11) at 0 FP.", 1),
        B("The panel Δ-plateau was empirically confirmed against 3 resolvable decoys (could be tightened genome-wide).", 0),
        B("DE_max set conservatively, not swept for its own plateau.", 0),
        B("Genome-wide families are dominated by uncharacterized LOC arrays (no orthology DB) → co-location is the proxy.", 0),
        B("~20 residual cross-chrom families at production resolution = genuine dispersed retrocopies (correct positives).", 0),
        B("RABL2's recovery is quorum-carried (recent-paralog regime: MIN_READS, not a per-read tie).", 0),
    ])

    add_content_slide(prs, "Summary — an airtight, RNA-derived family definition", [
        B("Family = connected component of the read-conflict graph  =  recent-paralogy ∩ read-confusability.", 0, NAVY),
        B("No tuned similarity threshold — the boundary is the HiFi error-model resolution.", 0),
        B("Three structural robustnesses BY CONSTRUCTION (domain-sharer · over-split · retrocopy).", 0),
        B("TP=7  TN=10  FP=0  FN=0  on the panel; genome-wide stable and de-novo-clean.", 0, TEAL),
        B("It is the exact unit the read-to-copy assignment problem operates on.", 0, TEAL),
    ])

    prs.save(OUT)
    print(f"[+] wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    print(f"[+] wrote {FIG_GW}")


if __name__ == "__main__":
    build()
