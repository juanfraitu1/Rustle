#!/usr/bin/env python
"""Build a slide deck (.pptx) explaining the multi-copy gene-family DEFINITION + how we CORROBORATE it.
Figures are generated from the project's own results. Output: bench/slides/multicopy_families.pptx (+ PNGs).

Run: /home/juanfra/miniforge3/bin/python bench/make_slides.py
"""
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mp
import networkx as nx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = os.path.join(os.path.dirname(__file__), "slides")
os.makedirs(OUT, exist_ok=True)
plt.rcParams.update({"font.size": 15, "font.family": "DejaVu Sans", "axes.spines.top": False,
                     "axes.spines.right": False, "figure.dpi": 150})
BLUE, RED, GREEN, ORANGE, PURPLE, GREY = "#377eb8", "#e41a1c", "#4daf4a", "#ff7f00", "#984ea3", "#999999"


def savefig(name):
    p = os.path.join(OUT, name)
    plt.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close()
    return p


# ---------------------------------------------------------------- figures
def fig_problem():
    fig, ax = plt.subplots(figsize=(9, 4.2))
    ax.axis("off")
    for x, lab, c in [(1.0, "copy A", BLUE), (6.0, "copy B", RED)]:
        ax.add_patch(plt.Rectangle((x, 2.2), 3.2, 0.7, fc=c, ec="black", alpha=.85))
        ax.text(x + 1.6, 2.55, lab, ha="center", va="center", color="white", fontsize=14, weight="bold")
    ax.text(5.1, 3.55, "99.x% identical paralogs", ha="center", fontsize=13, style="italic")
    # reads in the middle, ambiguous
    for k, y in enumerate([1.3, 0.9, 0.5]):
        ax.add_patch(plt.Rectangle((4.2, y), 1.6, 0.18, fc=GREY, ec="black"))
        ax.annotate("", (2.6, 2.2), (4.4, y + .1), arrowprops=dict(arrowstyle="->", color=BLUE, ls="--"))
        ax.annotate("", (7.6, 2.2), (5.6, y + .1), arrowprops=dict(arrowstyle="->", color=RED, ls="--"))
    ax.text(5.0, 0.0, "reads fit BOTH copies equally  →  minimap2 MAPQ = 0", ha="center", fontsize=14, color=RED, weight="bold")
    ax.set_xlim(0, 10); ax.set_ylim(-0.4, 4)
    return savefig("fig_problem.png")


def fig_conflict_graph():
    fig, ax = plt.subplots(figsize=(9, 4.6))
    G = nx.Graph()
    fam1 = ["A0", "A1", "A2"]; fam2 = ["B0", "B1"]
    single = ["S"]; dom = ["D0", "D1"]
    G.add_nodes_from(fam1 + fam2 + single + dom)
    G.add_edges_from([("A0", "A1"), ("A1", "A2"), ("A0", "A2"), ("B0", "B1")])   # ties
    pos = {"A0": (0, 1), "A1": (1, 1.6), "A2": (1, .4), "B0": (3, 1.3), "B1": (4, .8),
           "S": (5.7, 1.1), "D0": (3.3, -.7), "D1": (4.3, -.7)}
    nx.draw_networkx_nodes(G, pos, nodelist=fam1, node_color=BLUE, node_size=1500, ax=ax)
    nx.draw_networkx_nodes(G, pos, nodelist=fam2, node_color=GREEN, node_size=1500, ax=ax)
    nx.draw_networkx_nodes(G, pos, nodelist=single, node_color=GREY, node_size=1500, ax=ax)
    nx.draw_networkx_nodes(G, pos, nodelist=dom, node_color=GREY, node_size=1500, ax=ax)
    nx.draw_networkx_edges(G, pos, width=3, edge_color="black", ax=ax)
    nx.draw_networkx_labels(G, pos, font_color="white", font_weight="bold", font_size=12, ax=ax)
    ax.text(0.5, 2.15, "family 1\n(tandem array)", ha="center", color=BLUE, fontsize=12, weight="bold")
    ax.text(3.5, 1.8, "family 2\n(cross-chrom pair)", ha="center", color=GREEN, fontsize=12, weight="bold")
    ax.text(5.7, 1.7, "single-copy\n(no tie → not a family)", ha="center", color="black", fontsize=11)
    ax.text(3.8, -1.25, "domain-sharers\n(share 1 exon, NO tie → not merged)", ha="center", color="black", fontsize=11)
    ax.text(2.0, -2.0, "edge = a read ties between two loci   •   family = connected component   •   no threshold",
            ha="center", fontsize=12.5, style="italic")
    ax.axis("off"); ax.set_ylim(-2.3, 2.6); ax.set_xlim(-0.8, 6.6)
    return savefig("fig_conflict_graph.png")


def fig_detie():
    fig, (a1, a2) = plt.subplots(1, 2, figsize=(10, 3.6))
    for ax, (dA, dB, title, edge, col) in zip(
        [a1, a2],
        [(0.010, 0.012, "TIE  →  edge (same family)", True, GREEN),
         (0.001, 0.020, "RESOLVED  →  no edge", False, RED)]):
        ax.barh([1, 0], [dA, dB], color=[BLUE, RED], height=.5)
        ax.set_yticks([1, 0]); ax.set_yticklabels(["divergence to A", "divergence to B"])
        ax.axvline(0.05, color="black", ls=":", lw=1.5)
        ax.text(0.05, 1.7, "de_max ceiling", fontsize=10, ha="center")
        ax.set_xlim(0, 0.06); ax.set_xlabel("read divergence  (de)")
        ax.set_title(title, color=col, fontsize=13, weight="bold")
    a1.text(0.011, -0.7, "|Δde| ≤ δ  &  both fit", fontsize=11, color=GREEN)
    a2.text(0.012, -0.7, "|Δde| ≫ δ  → one clear winner", fontsize=11, color=RED)
    fig.suptitle("The de-tie criterion — the boundary is a property of the data, not a cutoff", fontsize=13)
    return savefig("fig_detie.png")


def fig_sim():
    fig, ax = plt.subplots(figsize=(8.5, 4))
    labels = ["O2 accuracy\n(resolvable)", "K=0 floor\ncertified Tied", "false\nmerges", "mis-\nassignments"]
    vals = [100, 100, 0, 0]; cols = [GREEN, BLUE, GREEN, GREEN]
    b = ax.bar(labels, vals, color=cols, ec="black")
    ax.bar_label(b, labels=["460/460", "120/120", "0", "0"], padding=4, fontsize=13, weight="bold")
    ax.set_ylim(0, 115); ax.set_ylabel("%")
    ax.set_title("Fully-simulated PLANTED genome (non-circular ground truth)", fontsize=13, weight="bold")
    ax.text(0.5, -22, "copies/divergence/CN planted • read labels known • O1: K0tandem 3/3, collapse 5/5, cnv 3/3, xchrom 2/2",
            transform=ax.transAxes, ha="center", fontsize=10.5, style="italic")
    return savefig("fig_sim.png")


def fig_dna():
    fig, ax = plt.subplots(figsize=(8.5, 4))
    b = ax.bar(["chance\n(1/K)", "our gate\n(held-out DNA)"], [43.2, 97.2], color=[GREY, GREEN], ec="black", width=.55)
    ax.bar_label(b, labels=["43.2%", "97.2%"], padding=4, fontsize=14, weight="bold")
    ax.set_ylim(0, 112); ax.set_ylabel("held-out DNA-column confirmation (%)")
    ax.annotate("2.2× over chance\n(K≥3 subset: 3.9×)", (1, 97.2), (0.35, 70), fontsize=12, color=GREEN,
                weight="bold", arrowprops=dict(arrowstyle="->", color=GREEN))
    ax.set_title("DNA-supervised decode — non-circular, REAL gorilla data\n(copies & columns from DNA, no RNA self-reference)",
                 fontsize=12.5, weight="bold")
    return savefig("fig_dna.png")


def fig_levels():
    fig, ax = plt.subplots(figsize=(8.5, 4))
    labels = ["any level\n(RNA∪DNA∪protein)", "DNA segdup\n(SEDEF)", "protein\northolog", "NONE\n(possible FP)"]
    vals = [85, 73, 30, 15]; cols = [GREEN, BLUE, ORANGE, RED]
    b = ax.bar(labels, vals, color=cols, ec="black")
    ax.bar_label(b, fmt="%d%%", padding=4, fontsize=13, weight="bold")
    ax.set_ylim(0, 100); ax.set_ylabel("% of RNA-defined families")
    ax.set_title("Cross-level corroboration of the RNA families", fontsize=13, weight="bold")
    ax.text(0.5, -0.30, "protein 'only 30%' is a FLOOR — 30% of duplications are non-coding / pseudogene (no protein ortholog)",
            transform=ax.transAxes, ha="center", fontsize=10, style="italic")
    return savefig("fig_levels.png")


# ---------------------------------------------------------------- deck
def build(figs):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    W = prs.slide_width

    def textbox(slide, text, left, top, width, height, size, bold=False, color=(20, 20, 20), align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame; tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line; p.alignment = align
            r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = RGBColor(*color)
        return tb

    def title_bar(slide, title):
        tb = textbox(slide, title, 0.6, 0.3, 12.1, 1.0, 30, bold=True, color=(20, 40, 90))

    def fig_slide(title, png, caption=None):
        s = prs.slides.add_slide(blank); title_bar(s, title)
        from PIL import Image
        iw, ih = Image.open(png).size
        maxw, maxh = 11.5, 5.2
        scale = min(maxw / (iw / 150), maxh / (ih / 150))
        w = (iw / 150) * scale
        s.shapes.add_picture(png, Inches((13.333 - w) / 2), Inches(1.45), width=Inches(w))
        if caption:
            textbox(s, caption, 0.8, 6.7, 11.7, 0.6, 14, italic=True, align=PP_ALIGN.CENTER, color=(90, 90, 90))
        return s

    def bullets_slide(title, bullets):
        s = prs.slides.add_slide(blank); title_bar(s, title)
        tb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.4))
        tf = tb.text_frame; tf.word_wrap = True
        for i, (txt, lvl, bold, col) in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ("•  " if lvl == 0 else "      – ") + txt
            p.level = lvl
            r = p.runs[0]; r.font.size = Pt(22 if lvl == 0 else 18); r.font.bold = bold
            r.font.color.rgb = RGBColor(*col); p.space_after = Pt(8)
        return s

    # 1 title
    s = prs.slides.add_slide(blank)
    textbox(s, "Multi-copy gene families at the RNA level", 0.8, 2.4, 11.7, 1.3, 40, bold=True, color=(20, 40, 90), align=PP_ALIGN.CENTER)
    textbox(s, "What the definition is  ·  and how we corroborate it", 0.8, 3.8, 11.7, 0.8, 22, italic=True, align=PP_ALIGN.CENTER, color=(90, 90, 90))
    textbox(s, "gorilla HiFi / Iso-Seq  ·  annotation-free  ·  threshold-free", 0.8, 4.6, 11.7, 0.6, 16, align=PP_ALIGN.CENTER, color=(120, 120, 120))

    # 2 problem
    fig_slide("The problem: near-identical copies confuse the aligner", figs["problem"],
              "When copies are ~identical, a read fits several equally → minimap2 leaves it at MAPQ 0. The copies are 'collapsed'.")

    # 3 definition
    s = bullets_slide("The definition", [
        ("A multi-copy gene family = a set of loci the READS cannot tell apart.", 0, True, (20, 40, 90)),
        ("Build a graph: nodes = loci; an edge = some read fits BOTH loci about equally (a 'tie').", 0, False, (20, 20, 20)),
        ("A family = a connected component of that read-conflict graph.", 0, True, (20, 20, 20)),
        ("A locus with no tie stands alone — it needs no resolution, so it is not a 'family'.", 0, False, (20, 20, 20)),
        ("No similarity cutoff — the boundary is decided by the data, not a threshold.", 0, True, (200, 30, 30)),
    ])

    # 4 conflict graph fig
    fig_slide("The read-conflict graph", figs["graph"],
              "Tied loci form families; single-copy genes and domain-sharers (one shared exon, no tie) are correctly left out.")

    # 5 de-tie
    fig_slide("Threshold-free: the de-tie criterion", figs["detie"],
              "Two placements tie iff their divergences are within δ AND both fit (≤ de_max). The tie IS the data, not a tuned cutoff.")

    # 6 what we use from the BAM
    bullets_slide("What we read from the BAM", [
        ("Placements — primary + secondary alignments (the multimappings) → the candidate loci & the edges.", 0, False, (20, 20, 20)),
        ("de tag / NM (edit distance) → the divergence that decides the tie.", 0, False, (20, 20, 20)),
        ("CIGAR (introns) → splice structure: assembles the copies + the read's junction fingerprint.", 0, False, (20, 20, 20)),
        ("MAPQ — both-zero corroborates a genuine multimapper.", 0, False, (20, 20, 20)),
        ("(Read bases at PSV positions are used LATER — to assign a read to a specific copy, not to define the family.)", 0, False, (110, 110, 110)),
    ])

    # 7 corroboration overview
    bullets_slide("How we corroborate it — four independent axes", [
        ("Simulated PLANTED genome — non-circular ground truth (we know every answer).", 0, True, (20, 120, 20)),
        ("DNA-supervised decode — real gorilla data, copies/columns from DNA, held-out CV.", 0, True, (30, 80, 160)),
        ("Cross-level — does an RNA family correspond to a DNA segduplication / protein orthogroup?", 0, True, (200, 110, 0)),
        ("Error controls — is a 'copy difference' real, or sequencing / alignment error?", 0, True, (150, 30, 130)),
    ])

    # 8 sim
    fig_slide("Corroboration 1 — simulated ground truth", figs["sim"],
              "Plant the genome, label every read, run the unmodified pipeline. 100% accuracy where resolvable; the identical floor certified Tied; zero false merges.")

    # 9 dna
    fig_slide("Corroboration 2 — DNA-supervised (non-circular, real data)", figs["dna"],
              "Copies & distinguishing columns come from DNA, never from the RNA itself; held-out DNA columns confirm 97.2% of calls.")

    # 10 levels
    fig_slide("Corroboration 3 — the RNA families are real duplications", figs["levels"],
              "73% are independently confirmed DNA segduplications (SEDEF); 85% confirmed at some level. The RNA boundary is not an artifact.")

    # 11 error
    bullets_slide("Corroboration 4 — 'how do you know it isn't an error?'", [
        ("Recurrence — errors are random & independent; a variant seen in k reads has error-probability ε^k → ~0.", 0, False, (20, 20, 20)),
        ("Calibration — every call carries a p-value; the gate ABSTAINS rather than guess (no 1/k).", 0, False, (20, 20, 20)),
        ("Linkage — real copy-variants co-segregate into haplotypes; errors don't.", 0, False, (20, 20, 20)),
        ("Alignment error — MEASURED vs DNA truth: genuine artifacts are ~1% of columns and change 0.03% of calls.", 0, True, (150, 30, 130)),
        ("So we MODEL error (calibrated, abstain) — we never correct or trust a single read.", 0, True, (20, 40, 90)),
    ])

    # 12 summary
    bullets_slide("Summary", [
        ("Definition: a family = a connected component of the read-conflict graph — threshold-free, annotation-free.", 0, True, (20, 40, 90)),
        ("It picks out exactly the loci where copy assignment is needed; excludes single-copy & domain-sharers by construction.", 0, False, (20, 20, 20)),
        ("Corroborated four independent, largely non-circular ways: planted sim, DNA decode, cross-level, error controls.", 0, True, (20, 120, 20)),
        ("View it in IGV: copy_assign --gtf (transcripts, families flagged) + reads coloured by assigned copy.", 0, False, (110, 110, 110)),
    ])

    path = os.path.join(OUT, "multicopy_families.pptx")
    prs.save(path)
    return path


def main():
    figs = {"problem": fig_problem(), "graph": fig_conflict_graph(), "detie": fig_detie(),
            "sim": fig_sim(), "dna": fig_dna(), "levels": fig_levels()}
    path = build(figs)
    print(f"wrote {path}  ({len(figs)} figures, 12 slides)")


if __name__ == "__main__":
    main()
