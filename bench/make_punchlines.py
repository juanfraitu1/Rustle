#!/usr/bin/env python
"""Punchline deck — one idea per slide for the whole method. Big text, minimal graphics: a talking deck.
Every claim is distilled from already-verified session artifacts (VG evidence deck, DEFINITIONS_FORMAL.md,
SOTO_CONCORDANCE.md, HUMAN_CROSSSPECIES.md, VALIDATION_AND_STATUS.md), with the honest hedges kept.

Output: bench/slides/method_punchlines.pptx
Run: /home/juanfra/miniforge3/bin/python bench/make_punchlines.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = os.path.join(os.path.dirname(__file__), "slides")
os.makedirs(OUT, exist_ok=True)
NAVY = (20, 40, 90); DARK = (35, 35, 40); GREY = (110, 115, 125); LGREY = (150, 152, 160)
RED = (200, 45, 45); BLUE = (40, 95, 175); PURPLE = (120, 55, 150); GREEN = (30, 120, 55); ORANGE = (200, 120, 15)
TEAL = (20, 130, 130)

# (kicker, question, punchline, support, accent)
SLIDES = [
    ("THE PROBLEM",
     "When gene copies are ~identical, what happens?",
     "A read fits several copies equally, so the aligner gives up — MAPQ 0 — and the copies “collapse” into one. We define and RESOLVE these multi-copy gene families straight from the RNA reads. We are not an assembler.",
     "Short reads recover only 0.85% of the paralog-specific variants in these regions (Soto 2025) — long reads + a per-read model are the way in.",
     RED),

    ("O1  ·  WHAT IS A MULTI-COPY GENE FAMILY?",
     "What is a multi-copy gene family — for us?",
     "A set of mutually-homologous loci: a cohesive, dense cluster (a γ-quasi-clique) in the transcribed-homology graph, over ≥ 2 distinct loci.",
     "The boundary is a relational homology tie — no “≥ X% identical” cutoff, annotation-free. Copy number = the fewest copy-paths that explain the reads (a minimum path cover; a lower bound).",
     NAVY),

    ("O2  ·  COPY ASSIGNMENT",
     "Why not Eichler's rule — assign a copy only if the best alignment beats the runner-up by ≥ 10 — or just Liftoff?",
     "Near-identical copies never differ by that margin, so Eichler's whole-read rule resolves 0% of the COLLAPSED paralogs. Liftoff only re-places copies the assembly already has — and collapses the near-identical ones. We score the PSV columns PER READ, and assign-or-abstain with a certificate. Never 1/k.",
     "AS ≥ 10: 0% at every K in the collapsed regime (it works fine on divergent copies — blind only sub-threshold). Scoring only the PSVs: 60% (K=1) → 100% (K≥2), 0.99–1.00 accuracy. Liftoff missed copies we recover.",
     BLUE),

    ("THE ENGINE  ·  VARIATION GRAPH",
     "How does the variation graph “thread” the PSVs?",
     "All copies collapse onto ONE graph: the shared backbone is the spine, the copy differences are BUBBLES (the PSVs). A copy is a consistent PATH of alleles; a read walks the graph and VOTES at each bubble it spans → assign or abstain.",
     "A SUN — a copy's private allele — pins that copy from a single read. Multimapping reads become SHARED evidence, threaded, not divided 1/k.",
     PURPLE),

    ("RIGOR  ·  IS THE PSV REAL?",
     "How do we know a PSV is a real copy-difference, not a sequencing or alignment error?",
     "It must survive five independent gates — an error clears none of them:",
     None,  # special: 5-gate slide
     GREEN),

    ("O4  ·  REFERENCE-ABSENT COPIES",
     "How do we find copies that aren't even in the reference?",
     "When the reads need MORE copies than the reference hosts — χ(H) at a locus exceeds its annotated copies — the ambiguity can't be explained by what's there. That surplus IS a reference-absent copy.",
     "The same assignment engine, turned OUTWARD: detect-and-flag. (Whether the surplus is a new copy vs an allele needs DNA.)",
     ORANGE),

    ("O3  ·  ALLELE-SPECIFIC JUNCTIONS",
     "Does the same idea reach beyond copy number?",
     "Yes — the SAME per-molecule threading links an ALLELE to the splice JUNCTION it uses, on ONE read. Which variant a molecule carries, and which isoform it splices to, are read off the same molecule — so there is nothing to phase, and nothing to divide 1/k. That is allele-specific splicing, straight from single reads.",
     "A different question — cis-regulation, not copy number — from the SAME engine (a PSV bubble separates copies; an allele bubble separates haplotypes; a junction is an edge choice). Shipped: 54 strand-bias-clean links; the genome-false over-calls were checked against DNA and retracted — reported honestly.",
     TEAL),

    ("WHY IT HOLDS UP",
     "Is any of this real, or wishful thinking?",
     "Planted sim (non-circular): identical copies → 0 invented PSVs (120/120 certified tied); real PSVs recovered. Independent DNA (Soto, Cell 2025): we AGREE on his human-specific families. Same code on HUMAN testis: it tracks the real expansions — MAGEA 2→11, TSPY 5→33 — not the gorilla numbers.",
     "One line: a family = read-threaded paths through one PSV-bubble graph; assignment is assign-or-abstain over provable structure (χ(H) = min copies; facility location). No thresholds, no 1/k.",
     NAVY),
]

SEQ_ROWS = [
    ("copy", "a genomic instance of a gene — a path through the graph"),
    ("isoform", "a splice variant of ONE gene — a path through the junctions"),
    ("exon / intron", "kept in the mRNA  /  spliced out"),
    ("splice junction", "where two exons join (a branch point)"),
    ("allele", "the base at a variant position"),
    ("PSV", "a position where copies differ (a bubble)"),
    ("SUN", "a PSV private to ONE copy  (SUN ⊆ PSV) → one read pins it"),
]
METHOD_ROWS = [
    ("variation graph", "spine + PSV bubbles + junction branches; copies/isoforms = paths"),
    ("bubble", "where paths diverge — a PSV or a junction"),
    ("family  (O1)", "a cohesive homology cluster of loci (γ-quasi-clique)"),
    ("copy number  χ_H", "fewest copy-paths covering the reads — a lower bound"),
    ("assignment  (O2)", "which copy-path a read is on (facility location; no 1/k)"),
    ("reference-absent  (O4)", "a copy the individual has but the assembly lacks"),
    ("allele-specific jn  (O3)", "an allele linked to its junction on one read"),
    ("K=0 floor", "identical copies carry no PSV → certified TIED, not guessed"),
]
EVO_ROWS = [
    ("paralog", "two genes related by DUPLICATION (within a lineage)"),
    ("ortholog", "two genes related by SPECIATION (the same gene in two species)"),
    ("multi-copy gene family", "≥ 2 paralogs present in one genome — the state we measure"),
    ("segmental duplication (SD)", "a duplicated DNA block (≥ 1 kb, high identity); SD98 = ≥ 98%"),
    ("expansion", "copies GAINED on a lineage vs its ancestor (needs an outgroup)"),
    ("contraction", "copies LOST on a lineage"),
    ("pseudogene", "a copy that lost function (processed = intronless retrocopy)"),
    ("retrocopy / retrogene", "an mRNA-derived, intronless copy"),
    ("exonization", "non-coding sequence (intron / Alu) becoming a new EXON"),
]

GATES = [
    ("RECURRENCE", "an error hits one read (prob ε); a real PSV recurs at the same column → εᵏ → 0"),
    ("CALIBRATION", "a per-read significance test; if it isn't significant we ABSTAIN — we don't guess (no 1/k)"),
    ("LINKAGE", "real copy-variants co-segregate into ONE consistent path; independent errors don't"),
    ("EDITING ≠ PSV", "A→I editing (A→G / T→C) is flagged and down-weighted, not counted as a copy difference"),
    ("OPTIMAL ALIGNMENT", "a misplaced gap invents a PHANTOM PSV — our exact banded DP aligner is provably optimal"),
]


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def tb(slide, text, l, t, w, h, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, spacing=None):
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line; p.alignment = align
            if spacing: p.line_spacing = spacing
            r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = RGBColor(*color)
        return box

    def accent(slide, color):
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.34), Inches(7.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*color); bar.line.fill.background()

    # ---- title
    s = prs.slides.add_slide(blank)
    accent(s, NAVY)
    tb(s, "The method, in punchlines", 0.9, 2.6, 11.6, 1.3, 44, NAVY, bold=True)
    tb(s, "Multi-copy gene families at the RNA level — one idea per slide", 0.9, 3.9, 11.6, 0.8, 22, GREY, italic=True)

    for kicker, question, punch, support, accent_c in SLIDES:
        s = prs.slides.add_slide(blank)
        accent(s, accent_c)
        tb(s, kicker, 0.85, 0.5, 11.8, 0.5, 15, accent_c, bold=True)
        tb(s, question, 0.85, 1.15, 11.8, 1.0, 21, GREY, italic=True)
        if support is None:
            # the five-gate slide
            tb(s, punch, 0.85, 2.15, 11.8, 0.7, 24, NAVY, bold=True)
            y = 3.15
            for name, desc in GATES:
                tb(s, name, 0.95, y, 3.05, 0.6, 17, accent_c, bold=True)
                tb(s, desc, 4.05, y, 8.5, 0.7, 15.5, DARK)
                y += 0.78
        else:
            tb(s, punch, 0.85, 2.35, 11.9, 3.0, 29, NAVY, bold=True, spacing=1.05)
            tb(s, support, 0.85, 5.9, 11.9, 1.4, 15.5, GREY, italic=True, spacing=1.05)

    # ---- backup slide: one framework (how the graphs relate) ----
    import make_one_framework
    of = os.path.join(OUT, "one_framework.png")
    if not os.path.exists(of):
        make_one_framework.build()
    s = prs.slides.add_slide(blank); accent(s, NAVY)
    tb(s, "BACKUP  ·  one framework, not four graphs", 0.85, 0.5, 11.8, 0.5, 16, NAVY, bold=True)
    tb(s, "if asked how the graphs relate: they are ONE object — the variation graph — at different scopes / views",
       0.85, 1.12, 11.8, 0.6, 15, GREY, italic=True)
    from PIL import Image as _Im1
    iw, ih = _Im1.open(of).size
    scale = min(12.7 / (iw / 150), 5.4 / (ih / 150)); w = (iw / 150) * scale
    s.shapes.add_picture(of, Inches((13.333 - w) / 2), Inches(1.8), width=Inches(w))

    # ---- backup slide: the two formal objects ----
    import make_formal_objects
    fo = os.path.join(OUT, "formal_objects.png")
    if not os.path.exists(fo):
        make_formal_objects.build()
    s = prs.slides.add_slide(blank); accent(s, NAVY)
    tb(s, "BACKUP  ·  the two formal objects", 0.85, 0.5, 11.8, 0.5, 16, NAVY, bold=True)
    tb(s, "if asked for the formalism: what exactly IS a family, and what is the copy number?",
       0.85, 1.12, 11.8, 0.6, 15, GREY, italic=True)
    from PIL import Image
    iw, ih = Image.open(fo).size
    maxw, maxh = 12.3, 4.7
    scale = min(maxw / (iw / 150), maxh / (ih / 150)); w = (iw / 150) * scale
    s.shapes.add_picture(fo, Inches((13.333 - w) / 2), Inches(1.95), width=Inches(w))
    tb(s, "family = γ-quasi-clique (dense homology cluster — not a chain, not a strict clique)    ·    copy number = χ(H) = fewest copies that explain the reads = min clique cover (a lower bound)",
       0.6, 6.75, 12.1, 0.6, 13, NAVY, align=PP_ALIGN.CENTER)

    # ---- backup: glossary (two slides) ----
    def gloss_col(slide, x, y, w, header, hcolor, rows):
        tb(slide, header, x, y, w, 0.4, 13.5, hcolor, bold=True)
        box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.5), Inches(w), Inches(5.4))
        tf = box.text_frame; tf.word_wrap = True
        for i, (term, defn) in enumerate(rows):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r1 = p.add_run(); r1.text = term; r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = RGBColor(*NAVY)
            r2 = p.add_run(); r2.text = "  —  " + defn; r2.font.size = Pt(10.5); r2.font.color.rgb = RGBColor(88, 92, 102)
            p.space_after = Pt(7)

    s = prs.slides.add_slide(blank); accent(s, NAVY)
    tb(s, "BACKUP  ·  glossary  (1/2) — sequence & method", 0.85, 0.5, 11.8, 0.5, 16, NAVY, bold=True)
    gloss_col(s, 0.55, 1.4, 6.0, "THE SEQUENCE WORLD", BLUE, SEQ_ROWS)
    gloss_col(s, 6.95, 1.4, 6.1, "OUR METHOD OBJECTS", GREEN, METHOD_ROWS)

    s = prs.slides.add_slide(blank); accent(s, NAVY)
    tb(s, "BACKUP  ·  glossary  (2/2) — evolution", 0.85, 0.5, 11.8, 0.5, 16, NAVY, bold=True)
    gloss_col(s, 0.55, 1.4, 12.4, "THE EVOLUTION WORLD", ORANGE, EVO_ROWS)
    box = s.shapes.add_shape(1, Inches(0.5), Inches(5.7), Inches(12.4), Inches(1.35))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(244, 247, 251); box.line.color.rgb = RGBColor(*NAVY); box.line.width = Pt(1)
    tb(s, "the relationships that trip people:", 0.75, 5.85, 12.0, 0.4, 12, NAVY, bold=True)
    tb(s, "SUN ⊆ PSV     ·     paralog = the relationship,  multi-copy family = the state we measure,  SD = the DNA mechanism,  expansion = the cross-species change     ·     copy (bubbles) vs isoform (junctions) = two axes of ONE variation graph",
       0.75, 6.28, 12.0, 0.8, 11, DARK)

    for path in (os.path.join(OUT, "method_punchlines.pptx"), os.path.join(OUT, "method_punchlines.new.pptx")):
        try:
            prs.save(path); return path
        except PermissionError:
            continue
    return None


if __name__ == "__main__":
    print("wrote", build())
