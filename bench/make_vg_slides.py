#!/usr/bin/env python
"""Advisor deck (.pptx): a VARIATION-GRAPH-centric story for RNA multi-copy gene families.
Answers three questions: (1) WHY variation graphs help, (2) HOW threading keeps/uses PSVs,
(3) HOW we know a PSV is REAL (not error/artifact) — then PROOF + honest limits.

Every number is verified against a real file (see speaker notes, loaded from deck.json).
Reuses the house style of bench/make_slides.py. Output: bench/slides/vg_copies_as_paths.pptx (+ PNGs).

Run: /home/juanfra/miniforge3/bin/python bench/make_vg_slides.py
"""
import json
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mp
from matplotlib.patches import FancyArrowPatch, Circle, FancyBboxPatch
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(__file__)
OUT = os.path.join(HERE, "slides")
os.makedirs(OUT, exist_ok=True)
plt.rcParams.update({"font.size": 15, "font.family": "DejaVu Sans", "axes.spines.top": False,
                     "axes.spines.right": False, "figure.dpi": 150})
BLUE, RED, GREEN, ORANGE, PURPLE, GREY = "#377eb8", "#e41a1c", "#4daf4a", "#ff7f00", "#984ea3", "#999999"
NAVY = (20, 40, 90)
DARK = (20, 20, 20)
GDARK = (20, 120, 20)
RDARK = (200, 30, 30)
SPINE = "#c9c9c9"
YELLOW = "#ffe680"

# vg_deck_content.json holds the adversarially-verified content + speaker-note source citations
# (produced by the evidence workflow; committed next to this script so the deck is reproducible).
DECK = json.load(open(os.path.join(HERE, "vg_deck_content.json")))
NOTES = {s["n"]: s["speaker_notes"] for s in DECK["slides"]}


def savefig(name):
    p = os.path.join(OUT, name)
    plt.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close()
    return p


# ----------------------------------------------------------------- figures
def fig_flip():
    """WHY VG: N near-identical copies collapse into ONE graph; shared backbone -> spine, differences -> bubbles."""
    DARK, NAVY, GDARK, RDARK = "#1a1a1a", "#14285a", "#147814", "#c81e1e"
    fig, ax = plt.subplots(figsize=(12, 4.4))
    ax.axis("off")
    ax.set_xlim(0, 13); ax.set_ylim(-0.6, 4.4)
    cols = [RED, BLUE, GREEN, PURPLE]
    psv_x = [1.4, 2.8, 4.2]  # positions where the copies differ
    # LEFT: N stacked linear copies, ~identical, with a few colored PSV ticks
    ax.text(2.85, 4.15, "N copies  ·  99%+ identical", ha="center", fontsize=13, weight="bold", color=DARK)
    for i, c in enumerate(cols):
        y = 3.3 - i * 0.72
        ax.add_patch(mp.FancyBboxPatch((0.4, y), 4.6, 0.42, boxstyle="round,pad=0.02,rounding_size=0.08",
                                       fc="#eef2f7", ec=c, lw=2))
        for k, px in enumerate(psv_x):
            ax.add_patch(plt.Rectangle((0.4 + px * 0.95, y), 0.09, 0.42, fc=cols[(i + k) % 4], ec="none"))
    ax.text(2.7, 0.35, "as N separate sequences → a read that fits several is a nuisance (split 1/k)",
            ha="center", fontsize=10.5, style="italic", color=RDARK)
    # ARROW
    ax.add_patch(FancyArrowPatch((5.25, 2.1), (6.55, 2.1), arrowstyle="-|>", mutation_scale=28, lw=3, color="#555"))
    ax.text(5.9, 2.5, "collapse", ha="center", fontsize=12, weight="bold", color="#555")
    # RIGHT: one graph — gray spine + colored PSV bubbles
    ax.text(9.7, 4.15, "ONE variation graph", ha="center", fontsize=13, weight="bold", color=NAVY)
    sx0, sx1, sy = 7.0, 12.4, 2.1
    ax.add_patch(plt.Rectangle((sx0, sy - 0.13), sx1 - sx0, 0.26, fc=SPINE, ec="none"))
    ax.text((sx0 + sx1) / 2, sy - 0.55, "shared backbone (spine)", ha="center", fontsize=10.5, color="#777")
    bx = [8.1, 9.7, 11.3]
    for x in bx:  # each PSV = a bubble (two alleles off the spine)
        ax.plot([x, x], [sy, sy + 0.8], color="#bbb", lw=1)
        ax.plot([x, x], [sy, sy - 0.8], color="#bbb", lw=1)
        ax.add_patch(Circle((x, sy + 0.8), 0.17, fc=BLUE, ec="black", lw=1, zorder=3))
        ax.add_patch(Circle((x, sy - 0.8), 0.17, fc=RED, ec="black", lw=1, zorder=3))
    ax.text(9.7, sy + 1.25, "copy differences = PSV bubbles", ha="center", fontsize=11, color=NAVY, weight="bold")
    # a multimapping read threading the graph
    ax.add_patch(FancyArrowPatch((7.6, sy - 1.45), (11.8, sy - 1.45), arrowstyle="-", lw=6, color=GREEN, alpha=.8))
    ax.text(9.7, sy - 1.95, "one multimapping read = SHARED evidence for every copy-path it threads  (never 1/k)",
            ha="center", fontsize=10.5, color=GDARK, weight="bold")
    return savefig("vg_fig_flip.png")


def _node(ax, x, y, letter, fc, r=0.26):
    ax.add_patch(Circle((x, y), r, fc=fc, ec="black", lw=1.3, zorder=4))
    ax.text(x, y, letter, ha="center", va="center", fontsize=13, weight="bold", zorder=5,
            color="black" if fc in (YELLOW, "white") else "white")


def fig_threading():
    """HOW threading keeps PSVs: PSV=bubble, copy=path of alleles, read votes at each bubble -> assign or abstain."""
    DARK, NAVY, GDARK, RDARK = "#1a1a1a", "#14285a", "#147814", "#c81e1e"
    fig, ax = plt.subplots(figsize=(12, 5.4))
    ax.axis("off")
    ax.set_xlim(0, 12.4); ax.set_ylim(-3.0, 4.0)
    sy, TOPY, BOTY = 1.4, 2.6, 0.2
    ax.add_patch(plt.Rectangle((0.6, sy - 0.15), 10.0, 0.30, fc=SPINE, ec="none"))
    ax.text(0.65, sy + 0.35, "shared backbone", fontsize=10.5, color="#777")
    ax.text(0.6, 3.7, "each COPY = a consistent PATH of PSV alleles", fontsize=12, color=DARK, weight="bold")
    bx = [2.4, 4.8, 7.2, 9.6]
    top = [("A", BLUE), ("C", BLUE), ("G", BLUE), ("C", BLUE)]          # copy 1 (blue) path
    bot = [("G", RED), ("T", RED), ("T", YELLOW), ("A", RED)]           # copy 2 (red); bubble 3 = its private SUN
    for i, x in enumerate(bx):
        ax.plot([x, x], [sy + 0.15, TOPY], color="#bbb", lw=1, zorder=1)
        ax.plot([x, x], [sy - 0.15, BOTY], color="#bbb", lw=1, zorder=1)
        _node(ax, x, TOPY, top[i][0], top[i][1])
        _node(ax, x, BOTY, bot[i][0], bot[i][1])
    # two copy-paths threading their chosen alleles
    ax.plot(*zip(*([(0.6, sy)] + [(x, TOPY) for x in bx] + [(10.8, sy)])), color=BLUE, lw=2.4, alpha=.85, zorder=2)
    ax.plot(*zip(*([(0.6, sy)] + [(x, BOTY) for x in bx] + [(10.8, sy)])), color=RED, lw=2.4, alpha=.85, zorder=2)
    ax.text(11.0, TOPY, "copy 1", color=BLUE, fontsize=11.5, weight="bold", va="center")
    ax.text(11.0, BOTY, "copy 2", color=RED, fontsize=11.5, weight="bold", va="center")
    # SUN callout on bubble-3 bottom node (copy 2's private allele) — placed clear of the read arrows
    ax.annotate("SUN = copy 2's PRIVATE allele\n→ ONE read pins copy 2", (bx[2], BOTY), (8.15, -0.55),
                fontsize=10, color="#8a6d00", weight="bold", ha="left", va="center",
                arrowprops=dict(arrowstyle="->", color="#8a6d00"))
    # the read: threads bubbles 1-3, matching copy 2's alleles (incl. the private SUN) -> assigned copy 2
    ry = -1.7
    ax.add_patch(FancyArrowPatch((1.9, ry), (7.9, ry), arrowstyle="-", lw=8, color=GREEN, alpha=.65))
    for x in bx[:3]:
        ax.add_patch(FancyArrowPatch((x, ry + 0.15), (x, BOTY - 0.30), arrowstyle="-|>",
                     mutation_scale=13, lw=1.6, color=GREEN, alpha=.85, ls=(0, (2, 1))))
    ax.text(8.15, ry, "→ assigned copy 2", color=RED, fontsize=12.5, weight="bold", va="center")
    ax.text(4.9, ry - 0.42, "a READ votes copy 2's allele at every bubble it spans (incl. the private SUN)",
            ha="center", fontsize=10.5, color=GDARK, weight="bold")
    ax.text(6.2, -2.65, "consistent votes → ASSIGN      ·      no bubble spanned / conflicting votes → ABSTAIN      ·      never 1/k",
            ha="center", fontsize=11.5, color=NAVY, weight="bold")
    return savefig("vg_fig_threading.png")


def fig_realness():
    """Q3: is the PSV real or an error/artifact? five named legs, each with a verified number."""
    DARK, NAVY, GDARK, RDARK = "#1a1a1a", "#14285a", "#147814", "#c81e1e"
    fig, ax = plt.subplots(figsize=(13, 3.9))
    ax.axis("off")
    ax.set_xlim(0, 13); ax.set_ylim(0, 4)
    cards = [
        ("CALIBRATION", BLUE, "per-read p-value vs its\nOWN error rate; assign iff\np < α/(n−1) (Bonferroni)",
         "misassignment ≤ α\n(theorem)", "max p = 6.9e-4\n< α = 1e-3"),
        ("RECURRENCE", GREEN, "a base error is independent\nper read; a real PSV recurs\nat the same column",
         "joint error prob", "εᵏ → 0"),
        ("DNA-MEASURED", PURPLE, "artifact columns adjudicated\nby held-out CLEAN DNA\n(a different molecule)",
         "artifacts are inert", "0.03% of calls\n(6 / 17,376)"),
        ("EDITING ≠ PSV", ORANGE, "A→I mimics a variant:\nA→G (+) / T→C (−) flagged\n(Clair3-RNA) & down-weighted",
         "named guard,\nno DNA needed", "ε_edit → 0.2"),
        ("RECOMBINANT-\nABSTAIN", RED, "a read whose private SUNs\npoint to ≥2 copies belongs\nto NO single copy",
         "theorem boundary\n(pure-monotone)", "→ ABSTAIN"),
    ]
    w = 2.44
    for i, (title, c, mech, mid, num) in enumerate(cards):
        x = 0.15 + i * 2.56
        ax.add_patch(FancyBboxPatch((x, 0.15), w, 3.7, boxstyle="round,pad=0.02,rounding_size=0.06",
                                    fc="#f7f9fc", ec=c, lw=2.2))
        ax.add_patch(FancyBboxPatch((x, 3.15), w, 0.7, boxstyle="round,pad=0.02,rounding_size=0.06",
                                    fc=c, ec=c, lw=2.2))
        ax.text(x + w / 2, 3.5, title, ha="center", va="center", fontsize=11.5, weight="bold", color="white")
        ax.text(x + w / 2, 2.35, mech, ha="center", va="center", fontsize=9.3, color=DARK)
        ax.text(x + w / 2, 1.45, mid, ha="center", va="center", fontsize=9.2, style="italic", color="#666")
        ax.text(x + w / 2, 0.62, num, ha="center", va="center", fontsize=11.5, weight="bold", color=c)
    return savefig("vg_fig_realness.png")


def fig_sim():
    """PROOF: planted non-circular ground truth (current sim, post-hidive)."""
    DARK, NAVY, GDARK, RDARK = "#1a1a1a", "#14285a", "#147814", "#c81e1e"
    fig, ax = plt.subplots(figsize=(8.6, 4.1))
    labels = ["multimapping (MAPQ-0)\nassigned to CORRECT copy", "K=0 identical copies\ncertified TIED (not guessed)"]
    vals = [99.8, 100.0]
    b = ax.bar(labels, vals, color=[GREEN, BLUE], ec="black", width=.55)
    ax.bar_label(b, labels=["579/580", "120/120"], padding=4, fontsize=13, weight="bold")
    ax.set_ylim(0, 118); ax.set_ylabel("%")
    ax.set_title("Planted genome, every read labelled — non-circular ground truth", fontsize=13, weight="bold")
    ax.text(0.5, -0.30, "0 false PSVs invented on identical copies   ·   0 false merges   ·   aligner-invariant (DP = poasta, 0 diffs)",
            transform=ax.transAxes, ha="center", fontsize=10.5, style="italic", color=GDARK)
    return savefig("vg_fig_sim.png")


def crop_top(src, dst, frac=0.32):
    """keep the top `frac` of a tall real figure (title + reads-threaded band), drop the busy footer."""
    im = Image.open(src)
    w, h = im.size
    im.crop((0, 0, w, int(h * frac))).save(dst)
    return dst


def crop_box(src, dst, box):
    Image.open(src).crop(box).save(dst)
    return dst


# ----------------------------------------------------------------- deck
def build(assets):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def textbox(slide, text, left, top, width, height, size, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame; tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line; p.alignment = align
            r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = RGBColor(*color)
        return tb

    def title_bar(slide, title, kicker=None):
        if kicker:
            textbox(slide, kicker, 0.62, 0.22, 12.1, 0.4, 14, bold=True, color=(120, 130, 150))
        textbox(slide, title, 0.6, 0.52, 12.2, 0.9, 28, bold=True, color=NAVY)

    def notes(slide, n):
        slide.notes_slide.notes_text_frame.text = NOTES.get(n, "")

    def bullets(slide, items, left, top, width, height, size=15.5, gap=7):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame; tf.word_wrap = True
        for i, (txt, bold, col) in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "•  " + txt
            r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = RGBColor(*col)
            p.space_after = Pt(gap)
        return tb

    def picture(slide, png, left, top, width=None, height=None):
        kw = {}
        if width: kw["width"] = Inches(width)
        if height: kw["height"] = Inches(height)
        return slide.shapes.add_picture(png, Inches(left), Inches(top), **kw)

    def pic_fit(slide, png, left, top, maxw, maxh):
        iw, ih = Image.open(png).size
        scale = min(maxw / (iw / 150), maxh / (ih / 150))
        w = (iw / 150) * scale
        h = (ih / 150) * scale
        slide.shapes.add_picture(png, Inches(left + (maxw - w) / 2), Inches(top + (maxh - h) / 2), width=Inches(w))
        return w, h

    # ---- 0 TITLE
    s = prs.slides.add_slide(blank)
    textbox(s, "Copies as Paths", 0.8, 2.15, 11.7, 1.1, 46, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    textbox(s, "A variation-graph framework for RNA multi-copy gene families", 0.8, 3.35, 11.7, 0.7, 22,
            italic=True, align=PP_ALIGN.CENTER, color=(90, 90, 90))
    textbox(s, "the shared backbone collapses  ·  copy differences become bubbles  ·  a multimapping read is threaded (assign-or-abstain, never 1/k)",
            0.8, 4.15, 11.7, 0.8, 15, align=PP_ALIGN.CENTER, color=(120, 120, 120))
    textbox(s, "gorilla HiFi / Iso-Seq   ·   provable structure: χ(H) = min copy count, assignment = facility-location",
            0.8, 5.15, 11.7, 0.5, 14, align=PP_ALIGN.CENTER, color=(150, 150, 150))

    # ---- 1 WHY VG
    s = prs.slides.add_slide(blank)
    title_bar(s, "Why a variation graph — the flip on 1/k", "QUESTION 1")
    bullets(s, [
        ("The flip: N near-identical copies → ONE graph. The shared backbone collapses to a spine; only copy differences survive as PSV bubbles.",
         True, NAVY),
        ("A multimapping read = shared evidence threaded through every copy-path it touches → assignment is flow-decomposition / threading, never split 1/k.",
         False, DARK),
        ("Provable structure: min copy count = χ(H) (machine-checked); assignment = max-weight facility-location, (1−1/e)-OPT, EXACT under Strong Separation.",
         False, DARK),
        ("The whole-read AS≥10 rule (Eichler) resolves 0% of collapsed paralogs at every K — provably blind sub-threshold; scoring only the bubbles: 60% → 100% (K≥2), acc 0.99–1.00.",
         True, GDARK),
    ], 0.6, 1.55, 12.2, 2.7, size=15)
    pic_fit(s, assets["flip"], 0.6, 4.25, 12.1, 2.9)
    notes(s, 1)

    # ---- 2 THREADING
    s = prs.slides.add_slide(blank)
    title_bar(s, "Threading — how the read keeps and uses the PSVs", "QUESTION 2")
    bullets(s, [
        ("PSV = a bubble; a copy = a consistent path of alleles; a read threads the graph and VOTES at each bubble it spans → assign or abstain.",
         True, NAVY),
        ("A SUN (Sudmant 2010) is a bubble where one copy's allele is PRIVATE → a single read pins that copy (a Strong-Separation witness).",
         False, DARK),
        ("The bubbles must be OPTIMAL or they are phantom: our exact banded DP is provably optimal (poasta was suboptimal — GSTM 1181<1331 — creating phantom PSVs).",
         False, DARK),
    ], 0.6, 1.5, 12.2, 1.9, size=14.5)
    pic_fit(s, assets["threading"], 0.4, 3.35, 7.3, 3.9)
    # real RABL2 strip + caption
    pic_fit(s, assets["rabl2_strip"], 7.5, 3.45, 5.7, 2.2)
    textbox(s, "REAL gorilla RABL2 (5 copies, cross-chrom): reads threaded, coloured by assigned copy",
            7.5, 5.75, 5.7, 0.5, 11.5, italic=True, align=PP_ALIGN.CENTER, color=(90, 90, 90))
    textbox(s, "235 bubbles (195 SUN) · K=5 fully resolvable · 222 reads → 191 assigned, 4 ambiguous, 27 no-cover, 0 tied",
            7.5, 6.25, 5.7, 0.8, 11.5, bold=True, align=PP_ALIGN.CENTER, color=NAVY)
    notes(s, 2)

    # ---- 3 REAL OR ERROR (the 5-card panel is the slide)
    s = prs.slides.add_slide(blank)
    title_bar(s, "Is the PSV real, or an error / artifact?", "QUESTION 3")
    textbox(s, "Every distinguishing PSV must survive five independent checks — each a named test or a theorem boundary, not a tuned filter:",
            0.6, 1.5, 12.2, 0.6, 15, bold=True, color=DARK)
    pic_fit(s, assets["realness"], 0.4, 2.2, 12.5, 4.2)
    textbox(s, "The load-bearing output is the min_p = 1.0 certificate that NAMES the unresolvable reads (30,620 of 30,806 tied rows) instead of guessing them.",
            0.6, 6.55, 12.2, 0.6, 13, italic=True, color=RDARK)
    notes(s, 3)

    # ---- 4 PROOF
    s = prs.slides.add_slide(blank)
    title_bar(s, "Proof it is real — RNA-internal vs DNA-independent, kept separate", "PROOF")
    bullets(s, [
        ("RNA-internal, planted (labels non-circular): identical copies → 0 invented PSVs, 120/120 certified TIED; multimapping regime 579/580 = 99.8% to the correct copy.",
         True, GDARK),
        ("Held-out PSVs (non-circular in the columns): assign on half a read's PSVs, the UNUSED half re-picks the same copy 80% — up to 6.4× the 1/K chance.",
         False, DARK),
        ("DNA-independent (non-circular in the molecule): held-out DNA columns confirm 97.2% (2.2×; K≥3 3.9×); 62% of families are ≥98% DNA segmental duplications (SEDEF).",
         True, (30, 80, 160)),
        ("Real annotated families recovered exactly: GSTM 3, MAGEA 2, DAZ 2, RBMY 6, TSPY 5, PCDHB 5; single-copy controls EEF1A1 / SRGAP2 → 0.",
         True, NAVY),
    ], 0.6, 1.55, 7.0, 4.6, size=15)
    pic_fit(s, assets["sim"], 7.55, 1.7, 5.6, 3.0)
    textbox(s, "Two non-circular axes:  RNA-internal (planted read-of-origin; held-out columns the call never saw)  ·  DNA-independent (a different molecule).  Neither uses minimap2-primary 'silver'.",
            7.5, 4.9, 5.75, 1.6, 12.5, italic=True, align=PP_ALIGN.CENTER, color=(90, 90, 90))
    notes(s, 4)

    # ---- 5 EXTERNAL ANCHOR: SOTO et al. 2025
    s = prs.slides.add_slide(blank)
    title_bar(s, "External anchor — concordance with an independent DNA catalog", "SOTO et al. 2025, Cell")
    bullets(s, [
        ("Independent, cross-lab, cross-modal: Soto et al. (Cell 2025, Dennis lab) cataloged human gene-family expansions from DNA / T2T long-read assemblies — different lab, different molecule, no shared aligner or silver standard.",
         True, NAVY),
        ("SPECIFICITY — on Soto's 13 human-specific families, our RNA method calls single-copy / absent in gorilla, on EXPRESSED loci (4–579 reads): 13/13 concordant, 0 fabricated expansions.",
         True, GDARK),
        ("RECOVERY — ancestral gorilla families resolved at concordant copy number; DAZ / RBMY / TSPY are genuine SD98 (≥98%) segdups — the exact objects Soto's pipeline clusters.",
         True, (30, 80, 160)),
        ("Honest scope: his per-family gorilla CN is not public → this is gene-set concordance, not a CN correlation; his families are human-specific by construction, so the correct test is specificity + recovery — both pass.",
         False, DARK),
    ], 0.6, 1.5, 12.3, 2.55, size=13.5)
    pic_fit(s, assets["soto"], 0.45, 4.15, 12.4, 3.05)
    s.notes_slide.notes_text_frame.text = (
        "Soto DC et al., Cell 188(19):5363-5383 (2025), doi 10.1016/j.cell.2025.06.037; bioRxiv 10.1101/2024.09.26.615256; "
        "data mydennislab/HSD_brain_evolution, Zenodo 10.5281/zenodo.15486469 (retrieved via PubMed). "
        "SPECIFICITY (soto_specificity.sh on GGO_mm.bam): 9 present human-specific loci return 0 multi-copy families with reads "
        "present (SRGAP2 248, ARHGAP11A 310, FRMPD2 363, GPR89A 70, GPRIN2 8, HYDIN 187, ROCK1 579, CD8B 4, DUSP22 277 primary "
        "reads); NPY4R/CFC1/NOTCH2NL absent in GGO_genomic.gff. RECOVERY chi(H) from reg_*_dp this session. SD98 = max SEDEF "
        "final.bed fracMatch (col 21) of a segdup overlapping the locus: DAZ 99.63, RBMY 99.74, TSPY 99.45 (>=98, Soto clusters); "
        "GSTM 95.29, MAGEA 94.85, PCDHB 88.86 (below SD98, more divergent). Soto's per-family gorilla CN is on the authors' "
        "internal share, not public, so this is gene-set concordance, not a CN correlation. Full record: bench/SOTO_CONCORDANCE.md.")

    # ---- 6 CROSS-SPECIES: NOT OVERFIT (human)
    s = prs.slides.add_slide(blank)
    title_bar(s, "Not overfit — the identical method on HUMAN data", "CROSS-SPECIES")
    bullets(s, [
        ("The 'overfit to gorilla / to this sample?' objection, answered directly: the IDENTICAL binary + recipe (-N 50 --secondary=yes) on a HUMAN testis Iso-Seq (ERR13885926, GENCODE, Sequel II HiFi) → T2T-CHM13 — different lab, individual, species.",
         True, NAVY),
        ("It tracks SPECIES-SPECIFIC copy number: MAGEA 2→11, TSPY 5→33 — the REAL human expansions (≈ the human annotation), not the gorilla numbers; RBMY concordant (6 = 6). Overfitting could not produce this.",
         True, (30, 80, 160)),
        ("Every recovered copy lands on an annotated human paralog (MAGEA1–12, RBMY1B/D/E/J/F, the TSPY array); MAGEA even split the adjacent CSAG family off correctly — real families, not mis-chains.",
         False, DARK),
        ("Honest: PCDHB / DAZ coverage-limited in this one library (77 / 16 reads); the near-identical human-specific duplicates (SRGAP2C, ARHGAP11B) sit at the K=0 frontier — the same limits we already disclose, not new ones.",
         False, DARK),
    ], 0.6, 1.5, 12.3, 2.55, size=13.5)
    pic_fit(s, assets["human"], 0.5, 4.15, 12.3, 3.05)
    s.notes_slide.notes_text_frame.text = (
        "Human = ERR13885926 (ENA), human testis full-length cDNA, GENCODE, PacBio Sequel II HiFi; 1,233,001 reads, median 888bp. "
        "Reference T2T-CHM13v2.0 (chm13v2.0.fa, complete chrY). Aligned minimap2 -ax splice:hq --eqx -Y -N 50 -p 0.1 --secondary=yes "
        "-> 96% mapped (1,179,288 primary; 7,942,153 secondary). copy_assign --min-copies 2 --skip-poa-diagnostic --homology-primary "
        "(same flags as gorilla), foreground/serial (human_families.sh). Human chi(H): RBMY 6 (RBMY1B/A1/D/E/J/F), TSPY 33 (~35 array, "
        "TSPY2/3/4/8/9/10 + LOCs), MAGEA 11 (MAGEA1-12; CSAG split off), GSTM 2 (GSTM2/5 expressed). PCDHB/DAZ coverage-limited "
        "(77/16 reads). Copy coords verified to overlap annotated genes. Tightest control (not yet run): Makova PRJNA911852 has matched "
        "human+gorilla testis Iso-Seq (SRR22838397.. / SRR22838403..) but as Sequel subreads (need CCS). Full record: bench/HUMAN_CROSSSPECIES.md.")

    # ---- 7 LIMITS
    s = prs.slides.add_slide(blank)
    title_bar(s, "Honest scope, the identifiability floor, and one line", "LIMITS")
    bullets(s, [
        ("The K=0 floor is STRUCTURAL, not a threshold: exonically-identical copies are genuinely RNA-unresolvable → planted identical copies 120/120 TIED; ~10% of families genome-wide carry zero bubble.",
         True, RDARK),
        ("No external true-origin accuracy on real data: 'resolvable' = 'threads a distinguishing bubble', not 'verified correct'. Truth-anchored legs = planted sim + machine-checked theory only.",
         False, DARK),
        ("Near-identical frontier disclosed, not hidden: a 5-copy collapse (div 0.001) over-merges to 4; clean recovery is a K≤2 phenomenon.",
         False, DARK),
        ("Named constants (relocated, not eliminated): α=1e-3 · refine id≥0.80 & cov≥0.50 · editing down-weight 0.2. The VG is an ASSIGNMENT object, not a read aligner (vg giraffe 94.5% < 98%).",
         False, DARK),
        ("ONE LINE: each family = read-threaded paths through one PSV-bubble graph — backbone collapses, differences become bubbles, assignment is assign-or-abstain threading with a per-read certificate, over provable structure (χ(H); facility-location) — no 1/k.",
         True, NAVY),
    ], 0.6, 1.55, 8.1, 5.0, size=14)
    pic_fit(s, assets["rgpd8"], 8.9, 1.7, 4.2, 3.0)
    textbox(s, "RGPD8: 7 copies, 0 bubbles → K=1.\nThe K=0 floor made visible: nothing to\nthread → all 2075 reads abstain (no-cover).",
            8.85, 4.75, 4.3, 1.3, 12, italic=True, align=PP_ALIGN.CENTER, color=(90, 90, 90))
    notes(s, 5)

    path = os.path.join(OUT, "vg_copies_as_paths.pptx")
    try:
        prs.save(path)
    except PermissionError:
        path = os.path.join(OUT, "vg_copies_as_paths.new.pptx")  # original open in PowerPoint -> write beside it
        prs.save(path)
    return path


def main():
    # custom schematics
    assets = {"flip": fig_flip(), "threading": fig_threading(), "realness": fig_realness(), "sim": fig_sim()}
    # real figures (cropped to the clean, consistent portion)
    assets["rabl2_strip"] = crop_top(os.path.join(HERE, "fig_o2_vg_fam39_RABL2.png"),
                                     os.path.join(OUT, "vg_rabl2_strip.png"), frac=0.365)
    # RGPD8 graph (0 bubbles) — use the graph-object view
    assets["rgpd8"] = os.path.join(HERE, "fig_o2_vg_fam1_RGPD8_graph.png")
    # Soto et al. 2025 concordance panel (generate if not already present)
    soto = os.path.join(OUT, "soto_concordance.png")
    if not os.path.exists(soto):
        import make_soto_concordance
        make_soto_concordance.build()
    assets["soto"] = soto
    # human cross-species panel (generate if not already present)
    human = os.path.join(OUT, "human_crossspecies.png")
    if not os.path.exists(human):
        import make_human_crossspecies
        make_human_crossspecies.build()
    assets["human"] = human
    path = build(assets)
    print(f"wrote {path}")
    for k, v in assets.items():
        print(f"  {k}: {v}")


if __name__ == "__main__":
    main()
