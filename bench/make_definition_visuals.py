#!/usr/bin/env python3
"""
make_definition_visuals.py — picture-first figures of the family definition on a REAL
family (the RFPL4A tandem array: 5 copies on NC_073244.2), from real reads.

  fig1_reads.png   : the 5 copies + real read coverage + example multimapper reads (arcs)
                     -> reads that fit several copies equally well can't be told apart
  fig2_exonunion.png: real isoform reads at one copy -> their exon-UNION = the copy model
  fig3_graph.png   : the ~R graph (ties) confirmed by ~B (shared backbone) = a FAMILY;
                     contrasted with OCLN~SEPTIN7 (~R yes, ~B no -> bridge cut)
Run: /home/juanfra/miniforge3/bin/python bench/make_definition_visuals.py
"""
import collections
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle
import numpy as np
import pysam

HERE = os.path.dirname(os.path.abspath(__file__))
BAM = "/home/juanfra/winloci_scratch/GGO.bam"
NAVY, TEAL, ORANGE, GREY = "#1F2D5A", "#127A6E", "#D96B27", "#9aa0a6"
CHROM = "NC_073244.2"
COPIES = [   # (name, start, end)  the RFPL4A tandem array
    ("RFPL4A", 67776054, 67779514),
    ("copy-2", 67785173, 67788657),
    ("copy-3", 67794315, 67797800),
    ("copy-4", 67803458, 67806943),
    ("copy-5", 67812601, 67816085),
]
# ~R tie counts (measured, from family_def_dna_pr_edges.tsv) — indices into COPIES
TIES = {(0,1):8,(0,2):5,(0,3):5,(0,4):5,(1,2):23,(1,3):19,(1,4):20,(2,3):18,(2,4):19,(3,4):16}
REG0, REG1 = 67774000, 67818000


def copy_of(pos):
    for i, (_, s, e) in enumerate(COPIES):
        if s <= pos <= e:
            return i
    return None


# ----------------------------------------------------------------- fig 1: reads
def fig1():
    bam = pysam.AlignmentFile(BAM, "rb")
    cov = np.zeros(REG1 - REG0)
    # per read name -> set of copies it aligns to (primary + secondary)
    read_copies = collections.defaultdict(set)
    examples = []   # (name, [(copy_idx, mid_local), ...]) for a few multimappers
    for r in bam.fetch(CHROM, REG0, REG1):
        if r.is_unmapped:
            continue
        for s, e in r.get_blocks():
            a, b = max(s, REG0), min(e, REG1)
            if b > a:
                cov[a - REG0:b - REG0] += 1
        ci = copy_of((r.reference_start + r.reference_end) // 2)
        if ci is not None and not r.is_supplementary:
            read_copies[r.query_name].add(ci)
    bam.close()
    mm = {n: cs for n, cs in read_copies.items() if len(cs) >= 2}

    fig, ax = plt.subplots(figsize=(13, 4.6))
    x = np.arange(REG0, REG1) / 1000.0
    ax.fill_between(x, 0, cov, color=TEAL, alpha=0.35, lw=0)
    ax.plot(x, cov, color=TEAL, lw=0.6)
    ymax = cov.max() * 1.15
    for i, (nm, s, e) in enumerate(COPIES):
        ax.add_patch(Rectangle((s / 1000, -ymax * 0.16), (e - s) / 1000, ymax * 0.10,
                     facecolor=NAVY, edgecolor="none"))
        ax.text((s + e) / 2000, -ymax * 0.255, f"copy {i+1}", ha="center", fontsize=9,
                weight="bold", color=NAVY)
    # example multimapper reads as arcs between the copies they hit
    mids = [(s + e) / 2000 for _, s, e in COPIES]
    shown = 0
    for nm, cs in mm.items():
        cs = sorted(cs)
        for a, b in zip(cs, cs[1:]):
            xa, xb = mids[a], mids[b]
            xm = (xa + xb) / 2; h = ymax * (0.45 + 0.12 * shown)
            ax.plot([xa, xm, xb], [ymax * 0.30, h, ymax * 0.30], color=ORANGE, lw=1.1, alpha=0.8)
        shown += 1
        if shown >= 6:
            break
    ax.text((REG0 + REG1) / 2000, ymax * 1.0,
            f"orange = the same read fitting two copies equally well  ({len(mm)} such multimapping reads here)",
            ha="center", fontsize=9.5, color=ORANGE, style="italic")
    ax.set_xlim(REG0 / 1000, REG1 / 1000); ax.set_ylim(-ymax * 0.30, ymax * 1.1)
    ax.set_yticks([]); ax.set_xlabel(f"{CHROM} position (kb)")
    ax.set_title("5 near-identical copies of one gene, each expressed (read coverage) — "
                 "and reads that can't tell them apart", fontsize=12, weight="bold", color=NAVY)
    for sp in ["top", "right", "left"]:
        ax.spines[sp].set_visible(False)
    fig.tight_layout(); fig.savefig(os.path.join(HERE, "fig1_reads.png"), dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- fig 2: exon-union
def fig2():
    bam = pysam.AlignmentFile(BAM, "rb")
    nm, s, e = COPIES[1]            # copy-2 (a LOC copy, lots of reads)
    reads = []
    for r in bam.fetch(CHROM, s - 2000, e + 2000):
        if r.is_unmapped or r.is_supplementary:
            continue
        b = [(bs, be) for bs, be in r.get_blocks() if be > s - 2000 and bs < e + 2000]
        if b:
            reads.append(b)
        if len(reads) >= 14:
            break
    bam.close()
    # exon-union
    allb = sorted([iv for b in reads for iv in b])
    union = []
    for a, c in allb:
        if union and a <= union[-1][1]:
            union[-1] = (union[-1][0], max(union[-1][1], c))
        else:
            union.append((a, c))

    fig, ax = plt.subplots(figsize=(11.5, 5.2))
    x0, x1 = s - 1500, e + 1500
    for row, b in enumerate(reads):
        y = len(reads) - row
        ax.plot([x0 / 1000, x1 / 1000], [y, y], color="#ddd", lw=0.5, zorder=0)
        for bs, be in b:
            ax.add_patch(Rectangle((bs / 1000, y - 0.28), (be - bs) / 1000, 0.56,
                         facecolor=GREY, edgecolor="none"))
    ax.text(x0 / 1000, len(reads) + 1.1, "isoform reads at one copy (each row = one read; blocks = its exons)",
            fontsize=10, color=GREY)
    for bs, be in union:
        ax.add_patch(Rectangle((bs / 1000, -0.9), (be - bs) / 1000, 0.7,
                     facecolor=ORANGE, edgecolor="none"))
        if be - bs > 200:
            ax.plot([(be) / 1000, (be) / 1000], [-0.9, -0.2], color="#fff", lw=0)
    # connect union exons with thin intron line
    for (a1, b1), (a2, b2) in zip(union, union[1:]):
        ax.plot([b1 / 1000, a2 / 1000], [-0.55, -0.55], color=ORANGE, lw=0.8)
    ax.text(x0 / 1000, -1.5, "exon-UNION = the copy model  (sum the exons of all isoforms → 'the full gene minus introns')",
            fontsize=10.5, color=ORANGE, weight="bold")
    ax.set_xlim(x0 / 1000, x1 / 1000); ax.set_ylim(-2.2, len(reads) + 2)
    ax.set_yticks([]); ax.set_xlabel(f"{CHROM} position (kb)")
    ax.set_title("Summing exons: many isoforms of one copy → one exon-union",
                 fontsize=12.5, weight="bold", color=NAVY)
    for sp in ["top", "right", "left"]:
        ax.spines[sp].set_visible(False)
    fig.tight_layout(); fig.savefig(os.path.join(HERE, "fig2_exonunion.png"), dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- fig 3: the graph + ~B
def fig3():
    fig, ax = plt.subplots(1, 2, figsize=(13, 5.4))
    # left: RFPL family — 5 nodes in a ring, ~R edges (ties), all ~B-confirmed
    n = 5
    ang = np.linspace(90, 90 + 360, n, endpoint=False) * np.pi / 180
    pos = [(np.cos(a), np.sin(a)) for a in ang]
    for (i, j), w in TIES.items():
        (xi, yi), (xj, yj) = pos[i], pos[j]
        ax[0].plot([xi, xj], [yi, yj], color=TEAL, lw=0.6 + w / 8, alpha=0.55, zorder=1)
        ax[0].text((xi + xj) / 2, (yi + yj) / 2, str(w), fontsize=7.5, color=TEAL,
                   ha="center", va="center")
    for i, (x, y) in enumerate(pos):
        ax[0].scatter([x], [y], s=900, color=NAVY, zorder=3)
        ax[0].text(x, y, f"copy\n{i+1}", color="white", ha="center", va="center",
                   fontsize=8, weight="bold", zorder=4)
    ax[0].set_title("RFPL4A array: every pair is read-coupled (edge = #tied reads)\n"
                    "AND every copy shares the backbone → ONE FAMILY (5 copies)",
                    fontsize=11, weight="bold", color=NAVY)
    ax[0].set_xlim(-1.5, 1.5); ax[0].set_ylim(-1.5, 1.6); ax[0].axis("off")
    ax[0].text(0, -1.5, "~R (reads tie them)  ✓   +   ~B (backbone ≥ 0.95)  ✓", ha="center",
               fontsize=10, color=TEAL, weight="bold")

    # right: the bridge — OCLN ~ SEPTIN7: ~R yes, ~B no -> cut
    p = {"OCLN": (-0.8, 0), "SEPTIN7": (0.8, 0)}
    ax[1].plot([-0.8, 0.8], [0, 0], color=ORANGE, lw=3, ls=(0, (4, 4)), zorder=1)
    ax[1].text(0, 0.18, "3,369 reads tie  → ~R edge", ha="center", fontsize=9, color=TEAL)
    ax[1].text(0, -0.22, "but backbones align only 0.05  → ~B FAILS", ha="center", fontsize=9, color=ORANGE)
    ax[1].scatter([-0.8], [0], s=1500, color=NAVY); ax[1].scatter([0.8], [0], s=1500, color=NAVY)
    ax[1].text(-0.8, 0, "OCLN", color="white", ha="center", va="center", fontsize=8.5, weight="bold")
    ax[1].text(0.8, 0, "SEPTIN7", color="white", ha="center", va="center", fontsize=8, weight="bold")
    ax[1].scatter([0], [0], s=420, marker="X", color="#c00", zorder=5)
    ax[1].set_title("A bridge: reads tie them (~R) but they share NO backbone (~B)\n"
                    "→ the edge is cut → NOT a family", fontsize=11, weight="bold", color=NAVY)
    ax[1].set_xlim(-1.6, 1.6); ax[1].set_ylim(-1, 1); ax[1].axis("off")
    ax[1].text(0, -0.7, "the two relations must BOTH hold", ha="center", fontsize=10,
               color=ORANGE, weight="bold")
    fig.tight_layout(); fig.savefig(os.path.join(HERE, "fig3_graph.png"), dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- fig 9: results table
def fig9():
    fig, ax = plt.subplots(figsize=(13, 6.6)); ax.axis("off"); ax.set_xlim(0, 12); ax.set_ylim(0, 13)
    ax.text(6, 12.5, "What the method finds — precision / recall on the 17-candidate panel",
            ha="center", fontsize=14, weight="bold", color=NAVY)
    # FOUND table
    ax.text(0.2, 11.6, "FAMILIES FOUND  (true positives)", fontsize=11.5, weight="bold", color=TEAL)
    found = [
        ("RABL2  (RABL2A ~ RABL2B)", "recent paralog, separate chromosomes", "47"),
        ("AK6  (~ LOC copy)", "high-identity cross-chromosome copy", "24"),
        ("CCDC196  (~ LOC copy)", "high-identity cross-chromosome copy", "24"),
        ("MAGEA array  (4 co-located copies)", "tandem array of de-novo loci", "75–303"),
    ]
    ax.add_patch(Rectangle((0.2, 10.7), 11.6, 0.55, facecolor=NAVY))
    for x, t in [(0.4, "gene family"), (5.6, "type"), (10.2, "tied reads (~R)")]:
        ax.text(x, 10.97, t, color="white", fontsize=9.5, weight="bold", va="center")
    for i, (fam, typ, ev) in enumerate(found):
        y = 10.15 - i * 0.62
        ax.add_patch(Rectangle((0.2, y - 0.27), 11.6, 0.54, facecolor="#eef5f4" if i % 2 else "#fff",
                     edgecolor="#dde"))
        ax.text(0.4, y, fam, fontsize=9.5, va="center", color=NAVY, weight="bold")
        ax.text(5.6, y, typ, fontsize=9, va="center", color="#333")
        ax.text(10.6, y, ev, fontsize=9.5, va="center", color=TEAL, weight="bold")
    # headline P/R
    ax.add_patch(FancyBboxPatch((0.2, 6.0), 11.6, 1.0, boxstyle="round,pad=0.05",
                 facecolor="#eef5f4", edgecolor=TEAL, linewidth=1.5))
    ax.text(6, 6.5, "Precision = 7 / 7 = 1.00  (0 false families)        Recall = 7 / 7 = 1.00  (0 real families missed)",
            ha="center", fontsize=13, weight="bold", color=NAVY)
    # REJECTED summary
    ax.text(0.2, 5.3, "CORRECTLY REJECTED  (true negatives — neither relation, or read-resolvable)",
            fontsize=11, weight="bold", color=ORANGE)
    rej = [
        "APOBEC3, RFPL — recent paralogs, but the reads place uniquely (read-resolvable → not ~R)",
        "EEF1A1, CNN2 — processed-pseudogene retrocopies (0–1 ties, no quorum)",
        "CREB1~METTL21A, GCA~KCNH7, CASP8~FLACC1, ASDURF~ASNSD1, GPR39~LYPD1 — domain-sharers (0 ties → not ~R)",
        "MAGEA annotated genes — resolvable",
    ]
    for i, r in enumerate(rej):
        ax.text(0.5, 4.7 - i * 0.55, "•  " + r, fontsize=9.5, color="#333", va="center")
    # honest genome-wide / DNA footnote
    ax.add_patch(Rectangle((0.2, 0.2), 11.6, 1.7, facecolor="#f6f6f6", edgecolor="#ddd"))
    ax.text(0.4, 1.55, "Beyond the panel (honest context):", fontsize=10, weight="bold", color=GREY)
    ax.text(0.4, 1.05,
            "• Genome-wide: 196 validated families (de-novo loci).",
            fontsize=9.5, color="#333")
    ax.text(0.4, 0.55,
            "• vs an independent DNA-homology truth: precision 0.64 (0.72 crediting real homology just under the bar); recall is "
            "expression-limited — the two\n   definitions measure different things (read-confusability vs sequence homology), so this is an "
            "orthogonal comparison, not an error rate.",
            fontsize=9.0, color="#333")
    fig.tight_layout(); fig.savefig(os.path.join(HERE, "fig9_results.png"), dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- fig 6: alignment algorithm
def fig6():
    rng = np.random.default_rng(1)
    fig, ax = plt.subplots(1, 3, figsize=(13.5, 4.0))
    # 1. SEED: minimizers shared between read and reference
    ax[0].set_title("1. SEED\nshared minimizers (k=15, w=5)", fontsize=11, weight="bold", color=NAVY)
    ax[0].plot([0, 10], [2, 2], color=NAVY, lw=4); ax[0].plot([0, 10], [0, 0], color=GREY, lw=4)
    ax[0].text(-0.3, 2, "read", ha="right", va="center", fontsize=9, color=NAVY)
    ax[0].text(-0.3, 0, "copy", ha="right", va="center", fontsize=9, color=GREY)
    for p in [1.2, 2.6, 4.1, 5.0, 6.8, 8.3]:
        ax[0].plot([p, p + rng.uniform(-0.15, 0.15)], [2, 0], color=TEAL, lw=1.3, alpha=0.8)
        ax[0].scatter([p], [2], s=22, color=TEAL); ax[0].scatter([p], [0], s=22, color=TEAL)
    ax[0].set_xlim(-1.5, 11); ax[0].set_ylim(-1.2, 3.2); ax[0].axis("off")
    # 2. CHAIN: collinear anchors
    ax[1].set_title("2. CHAIN\ncollinear anchors (DP)", fontsize=11, weight="bold", color=NAVY)
    xs = np.array([1.2, 2.6, 4.1, 5.0, 6.8, 8.3]); ys = xs * 0.9 + rng.uniform(-0.15, 0.15, len(xs))
    ax[1].plot(xs, ys, "-", color=TEAL, lw=1.0, alpha=0.6)
    ax[1].scatter(xs, ys, s=40, color=NAVY, zorder=3)
    ax[1].scatter([3.5, 6.0], [7.5, 1.2], s=30, color="#ccc")  # off-diagonal noise anchors
    ax[1].set_xlabel("read position"); ax[1].set_ylabel("copy position")
    ax[1].set_xlim(0, 9.5); ax[1].set_ylim(0, 9.5)
    # 3. ALIGN: banded base-level DP fills the gaps between anchors
    ax[2].set_title("3. ALIGN\nbanded affine-gap DP (only the gaps)", fontsize=11, weight="bold", color=NAVY)
    ax[2].add_patch(Rectangle((0, 0), 8, 8, facecolor="#f3f3f3", edgecolor="none"))
    ax[2].plot([0.5, 7.5], [0.5, 7.5], color=NAVY, lw=2)
    for off in (-1, 1):
        ax[2].plot([0.5, 7.5], [0.5 + off, 7.5 + off], color=TEAL, lw=0.7, ls="--", alpha=0.6)
    ax[2].text(4, 6.6, "band = chain ± a few bp\n(not the full N×M matrix)", fontsize=8.5,
               color=NAVY, ha="center")
    ax[2].set_xlabel("read"); ax[2].set_ylabel("copy"); ax[2].set_xlim(0, 8); ax[2].set_ylim(0, 8)
    ax[2].set_xticks([]); ax[2].set_yticks([])
    fig.suptitle("How the alignment is computed — minimap2 'seed–chain–align', NOT a full Needleman–Wunsch matrix",
                 fontsize=12.5, weight="bold", color=NAVY, y=1.04)
    fig.text(0.5, -0.10,
             "scoring (splice:hq preset):  match +1 · mismatch −4 · gap two-piece affine (open −6/−24, extend −1/0) · "
             "non-canonical splice −5 · intron ≤ 200 kb\n"
             "→ mismatch ≫ match, so the aligner minimises differences; two-piece gaps make long introns cheap; "
             "canonical GT–AG splice sites are preferred.  The only number WE add downstream is Δ.",
             ha="center", fontsize=9.5, color="#333")
    fig.tight_layout(); fig.savefig(os.path.join(HERE, "fig6_algo.png"), dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- fig 7: error vs recurrence
def fig7():
    rng = np.random.default_rng(2)
    fig, ax = plt.subplots(figsize=(12.5, 5.2))
    nreads, L = 11, 40
    err_col, psv_col = 12, 27
    for row in range(nreads):
        y = nreads - row
        ax.add_patch(Rectangle((0, y - 0.32), L, 0.64, facecolor="#eef", edgecolor="#ccd"))
    # one-off error: a single read differs at err_col
    ax.add_patch(Rectangle((err_col, nreads - 4 - 0.32), 1, 0.64, facecolor="#c00"))
    # PSV: a consistent subset (copies) differ at psv_col
    for row in [1, 3, 5, 8]:
        ax.add_patch(Rectangle((psv_col, nreads - row - 0.32), 1, 0.64, facecolor=ORANGE))
    ax.annotate("1 read only, random position\n→ SEQUENCING ERROR → ignored",
                xy=(err_col + 0.5, nreads - 4), xytext=(err_col - 9, nreads + 1.6),
                fontsize=10, color="#c00", weight="bold",
                arrowprops=dict(arrowstyle="->", color="#c00"))
    ax.annotate("same difference in MANY reads, same position\n→ a REAL variant (a PSV)",
                xy=(psv_col + 0.5, nreads - 5), xytext=(psv_col - 4, -2.3),
                fontsize=10, color=ORANGE, weight="bold",
                arrowprops=dict(arrowstyle="->", color=ORANGE))
    ax.text(L / 2, nreads + 3.0, "reads stacked at one position (each row = one read; column = one base)",
            ha="center", fontsize=10, color=GREY)
    ax.set_xlim(-10, L + 1); ax.set_ylim(-3.2, nreads + 4); ax.axis("off")
    ax.set_title("Is a difference an ERROR or REAL?  —  decided by RECURRENCE across reads, never one read",
                 fontsize=12.5, weight="bold", color=NAVY)
    fig.text(0.5, -0.04,
             "HiFi error ≈ 0.1–0.5 % · de is GAP-COMPRESSED (a homopolymer indel — the main HiFi error — counts once) · "
             "an edge needs a QUORUM (k ≥ 3 tied reads) · Δ = 0.005 ≈ 4σ of single-read divergence noise → "
             "one read's errors cannot make a false tie.", ha="center", fontsize=9.5, color="#333")
    fig.tight_layout(); fig.savefig(os.path.join(HERE, "fig7_error.png"), dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- fig 8: CIGAR -> exons
def fig8():
    fig, ax = plt.subplots(figsize=(12.5, 4.4))
    # real read: 79M 1755N 295M 509N 823M  at copy-2 (positions in kb, shifted to 0)
    base = 67785198
    exons = [(67785198, 67785277, "79M"), (67787032, 67787327, "295M"), (67787836, 67788659, "823M")]
    introns = [(67785277, 67787032, "1755N"), (67787327, 67787836, "509N")]
    def X(p): return (p - base) / 1000.0
    y = 1.0
    for (s, e, lab) in exons:
        ax.add_patch(Rectangle((X(s), y - 0.22), X(e) - X(s), 0.44, facecolor=NAVY, edgecolor="none"))
        ax.text((X(s) + X(e)) / 2, y + 0.42, lab, ha="center", fontsize=9, color=NAVY, weight="bold")
        ax.text((X(s) + X(e)) / 2, y, "exon", ha="center", va="center", color="white", fontsize=8)
    for (s, e, lab) in introns:
        ax.plot([X(s), X(e)], [y, y], color=GREY, lw=1.0)
        ax.text((X(s) + X(e)) / 2, y - 0.30, f"{lab}  (intron, skipped)", ha="center", fontsize=8.5, color=GREY)
    ax.text(X(base) - 0.15, y, "one read →", ha="right", va="center", fontsize=10, weight="bold", color=NAVY)
    ax.text(X(base), y + 1.05,
            "the read's CIGAR is its blueprint:   M = aligned to genome = EXON     N = skipped = INTRON",
            fontsize=10.5, color=NAVY, weight="bold")
    ax.text(X(base), -0.15,
            "exons = the M-blocks.  A splice junction (an N boundary) is kept only if MANY reads place it identically "
            "(recurrence) AND it is canonical GT–AG;", fontsize=9.5, color="#333")
    ax.text(X(base), -0.55,
            "a one-read junction is dropped as mis-splicing.   Then the copy model = UNION of all reads' exon blocks (prev. slide).",
            fontsize=9.5, color=ORANGE, weight="bold")
    ax.set_xlim(X(base) - 0.6, X(67788659) + 0.3); ax.set_ylim(-0.9, 2.4); ax.axis("off")
    ax.set_title("How a read is digested into exons (and how junctions are trusted)",
                 fontsize=12.5, weight="bold", color=NAVY)
    fig.tight_layout(); fig.savefig(os.path.join(HERE, "fig8_cigar.png"), dpi=150, bbox_inches="tight")
    plt.close(fig)


REPCDNA = "/home/juanfra/winloci_scratch/rep_cdna.fa"


def load_seqs(names):
    want = set(names); seqs = {}; cur = None; buf = []
    for line in open(REPCDNA):
        if line.startswith(">"):
            if cur in want:
                seqs[cur] = "".join(buf)
            cur = line[1:].strip(); buf = []
        elif cur in want:
            buf.append(line.strip())
    if cur in want:
        seqs[cur] = "".join(buf)
    return seqs


# ----------------------------------------------------------------- fig 4: read alignment -> de
def fig4():
    rng = np.random.default_rng(0)
    fig, ax = plt.subplots(figsize=(12.5, 4.6))
    L = 5269
    # the read bar (data coords)
    ax.add_patch(Rectangle((0, 2.58), L, 0.42, facecolor=NAVY, edgecolor="none"))
    ax.text(L / 2, 2.79, "one HiFi read  (5,269 aligned bp)", ha="center", va="center",
            color="white", fontsize=11.5, weight="bold")
    ax.annotate("", xy=(L * 0.30, 1.70), xytext=(L * 0.42, 2.55),
                arrowprops=dict(arrowstyle="-|>", color=TEAL, lw=1.5))
    ax.annotate("", xy=(L * 0.70, 0.70), xytext=(L * 0.58, 2.55),
                arrowprops=dict(arrowstyle="-|>", color=ORANGE, lw=1.5))
    ax.text(L * 0.5, 2.30, "align to each copy (minimap2)", ha="center",
            fontsize=9.5, style="italic", color="#666")
    for k, (nm, de, col) in enumerate([("copy 3", 0.0187, TEAL), ("copy 4", 0.0180, ORANGE)]):
        y = 1.5 - k * 1.0
        ax.add_patch(Rectangle((0, y - 0.16), L, 0.32, facecolor="#eee", edgecolor="#bbb"))
        nmis = round(de * L)
        ax.vlines(rng.uniform(20, L - 20, nmis), y - 0.16, y + 0.16, color=col, lw=0.8, alpha=0.9)
        ax.text(-90, y, f"→ {nm}", ha="right", va="center", fontsize=11, weight="bold", color=col)
        ax.text(L + 80, y, f"de = {de:.4f}   (≈ {nmis} differences)", ha="left", va="center",
                fontsize=10, color=col)
    ax.text(L / 2, -0.20,
            "de = gap-compressed per-base divergence = (mismatches + indels) / aligned bases",
            ha="center", fontsize=10, style="italic", color="#333")
    ax.text(L / 2, -0.78,
            "|0.0187 − 0.0180| = 0.0007 ≤ Δ = 0.005   →   the read fits BOTH equally well   →   TIED   →   one vote for edge  copy3 — copy4",
            ha="center", fontsize=10.5, color=NAVY, weight="bold")
    ax.set_xlim(-750, L + 1400); ax.set_ylim(-1.1, 3.3); ax.axis("off")
    ax.set_title("How a read aligns → the number  de  → a 'tie'  (the input to ~R)",
                 fontsize=13, weight="bold", color=NAVY)
    fig.savefig(os.path.join(HERE, "fig4_align.png"), dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- fig 5: copy-copy dotplot -> ~B
def dotplot(ax, s1, s2, k=12, color=NAVY):
    idx = collections.defaultdict(list)
    for j in range(len(s2) - k):
        idx[s2[j:j + k]].append(j)
    xs, ys = [], []
    for i in range(0, len(s1) - k, 3):
        for j in idx.get(s1[i:i + k], ()):
            xs.append(i); ys.append(j)
    ax.scatter(xs, ys, s=1.2, color=color, alpha=0.5, rasterized=True)
    return len(set(xs))  # query positions with a match


def fig5():
    seqs = load_seqs(["LOC129528712", "LOC129528713", "OCLN", "SEPTIN7"])
    fig, ax = plt.subplots(1, 2, figsize=(13, 5.6))
    # copies: strong diagonal = colinear backbone over the whole length
    a, b = seqs["LOC129528712"], seqs["LOC129528713"]
    m = dotplot(ax[0], a, b, color=TEAL)
    cov = m / (len(a) / 3)
    ax[0].set_title(f"Two copies of the array (align them)\nstrong diagonal → shared backbone over ~{cov:.0%} of each  → ~B ✓",
                    fontsize=11.5, weight="bold", color=NAVY)
    ax[0].set_xlabel("copy A (bp)"); ax[0].set_ylabel("copy B (bp)")
    ax[0].set_xlim(0, len(a)); ax[0].set_ylim(0, len(b))
    # bridge: no diagonal -> no backbone
    o, sp = seqs["OCLN"], seqs["SEPTIN7"]
    dotplot(ax[1], o, sp, color=ORANGE)
    ax[1].set_title("OCLN vs SEPTIN7 (the bridge)\nno diagonal → no shared backbone (cov 0.05) → ~B ✗",
                    fontsize=11.5, weight="bold", color=ORANGE)
    ax[1].set_xlabel("OCLN (bp)"); ax[1].set_ylabel("SEPTIN7 (bp)")
    ax[1].set_xlim(0, len(o)); ax[1].set_ylim(0, len(sp))
    for a_ in ax:
        a_.grid(alpha=0.2)
    fig.suptitle("Aligning the COPIES (minimap2): a dot-plot — every dot is a shared 12-mer; a diagonal is a shared backbone",
                 fontsize=12.5, weight="bold", color=NAVY, y=1.02)
    fig.tight_layout(); fig.savefig(os.path.join(HERE, "fig5_dotplot.png"), dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- picture-first deck
def build_deck(out_name="family_definition_visual.pptx"):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from PIL import Image
    cNAVY, cTEAL, cORANGE, cGREY = RGBColor(0x1F,0x2D,0x5A), RGBColor(0x12,0x7A,0x6E), RGBColor(0xD9,0x6B,0x27), RGBColor(0x55,0x55,0x55)

    prs = Presentation(); prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

    def fig_slide(title, img, caption, tcolor=cNAVY):
        img = os.path.join(HERE, img)
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.7))
        tp = tb.text_frame.paragraphs[0]; tp.text = title
        tp.font.size = Pt(24); tp.font.bold = True; tp.font.color.rgb = tcolor
        w, h = Image.open(img).size
        scale = min(Inches(12.4) / w, Inches(5.35) / h)
        iw, ih = int(w*scale), int(h*scale)
        s.shapes.add_picture(img, int((Inches(13.33)-iw)/2), Inches(1.0), width=iw, height=ih)
        cb = s.shapes.add_textbox(Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.7))
        cp = cb.text_frame; cp.word_wrap = True; cp0 = cp.paragraphs[0]
        cp0.text = caption; cp0.font.size = Pt(15); cp0.font.color.rgb = cORANGE
        cp0.font.bold = True; cp0.alignment = PP_ALIGN.CENTER
        return s

    # title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bx = s.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(12), Inches(2.6)); tf = bx.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.text = "What is a multi-copy gene family?"
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = cNAVY
    p2 = tf.add_paragraph(); p2.text = "Loci whose reads can't be told apart  (~R)   AND   whose copies share a backbone  (~B)."
    p2.font.size = Pt(20); p2.font.color.rgb = cTEAL
    p3 = tf.add_paragraph(); p3.text = "Shown on one real family — the RFPL4A array, 5 tandem copies."
    p3.font.size = Pt(16); p3.font.color.rgb = cGREY

    fig_slide("Reads can't tell the copies apart  →  the relation ~R", "fig1_reads.png",
              "5 copies, all expressed; 32 reads fit several copies equally well (orange) — each such read links two copies.")
    fig_slide("How a read aligns: the number 'de', and what a 'tie' is", "fig4_align.png",
              "minimap2 aligns the read to each copy; de = differences per base. Two copies within Δ of each other = a tie = one ~R vote.")
    fig_slide("Sum the exons of all isoforms  →  one copy model", "fig2_exonunion.png",
              "A copy's many isoforms (grey) sum to one exon-union (orange) — so splicing can't masquerade as extra copies.")
    fig_slide("Aligning the copies: a dot-plot  →  backbone or not (~B)", "fig5_dotplot.png",
              "We run minimap2 to align the copy models; a diagonal = a shared backbone over most of both. Bridge: no diagonal → ~B fails.")
    fig_slide("Ties → edges → graph (~R), confirmed by the backbone (~B)  =  FAMILY", "fig3_graph.png",
              "Left: every pair ties AND shares a backbone → one 5-copy family.   Right: a bridge ties but shares no backbone → cut.")

    # definition (one line, big)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.8))
    tb.text_frame.paragraphs[0].text = "The definition"
    tb.text_frame.paragraphs[0].font.size = Pt(26); tb.text_frame.paragraphs[0].font.bold=True; tb.text_frame.paragraphs[0].font.color.rgb=cNAVY
    bx = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4)); tf = bx.text_frame; tf.word_wrap=True
    for i,(t,c,sz) in enumerate([
        ("A multi-copy gene family = a set of de-novo loci that is", cNAVY, 22),
        ("read-coupled (~R: ≥3 reads tie each link)   AND   backbone-connected (~B: copies align reciprocally).", cTEAL, 20),
        ("", cGREY, 10),
        ("Built only from reads:  align → loci → count ties → ~R graph → components → exon-union copies → align → ~B → family.", cGREY, 15),
    ]):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=t; p.font.size=Pt(sz); p.font.color.rgb=c; p.space_after=Pt(10)
        if i<2: p.font.bold=True

    # results table
    fig_slide("Results: gene families found  (precision / recall)", "fig9_results.png",
              "On the 17-candidate panel with hand-checked ground truth: precision = recall = 1.00 (0 false families, 0 missed).")

    # ---- appendix: "how it actually works" (for the literal questions) ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bx = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.5)); tf = bx.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.text = "Appendix — how it actually works"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = cNAVY
    p2 = tf.add_paragraph(); p2.text = "the alignment algorithm & scoring · error vs. real difference · digesting a read into exons"
    p2.font.size = Pt(17); p2.font.color.rgb = cGREY
    p3 = tf.add_paragraph()
    p3.text = "What is minimap2's vs ours:  minimap2 does the ALIGNMENT (reads → placements, de, CIGAR — we do not implement an aligner). Our method is everything downstream — the two relations, the graph, the definition."
    p3.font.size = Pt(13); p3.font.color.rgb = cORANGE; p3.font.bold = True
    fig_slide("How the alignment works (minimap2's job, not ours; not Needleman–Wunsch)", "fig6_algo.png",
              "minimap2 — not us — does seed–chain–align (minimizer seeds → chaining → banded affine-gap DP). We only consume its output.")
    fig_slide("Error or real? — by recurrence across reads", "fig7_error.png",
              "A difference in one read = error (ignored); the same difference in many reads = a real variant. Edges need a quorum (k≥3).")
    fig_slide("Digesting a read into exons (CIGAR → exons → union)", "fig8_cigar.png",
              "Each read's CIGAR marks exons (M) and introns (N); junctions are trusted by recurrence + canonical GT–AG; exons are then unioned.")

    try:
        prs.save(os.path.join(HERE, out_name))
        print(f"[+] wrote {out_name} ({len(prs.slides._sldIdLst)} slides)")
    except PermissionError:
        alt = out_name.replace(".pptx", "_FULL.pptx")
        prs.save(os.path.join(HERE, alt))
        print(f"[!] {out_name} was locked; wrote {alt} ({len(prs.slides._sldIdLst)} slides) instead")


if __name__ == "__main__":
    fig1(); fig2(); fig3(); fig4(); fig5(); fig6(); fig7(); fig8(); fig9()
    print("[+] wrote fig1..fig9")
    build_deck()
