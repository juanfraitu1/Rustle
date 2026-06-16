#!/usr/bin/env python3
"""
Assemble the rendered slide PNGs + speaker notes into a 16:9 PPTX (13 slides).

Each slideN.png is placed FULL-BLEED on a 16:9 slide; speaker notes are
attached per slide.

Run with: /home/juanfra/miniforge3/bin/python build_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))

# 16:9 at 13.33 x 7.5 inches (matches the PNG figsize)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Speaker notes per slide (1-indexed -> list index 0..12)
NOTES = [
    # 1 - title
    "",
    # 2 - fig2 mechanism
    "Flow assembly builds one splice graph from all reads and picks the "
    "FEWEST paths that explain coverage (parsimony) — so real but "
    "low-abundance / overlapping isoforms get merged or dropped. "
    "Read-coherence instead emits each molecule's own spliced path.",
    # 3 - fig1 igv locus
    "Real chr19 locus (NC_073243.2). The four StringTie variants splice a "
    "long intron straight over the internal exons; the read-coherence isoform "
    "retains them, and 15 spanning long-read alignments are the evidence. "
    "Eyeball-verifiable in IGV.",
    # 4 - identity = GROUP BY (text slide)
    "Clustering needs a distance + threshold and can merge distinct "
    "molecules; this is a discrete equivalence on junction chains (a GROUP "
    "BY), parameter-free and deterministic. Fragment folding is sub-path "
    "containment, not similarity. Junctions are matched exactly (fuzzy-merge "
    "off by default); the only tolerance is a few bp at the ragged 5'/3' ends.",
    # 5 - fig6 altsplice
    "Alternative splicing is NOT special-cased. A cassette-exon skip, an "
    "alternative 5' donor, an alternative 3' acceptor, or intron retention "
    "each change the junction list, so each is automatically a different "
    "exact key and a distinct transcript. The splice graph (the same one "
    "flow uses) represents the alternatives as branch points; read-coherence "
    "reads each molecule's realized path through it.",
    # 6 - fig4 gate collapse
    "Two precision mechanisms, both annotation-free. The realness gate keeps "
    "a read-chain isoform only if every junction is canonical, it has no "
    "RT-switch repeat signature, and read-depth >= min_cov. Degradation-aware "
    "collapse folds 5'/3'/internal fragments into their full-length parent "
    "(ISM -> FSM).",
    # 7 - fig3 architecture
    "Read-chain is held out of every per-bundle and global filter that could "
    "displace a flow transcript, then unioned back only as NOVEL intron chains "
    "and passed through the gate. So the output is a provable superset of the "
    "flow baseline: genome-wide gd\\rcg = 0 and st\\rcg = 0 (never loses a "
    "flow find or a StringTie guide).",
    # 8 - pseudocode (alignments -> transcripts)
    "This is figure 3 in code. Steps 1-3 are extract_transcripts_readchain: "
    "GROUP BY the exact intron chain (a hash key, an equivalence relation - "
    "not a distance/threshold), envelope the ragged 5'/3' ends, fold "
    "truncated fragments into their containing full-length chain. Step 4 is "
    "the additive merge: the flow assembly is the floor; read-chains are held "
    "out of the flow filters, then only NOVEL chains that pass the "
    "annotation-free realness gate (every junction canonical, no RT-switch, "
    "depth >= K) are added - so output is a provable superset of the flow "
    "output.",
    # 9 - fig5 results
    "Genome-wide (25 chromosomes). At the shipped default min_cov=3 (the "
    "precision/recall knee), read-coherence recovers +1,735 exact (FSM) "
    "annotated isoforms StringTie misses — about 3x the raw per-molecule "
    "approach. Of the gated read-chain extras: 100% canonical, 63% real, "
    "12% noise.",
    # 10 - positioning vs prior art (text slide)
    "Be upfront: grouping reads by intron chain is the standard long-read "
    "collapse (Cupcake, IsoSeq, FLAIR, TALON, IsoQuant, Mandalorion), and "
    "IsoQuant already showed per-read methods beat StringTie's graph/flow "
    "assembly on novel isoforms with lower FP. So the algorithm and the "
    "beat-StringTie result are prior art. What is ours: doing it as a "
    "provably-additive layer over a StringTie flow port (output superset of "
    "flow) in one tool, plus the annotation-free realness gate. The genuine "
    "novelty is the next slide — paralog-copy resolution via PSVs on a shared "
    "family graph.",
    # 11 - fig8 psv calc (calling a PSV)
    "How a PSV is determined. Homologous columns come from aligning the copies "
    "with POA (partial-order alignment, the graph generalization of "
    "Needleman-Wunsch); a column where the copies carry fixed different bases "
    "is a candidate PSV. It is accepted as real (not sequencing error or "
    "within-copy SNP) when the allele is consistent across essentially all "
    "reads of a copy and co-segregates with other PSVs; a family is declared "
    "identifiable only with >= K such columns, because K independent "
    "error-agreements are improbable at long-read error rates. Reads are then "
    "assigned to a copy by which PSV alleles they carry, plus the NM/AS gap.",
    # 12 - fig7 psv bridge (novel direction)
    "This is the novel direction and answers 'how would you use PSVs here'. "
    "For near-identical paralog copies that share the SAME junction chain, "
    "read-coherence alone can't separate them — same key, one bucket. PSVs "
    "(paralog-sequence variants: bases that differ consistently between "
    "copies) add a second axis: among reads in that bucket, partition by their "
    "PSV-allele vector to recover the individual copies. The transcript key "
    "extends from (intron chain) to (intron chain + PSV-allele vector): "
    "junctions give the structural identity, PSVs give the copy identity, both "
    "read off the SAME molecule. Closest prior art stops short: FLAIR2 phases "
    "per-read isoforms by HAPLOTYPE (alleles of one gene), Paraphase does "
    "PSV-based paralog resolution at the DNA/gene level — neither does "
    "paralog-copy-resolved isoform assembly. Caveat: bounded by "
    "identifiability — copies separate only when reads span enough PSVs above "
    "the error floor.",
    # 13 - fig9 contribution (the contribution in one slide)
    "The contribution in one slide: every existing method occupies a cell "
    "except de-novo, per-molecule isoform assembly resolved to paralog COPIES. "
    "read-coherence (per-molecule structure) plus PSVs (copy identity) on one "
    "spliced family variation graph fills it, with a provable identifiability "
    "boundary: copies are recoverable iff they differ at >= K PSV columns "
    "spanned by enough reads above the error floor; otherwise they are "
    "provably merged.",
]


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # fully blank

    for i in range(1, 14):
        png = os.path.join(HERE, f"slide{i}.png")
        if not os.path.exists(png):
            raise FileNotFoundError(png)

        slide = prs.slides.add_slide(blank_layout)

        # full-bleed image
        slide.shapes.add_picture(png, 0, 0, width=SLIDE_W, height=SLIDE_H)

        # speaker notes
        note = NOTES[i - 1]
        if note:
            slide.notes_slide.notes_text_frame.text = note

    out = os.path.join(HERE, "read_coherence_deck.pptx")
    try:
        prs.save(out)
        print("wrote", out)
    except PermissionError:
        # Canonical output is locked (e.g. open in PowerPoint on the Windows
        # side, leaving a ~$ owner-lock file). Don't lose the build: write to a
        # fallback path so the 13-slide deck is still produced, and report it.
        alt = os.path.join(HERE, "read_coherence_deck.new.pptx")
        prs.save(alt)
        print("WARNING:", out, "is locked (open elsewhere); wrote", alt,
              "instead. Close the file and rename it over the original.")


if __name__ == "__main__":
    main()
