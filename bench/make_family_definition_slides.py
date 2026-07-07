#!/usr/bin/env python3
"""
make_family_definition_slides.py — illustrate what a multi-copy gene family is,
using a real example from the refined catalog (Family 69, RABL2):
  1. definition slide
  2. birdseye quasi-clique view of one family
  3. zoom on one copy showing isoforms

Run: /home/juanfra/miniforge3/bin/python bench/make_family_definition_slides.py
Output: bench/family_definition_slides.pptx (+ 3 PNGs)
"""
import os
import csv
import math
import pysam
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
REFINE_TSV = os.path.join(HERE, "family_rna_refine.tsv")
META_TSV = os.path.join("/home/juanfra/winloci_scratch", "denovo_transcripts.meta.tsv")
SKEL_TSV = os.path.join("/home/juanfra/winloci_scratch", "denovo_skeletons.tsv")
EDGES_TSV = os.path.join(HERE, "denovo_family_edges.tsv")

FIG_DEF = os.path.join(HERE, "famdef_1_definition.png")
FIG_BIRD = os.path.join(HERE, "famdef_2_birdseye.png")
FIG_ISO = os.path.join(HERE, "famdef_3_isoforms.png")
FIG_RECOMB = os.path.join(HERE, "famdef_4_recombination.png")
FIG_IGV = os.path.join(HERE, "famdef_5_igv.png")
FIG_VALID = os.path.join(HERE, "famdef_6_validation.png")
FIG_PHASE = os.path.join(HERE, "famdef_7_phasing.png")
FIG_SILENT = os.path.join(HERE, "famdef_8_silent_copy.png")
OUT = os.path.join(HERE, "family_definition_slides.pptx")

NAVY = RGBColor(0x1F, 0x2D, 0x5A)
TEAL = RGBColor(0x12, 0x7A, 0x6E)
GREY = RGBColor(0x44, 0x44, 0x44)
ORANGE = RGBColor(0xD9, 0x6B, 0x27)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)

CN, CT, CG, CO, CP = "#1F2D5A", "#127A6E", "#444444", "#D96B27", "#8E44AD"


def load_meta():
    d = {}
    with open(META_TSV) as fh:
        for r in csv.DictReader(fh, delimiter="\t"):
            d[r["id"]] = dict(
                chrom=r["chrom"],
                start=int(r["start"]),
                end=int(r["end"]),
                strand=r["strand"],
                n_exon=int(r["n_exon"]),
                n_reads=int(r["n_reads"]),
            )
    return d


def load_skeletons():
    # key by (chrom, start, end, n_exon, n_reads) so isoforms with the same
    # coordinates but different splicing resolve to the correct skeleton.
    d = {}
    with open(SKEL_TSV) as fh:
        for r in csv.DictReader(fh, delimiter="\t"):
            c, s, e = r["chrom"], int(r["start"]), int(r["end"])
            nx, nr = int(r["n_exon"]), int(r["n_reads"])
            introns = []
            if r["introns"].strip():
                for tok in r["introns"].split(";"):
                    a, b = tok.split("-")
                    introns.append((int(a), int(b)))
            d[(c, s, e, nx, nr)] = introns
    return d


def exon_intervals(meta, skel, tid):
    m = meta[tid]
    c, s, e = m["chrom"], m["start"], m["end"]
    key = (c, s, e, m["n_exon"], m["n_reads"])
    introns = skel.get(key, [])
    exons = []
    cur = s
    for d, a in introns:
        exons.append((cur, d))
        cur = a
    exons.append((cur, e + 1))
    return exons


def load_family(fid="69"):
    members = []
    with open(REFINE_TSV) as fh:
        for r in csv.DictReader(fh, delimiter="\t"):
            if r["family_id"] == fid:
                members.append(r)
    return members


def load_edges(members):
    ids = {m["member_dn"] for m in members}
    edges = []
    with open(EDGES_TSV) as fh:
        for r in csv.DictReader(fh, delimiter="\t"):
            if r["a"] in ids and r["b"] in ids:
                edges.append((r["a"], r["b"], float(r["core_recip"])))
    return edges


# ----------------------------------------------------------------- figure 1: definition
def make_definition():
    fig, ax = plt.subplots(figsize=(12.0, 6.2))
    ax.set_xlim(0, 12); ax.set_ylim(0, 6.2); ax.axis("off")

    # chromosome strip
    ax.add_patch(Rectangle((0.8, 4.6), 10.4, 0.25, fc="#dfe7f2", ec=CN, lw=1.5, zorder=1))
    ax.text(6.0, 5.15, "chromosome", ha="center", fontsize=10, color=CN, style="italic")

    # copy A locus
    ax.add_patch(FancyBboxPatch((1.5, 4.25), 2.2, 0.95, boxstyle="round,pad=0.05,rounding_size=0.1",
                 fc=CT, ec=CN, lw=2.0, zorder=2))
    ax.text(2.6, 4.95, "copy A", ha="center", va="center", fontsize=13, color="white", fontweight="bold")
    ax.text(2.6, 4.55, "locus 1", ha="center", va="center", fontsize=9, color="white")

    # copy B locus
    ax.add_patch(FancyBboxPatch((7.0, 4.25), 2.2, 0.95, boxstyle="round,pad=0.05,rounding_size=0.1",
                 fc=CO, ec=CN, lw=2.0, zorder=2))
    ax.text(8.1, 4.95, "copy B", ha="center", va="center", fontsize=13, color="white", fontweight="bold")
    ax.text(8.1, 4.55, "locus 2", ha="center", va="center", fontsize=9, color="white")

    # homology edge
    ax.annotate("", xy=(7.0, 4.72), xytext=(3.7, 4.72),
                arrowprops=dict(arrowstyle="<->", color=CN, lw=2.5))
    ax.text(5.35, 5.0, "homology edge", ha="center", fontsize=11, color=CN, fontweight="bold")

    # definition box
    ax.add_patch(FancyBboxPatch((1.2, 0.75), 9.6, 3.2, boxstyle="round,pad=0.1,rounding_size=0.15",
                 fc="#f6f8fb", ec=CN, lw=2.0, zorder=1))

    # title at the top of the box
    ax.text(6.0, 3.75, "Multi-copy gene family", ha="center", fontsize=20, color=CN, fontweight="bold")

    # definition equation below title with clear separation
    ax.text(6.0, 2.55,
            "= two or more distinct genomic loci\n"
            "  that carry homologous copies of the same gene\n"
            "  and are expressed as RNA",
            ha="center", fontsize=14, color=CG, linespacing=1.5)

    # bullet distinctions
    ax.text(2.0, 1.55, "• copies = different loci (often different chromosomes)", ha="left", fontsize=11, color=CG)
    ax.text(2.0, 1.15, "• isoforms = same locus, different splicing", ha="left", fontsize=11, color=CG)

    fig.text(0.5, 0.02,
             "A family is a connected, cohesive subgraph (γ-quasi-clique) of the homology graph; "
             "each node is a distinct locus, each edge is a sequence-homology link.",
             ha="center", fontsize=9, style="italic", color=CG)
    fig.savefig(FIG_DEF, dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- figure 2: birdseye quasi-clique
def make_birdseye():
    members = load_family("69")
    meta = load_meta()
    edges = load_edges(members)

    fig, (axL, axR) = plt.subplots(1, 2, figsize=(13.3, 6.2),
                                   gridspec_kw={"width_ratios": [1.25, 1]})

    # ---- left: genomic view (one strip per chromosome) ----
    axL.set_xlim(0, 10)
    axL.set_ylim(0, 6.2)
    axL.axis("off")

    # sort by chromosome name for stable layout
    members_sorted = sorted(members, key=lambda m: (m["chrom"], int(m["start"])))
    colors = [CT, CO, CP, "#2874A6", "#117864"]
    y0 = 5.1
    dy = 1.05

    for i, m in enumerate(members_sorted):
        y = y0 - i * dy
        chrom = m["chrom"]
        s, e = int(m["start"]), int(m["end"])
        gene = m["member_gene"] if m["member_gene"] != "NA" else f"copy {i+1}"
        tid = m["member_dn"]
        n_exon = meta[tid]["n_exon"]

        # chromosome bar
        axL.add_patch(Rectangle((2.2, y - 0.08), 7.0, 0.16, fc="#dfe7f2", ec=CN, lw=1.2))
        # locus rectangle (width proportional to log10 length, capped)
        length = e - s
        width = min(2.8, max(0.6, math.log10(length) * 0.55))
        color = colors[i % len(colors)]
        axL.add_patch(FancyBboxPatch((3.0, y - 0.18), width, 0.36,
                     boxstyle="round,pad=0.02,rounding_size=0.05",
                     fc=color, ec=CN, lw=1.5, zorder=2))
        # label inside the locus rectangle (coordinates omitted to avoid clutter)
        label = f"{gene}"
        axL.text(3.0 + width / 2, y, label, ha="center", va="center",
                 fontsize=8, color="white", fontweight="bold")
        # chromosome label on the left
        axL.text(2.05, y, chrom, ha="right", va="center", fontsize=9, color=CN, fontweight="bold")

    axL.set_title("Family 69 — RABL2  (real refined family)",
                  fontsize=14, fontweight="bold", color=CN, pad=12)
    axL.text(5.0, 0.35,
             "5 distinct loci on 5 different chromosomes  ·  RABL2A and RABL2B are highlighted paralogs",
             ha="center", fontsize=10, color=CG, style="italic")

    # ---- right: homology graph ----
    axR.set_xlim(0, 6); axR.set_ylim(0, 6); axR.axis("off")
    axR.set_title("the homology graph (perfect clique)", fontsize=14, fontweight="bold", color=CN)

    # pentagon layout, large enough to avoid label collisions
    ids = [m["member_dn"] for m in members_sorted]
    labels = {}
    for i, m in enumerate(members_sorted):
        gene = m["member_gene"] if m["member_gene"] != "NA" else f"copy {i+1}"
        chrom = m["chrom"].replace("NC_", "").replace(".1", "").replace(".2", "")
        labels[m["member_dn"]] = f"{gene}\nchr {chrom}"

    center = (3.0, 3.0)
    radius = 2.1
    pos = {}
    for idx, tid in enumerate(ids):
        angle = math.pi / 2 - idx * (2 * math.pi / len(ids))
        pos[tid] = (center[0] + radius * math.cos(angle),
                    center[1] + radius * math.sin(angle))

    node_colors = {tid: colors[i % len(colors)] for i, tid in enumerate(ids)}

    # edges (complete graph)
    edge_seen = set()
    for a, b, cr in edges:
        if (b, a) in edge_seen:
            continue
        edge_seen.add((a, b))
        axR.plot([pos[a][0], pos[b][0]], [pos[a][1], pos[b][1]],
                 color=CN, lw=1.8, alpha=0.75, zorder=1)

    # nodes
    for tid, (x, y) in pos.items():
        fc = node_colors[tid]
        axR.add_patch(FancyBboxPatch((x - 0.62, y - 0.38), 1.24, 0.76,
                     boxstyle="round,pad=0.04,rounding_size=0.08",
                     fc=fc, ec=CN, lw=2.0, zorder=3))
        axR.text(x, y, labels[tid], ha="center", va="center",
                 fontsize=7.5, color="white", fontweight="bold", zorder=4, linespacing=1.05)

    # legend / note
    axR.text(3.0, 0.55,
             "all 10 possible edges present  →  perfect clique\n"
             "every locus is homologous to every other locus",
             ha="center", fontsize=9, color=CG, style="italic", linespacing=1.3)

    fig.savefig(FIG_BIRD, dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- figure 3: two loci, each with isoforms as vg paths
def make_isoforms():
    meta = load_meta()
    skel = load_skeletons()

    # Locus data: two family-69 copies that each have multiple isoforms
    locus_data = [
        {
            "name": "Locus 1  LOC101142457  (NC_073247.2)",
            "y": 1.9,
            "color_a": CT,
            "color_b": CO,
            "isoforms": [
                ("DN_NC_073247.2_167707352_7", "isoform a  (7 exons, 10 reads)"),
                ("DN_NC_073247.2_167707352_8", "isoform b  (8 exons, 6 reads)"),
            ],
        },
        {
            "name": "Locus 2  RABL2B  (NC_086018.1)",
            "y": -1.9,
            "color_a": CT,
            "color_b": CO,
            "isoforms": [
                ("DN_NC_086018.1_48818439_9", "isoform a  (9 exons, 4 reads)"),
                ("DN_NC_086018.1_48818439_10", "isoform b  (10 exons, 3 reads)"),
            ],
        },
    ]

    fig, ax = plt.subplots(figsize=(13.0, 7.2))
    ax.set_xlim(0, 13)
    ax.set_ylim(-3.6, 3.6)
    ax.axis("off")

    # draw each locus graph
    graph_anchors = []
    for locus in locus_data:
        anchor = _draw_locus_vg(ax, locus, meta, skel)
        graph_anchors.append(anchor)

    # homology edge between the two loci
    ax.annotate("", xy=(12.2, graph_anchors[1]["y"]),
                xytext=(12.2, graph_anchors[0]["y"]),
                arrowprops=dict(arrowstyle="<->", color=CN, lw=2.5))
    ax.text(12.55, (graph_anchors[0]["y"] + graph_anchors[1]["y"]) / 2,
            "homology\nedge", ha="left", va="center", fontsize=10,
            color=CN, fontweight="bold")

    ax.set_title("Two family copies, each producing multiple isoforms from its own variation graph",
                 fontsize=14, fontweight="bold", color=CN, pad=12)
    ax.text(6.2, 3.55,
            "Shared exons are central split-colored nodes; isoform-specific exons branch into bubbles",
            ha="center", fontsize=10, color=CG, style="italic")

    # legend
    ax.add_patch(FancyBboxPatch((10.4, 2.55), 0.35, 0.22,
                 boxstyle="round,pad=0.02,rounding_size=0.04", fc=CT, ec=CN))
    ax.text(10.9, 2.66, "isoform a", ha="left", va="center", fontsize=9, color=CG)
    ax.add_patch(FancyBboxPatch((10.4, 2.15), 0.35, 0.22,
                 boxstyle="round,pad=0.02,rounding_size=0.04", fc=CO, ec=CN))
    ax.text(10.9, 2.26, "isoform b", ha="left", va="center", fontsize=9, color=CG)
    ax.add_patch(FancyBboxPatch((10.4, 1.75), 0.17, 0.22,
                 boxstyle="round,pad=0.02,rounding_size=0.04", fc=CT, ec=CN))
    ax.add_patch(FancyBboxPatch((10.57, 1.75), 0.18, 0.22,
                 boxstyle="round,pad=0.02,rounding_size=0.04", fc=CO, ec=CN))
    ax.text(10.9, 1.86, "shared exon", ha="left", va="center", fontsize=9, color=CG)

    fig.savefig(FIG_ISO, dpi=150, bbox_inches="tight")
    plt.close(fig)


def _draw_locus_vg(ax, locus, meta, skel):
    """Draw one variation-graph panel for a locus. Returns a dict with y center and rightmost x."""
    tid_a, label_a = locus["isoforms"][0]
    tid_b, label_b = locus["isoforms"][1]
    col_a, col_b = locus["color_a"], locus["color_b"]
    y0 = locus["y"]

    exons_a = exon_intervals(meta, skel, tid_a)
    exons_b = exon_intervals(meta, skel, tid_b)

    # build node list with owners
    node_exons = []
    node_owners = []
    a_nodes = []
    b_nodes = []
    seen = {}
    for s, e in exons_a:
        key = (s, e)
        if key not in seen:
            seen[key] = len(node_exons)
            node_exons.append(key)
            node_owners.append("a")
        a_nodes.append(seen[key])
    for s, e in exons_b:
        key = (s, e)
        if key not in seen:
            seen[key] = len(node_exons)
            node_exons.append(key)
            node_owners.append("b")
        idx = seen[key]
        b_nodes.append(idx)
        if node_owners[idx] == "a":
            node_owners[idx] = "shared"

    # Manual layouts that keep genomic order and avoid overlap.
    # Nodes are numbered in the order they were first seen (isoform a first, then b).
    if "073247" in tid_a:
        # LOC101142457: 9 distinct exons
        # shared: 0-1-2-4-5-6, a-specific: 7 (final), b-specific: 3 (middle bubble), 8 (final)
        pos = {
            0: (1.0, y0),
            1: (2.1, y0),
            2: (3.2, y0),
            3: (4.3, y0 + 0.7),   # b-specific middle exon
            4: (5.4, y0),
            5: (6.5, y0),
            6: (7.6, y0),
            7: (9.0, y0 + 0.7),   # a-specific final exon
            8: (9.0, y0 - 0.7),   # b-specific final exon
        }
        a_path = [0, 1, 2, 4, 5, 6, 7]
        b_path = [0, 1, 2, 3, 4, 5, 6, 8]
    else:
        # RABL2B: 12 distinct exons
        # shared: 0-1-2-3-4-5-6, a-specific: 7-8, b-specific: 9-10-11
        pos = {
            0: (1.0, y0),
            1: (2.1, y0),
            2: (3.2, y0),
            3: (4.3, y0),
            4: (5.4, y0),
            5: (6.5, y0),
            6: (7.6, y0),
            7: (9.0, y0 + 0.7),
            8: (10.4, y0 + 0.7),
            9: (9.0, y0 - 0.7),
            10: (10.4, y0 - 0.7),
            11: (11.8, y0 - 0.7),
        }
        a_path = [0, 1, 2, 3, 4, 5, 6, 7, 8]
        b_path = [0, 1, 2, 3, 4, 5, 6, 9, 10, 11]

    def draw_node(idx, x, y):
        owner = node_owners[idx]
        width = 0.82
        height = 0.52
        if owner == "shared":
            ax.add_patch(FancyBboxPatch((x - width / 2, y - height / 2), width / 2, height,
                         boxstyle="round,pad=0.02,rounding_size=0.06",
                         fc=col_a, ec=CN, lw=1.5, zorder=3))
            ax.add_patch(FancyBboxPatch((x, y - height / 2), width / 2, height,
                         boxstyle="round,pad=0.02,rounding_size=0.06",
                         fc=col_b, ec=CN, lw=1.5, zorder=3))
        elif owner == "a":
            ax.add_patch(FancyBboxPatch((x - width / 2, y - height / 2), width, height,
                         boxstyle="round,pad=0.02,rounding_size=0.06",
                         fc=col_a, ec=CN, lw=1.5, zorder=3))
        else:
            ax.add_patch(FancyBboxPatch((x - width / 2, y - height / 2), width, height,
                         boxstyle="round,pad=0.02,rounding_size=0.06",
                         fc=col_b, ec=CN, lw=1.5, zorder=3))
        ax.text(x, y, f"E{idx + 1}", ha="center", va="center",
                fontsize=8, color="white", fontweight="bold", zorder=4)

    # splice-junction edges
    edge_pairs = set()
    for u, v in zip(a_path, a_path[1:]):
        edge_pairs.add((u, v))
    for u, v in zip(b_path, b_path[1:]):
        edge_pairs.add((u, v))
    for u, v in edge_pairs:
        x1, y1 = pos[u]
        x2, y2 = pos[v]
        style = "arc3,rad=0.12" if abs(y2 - y1) > 0.3 else "arc3,rad=0.0"
        ax.annotate("", xy=(x2 - 0.40, y2), xytext=(x1 + 0.40, y1),
                    arrowprops=dict(arrowstyle="-", color=CN, lw=1.4,
                                    connectionstyle=style), zorder=1)

    # path overlays (offset on shared backbone so both are visible)
    for path, col, label, y_legend, offset in [(a_path, col_a, label_a, y0 + 1.15, 0.14),
                                                (b_path, col_b, label_b, y0 - 1.15, -0.14)]:
        xs, ys = [], []
        for idx in path:
            x, y = pos[idx]
            y_eff = y + (offset if node_owners[idx] == "shared" else 0)
            xs.append(x)
            ys.append(y_eff)
        ax.plot(xs, ys, color=col, lw=6.5, alpha=0.35, zorder=2,
                solid_capstyle="round", solid_joinstyle="round")
        ax.plot(xs, ys, color=col, lw=2.8, alpha=0.85, zorder=2,
                solid_capstyle="round", solid_joinstyle="round")
        ax.text(0.35, y_legend, label, ha="left", va="center",
                fontsize=10, color=col, fontweight="bold")

    # draw nodes on top
    for idx, (x, y) in pos.items():
        draw_node(idx, x, y)

    # locus label above/below the path labels to avoid overlapping nodes
    label_y = y0 + 1.35 if y0 > 0 else y0 - 1.35
    ax.text(0.35, label_y, locus["name"], ha="left", va="center",
            fontsize=10, color=CN, fontweight="bold")

    # chromosome bar under this locus
    bar_y = y0 - 1.45
    ax.add_patch(Rectangle((0.6, bar_y), 11.8, 0.10, fc="#dfe7f2", ec=CN, lw=1.0))
    if "073247" in tid_a:
        ticks = [(1.0, 167.707), (5.4, 167.716), (9.0, 167.723)]
    else:
        ticks = [(1.0, 48.818), (5.4, 48.823), (11.8, 48.825)]
    for tx, coord in ticks:
        ax.plot([tx, tx], [bar_y, bar_y - 0.07], color=CG, lw=1.0)
        ax.text(tx, bar_y - 0.18, f"{coord:.3f}M", ha="center", fontsize=7, color=CG)

    rightmost = max(x for x, y in pos.values())
    return {"y": y0, "x_right": rightmost}


# ----------------------------------------------------------------- figure 4: recombinant reads don't merge copies
def make_recombination():
    fig, ax = plt.subplots(figsize=(12.5, 6.2))
    ax.set_xlim(0, 12.5)
    ax.set_ylim(0, 6.2)
    ax.axis("off")

    # --- top: two distinct loci on separate chromosomes ---
    y_chrom = 5.2
    # chromosome A strip
    ax.add_patch(Rectangle((1.0, y_chrom - 0.08), 4.0, 0.16, fc="#dfe7f2", ec=CN, lw=1.5))
    ax.text(3.0, y_chrom + 0.35, "chromosome A  (copy 1 locus)", ha="center", fontsize=10, color=CN)
    # copy 1 exons
    ax.add_patch(FancyBboxPatch((1.4, y_chrom - 0.28), 0.9, 0.56,
                 boxstyle="round,pad=0.02,rounding_size=0.06", fc=CT, ec=CN, lw=1.5))
    ax.text(1.85, y_chrom, "E1", ha="center", va="center", fontsize=10, color="white", fontweight="bold")
    ax.add_patch(FancyBboxPatch((3.0, y_chrom - 0.28), 0.9, 0.56,
                 boxstyle="round,pad=0.02,rounding_size=0.06", fc=CT, ec=CN, lw=1.5))
    ax.text(3.45, y_chrom, "E2", ha="center", va="center", fontsize=10, color="white", fontweight="bold")
    ax.annotate("", xy=(3.0, y_chrom), xytext=(2.3, y_chrom),
                arrowprops=dict(arrowstyle="-", color=CN, lw=2.0))

    # chromosome B strip
    ax.add_patch(Rectangle((7.5, y_chrom - 0.08), 4.0, 0.16, fc="#dfe7f2", ec=CN, lw=1.5))
    ax.text(9.5, y_chrom + 0.35, "chromosome B  (copy 2 locus)", ha="center", fontsize=10, color=CN)
    # copy 2 exons
    ax.add_patch(FancyBboxPatch((7.9, y_chrom - 0.28), 0.9, 0.56,
                 boxstyle="round,pad=0.02,rounding_size=0.06", fc=CO, ec=CN, lw=1.5))
    ax.text(8.35, y_chrom, "E1", ha="center", va="center", fontsize=10, color="white", fontweight="bold")
    ax.add_patch(FancyBboxPatch((9.5, y_chrom - 0.28), 0.9, 0.56,
                 boxstyle="round,pad=0.02,rounding_size=0.06", fc=CO, ec=CN, lw=1.5))
    ax.text(9.95, y_chrom, "E2", ha="center", va="center", fontsize=10, color="white", fontweight="bold")
    ax.annotate("", xy=(9.5, y_chrom), xytext=(8.8, y_chrom),
                arrowprops=dict(arrowstyle="-", color=CN, lw=2.0))

    # --- middle: combined vg view with recombinant edge ---
    # left copy nodes
    c1e1 = (2.5, 3.1)
    c1e2 = (2.5, 1.5)
    # right copy nodes
    c2e1 = (9.0, 3.1)
    c2e2 = (9.0, 1.5)

    # valid within-copy edges
    ax.annotate("", xy=(c1e2[0], c1e2[1] + 0.35), xytext=(c1e1[0], c1e1[1] - 0.35),
                arrowprops=dict(arrowstyle="-", color=CN, lw=2.5))
    ax.annotate("", xy=(c2e2[0], c2e2[1] + 0.35), xytext=(c2e1[0], c2e1[1] - 0.35),
                arrowprops=dict(arrowstyle="-", color=CN, lw=2.5))

    # invalid recombinant edge (copy1-E2 -> copy2-E1)
    ax.annotate("", xy=(c2e1[0] - 0.42, c2e1[1]), xytext=(c1e2[0] + 0.42, c1e2[1]),
                arrowprops=dict(arrowstyle="->", color="#C0392B", lw=2.5, ls="--",
                                connectionstyle="arc3,rad=0.15"))

    # nodes
    for (x, y), label, col in [(c1e1, "copy 1\nE1", CT), (c1e2, "copy 1\nE2", CT),
                                (c2e1, "copy 2\nE1", CO), (c2e2, "copy 2\nE2", CO)]:
        ax.add_patch(FancyBboxPatch((x - 0.55, y - 0.40), 1.1, 0.8,
                     boxstyle="round,pad=0.04,rounding_size=0.10", fc=col, ec=CN, lw=2.0, zorder=3))
        ax.text(x, y, label, ha="center", va="center", fontsize=9, color="white",
                fontweight="bold", linespacing=1.05, zorder=4)

    # forbidden symbol on the recombinant edge
    ax.text(5.75, 2.55, "✗", ha="center", va="center", fontsize=28, color="#C0392B",
            fontweight="bold", zorder=5)
    ax.text(5.75, 2.05, "recombinant\njunction", ha="center", va="center", fontsize=9,
            color="#C0392B", fontweight="bold", linespacing=1.1)

    # valid-path labels
    ax.text(2.5, 1.35, "valid path\n(copy 1)", ha="center", va="center", fontsize=9,
            color=CT, fontweight="bold", linespacing=1.1)
    ax.text(9.0, 1.35, "valid path\n(copy 2)", ha="center", va="center", fontsize=9,
            color=CO, fontweight="bold", linespacing=1.1)

    # --- bottom: rule box ---
    ax.add_patch(FancyBboxPatch((1.2, -0.35), 10.1, 1.05,
                 boxstyle="round,pad=0.08,rounding_size=0.10", fc="#fdf2f2", ec="#C0392B",
                 lw=2.0, zorder=1))
    ax.text(6.25, 0.45, "A locus is defined by collinear, supported intron chains", ha="center",
            fontsize=13, color=CN, fontweight="bold")
    ax.text(6.25, 0.00,
            "A read that splices copy-1-E2 into copy-2-E1 creates a non-collinear, unsupported junction. "
            "It is not added to either locus's variation graph, so the two copies remain separate.",
            ha="center", fontsize=10, color=CG, linespacing=1.3)

    ax.set_title("Cross-locus recombinant reads do not merge distinct copies",
                 fontsize=14, fontweight="bold", color=CN, pad=12)

    fig.savefig(FIG_RECOMB, dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- figure 5: IGV-style read view
def _read_blocks_and_introns(read):
    """Return aligned blocks [(s,e),...] and introns [(s,e),...] from a pysam read."""
    blocks = []
    introns = []
    pos = read.reference_start
    block_start = pos
    in_block = True
    for op, length in read.cigartuples:
        if op == 3:  # N: intron
            if in_block:
                blocks.append((block_start, pos))
                in_block = False
            introns.append((pos, pos + length))
            pos += length
            block_start = pos
        elif op in (0, 7, 8):  # M, =, X: match/mismatch
            if not in_block:
                block_start = pos
                in_block = True
            pos += length
        elif op in (1, 4, 5):  # I, S, H: no ref consumption
            pass
        elif op == 2:  # D: deletion consumes reference but no query
            pos += length
    if in_block:
        blocks.append((block_start, pos))
    return blocks, introns


def _intron_chain(read):
    _, introns = _read_blocks_and_introns(read)
    return tuple(introns)


def _draw_igv_locus(ax, chrom, region_start, region_end, reads, gene_exons, title):
    """Draw one IGV-style panel (axes uses genomic x coordinates and stacked y positions)."""
    ax.set_xlim(region_start, region_end)
    ax.axis("off")

    row_height = 0.15
    gap = 0.05
    gene_track_height = 0.35
    track_y = 0.0

    # --- gene track ---
    ax.add_patch(Rectangle((region_start, track_y - 0.02), region_end - region_start, 0.04,
                           fc="#aab7c8", ec="none", zorder=1))
    for s, e in gene_exons:
        ax.add_patch(FancyBboxPatch((s, track_y - gene_track_height / 2), e - s, gene_track_height,
                     boxstyle="round,pad=0.01,rounding_size=0.02",
                     fc=CT, ec=CN, lw=1.2, zorder=2))
    ax.text(region_start, track_y + gene_track_height / 2 + 0.18, title,
            ha="left", va="bottom", fontsize=10, color=CN, fontweight="bold")

    # --- reads ---
    from collections import defaultdict
    chain_groups = defaultdict(list)
    for r in reads:
        if r.is_unmapped:
            continue
        chain_groups[_intron_chain(r)].append(r)

    selected = []
    for chain, group in sorted(chain_groups.items(), key=lambda x: -len(x[1]))[:4]:
        selected.extend(group[:3])

    selected.sort(key=lambda r: (r.reference_start, -r.reference_end))
    row_ends = []
    placements = []
    for r in selected:
        y = None
        for i, end in enumerate(row_ends):
            if r.reference_start > end + 5000:
                y = track_y - gene_track_height / 2 - 0.25 - i * (row_height + gap)
                row_ends[i] = r.reference_end
                break
        if y is None:
            y = track_y - gene_track_height / 2 - 0.25 - len(row_ends) * (row_height + gap)
            row_ends.append(r.reference_end)
        placements.append((r, y))

    for r, y in placements:
        blocks, introns = _read_blocks_and_introns(r)
        for s, e in introns:
            ax.plot([s, e], [y, y], color="#7f8c8d", lw=1.0, zorder=1)
        for s, e in blocks:
            ax.add_patch(Rectangle((s, y - row_height / 2), e - s, row_height,
                         fc="#5dade2" if not r.is_reverse else "#aeb6bf",
                         ec="none", zorder=2))

    bottom_y = track_y - gene_track_height / 2 - 0.25 - (len(row_ends) - 1) * (row_height + gap) - row_height

    # coordinate axis
    tick_positions = [region_start, (region_start + region_end) // 2, region_end]
    axis_y = bottom_y - 0.12
    ax.plot([region_start, region_end], [axis_y, axis_y], color=CG, lw=1.0)
    for x in tick_positions:
        ax.plot([x, x], [axis_y, axis_y - 0.05], color=CG, lw=1.0)
        ax.text(x, axis_y - 0.10, f"{x / 1e6:.3f}M", ha="center", va="top", fontsize=8, color=CG)

    ax.set_ylim(axis_y - 0.25, track_y + gene_track_height / 2 + 0.35)


def make_igv():
    bam_path = "/home/juanfra/winloci_scratch/GGO.bam"
    bam = pysam.AlignmentFile(bam_path, "rb")

    loci = [
        {
            "title": "LOC101142457  (NC_073247.2)",
            "chrom": "NC_073247.2",
            "start": 167707352,
            "end": 167723174,
            "exons": [(167707352, 167708718), (167708884, 167708968), (167709298, 167709396),
                      (167716208, 167716288), (167717489, 167717519), (167721697, 167721857),
                      (167723010, 167723174)],
        },
        {
            "title": "RABL2B  (NC_086018.1)",
            "chrom": "NC_086018.1",
            "start": 48818439,
            "end": 48832011,
            "exons": [(48818439, 48818874), (48819040, 48819124), (48819454, 48819552),
                      (48819898, 48820010), (48825425, 48825505), (48826332, 48826412),
                      (48827614, 48827644), (48830502, 48830666), (48831815, 48832011)],
        },
    ]

    fig, axes = plt.subplots(2, 1, figsize=(13.0, 6.8))

    for ax, locus in zip(axes, loci):
        reads = list(bam.fetch(locus["chrom"], locus["start"], locus["end"]))
        _draw_igv_locus(ax, locus["chrom"], locus["start"], locus["end"],
                        reads, locus["exons"], locus["title"])

    fig.suptitle("IGV-style view: reads supporting the two loci",
                 fontsize=14, fontweight="bold", color=CN, y=0.98)
    fig.text(0.5, 0.01,
             "Reads are colored by strand (blue = forward, grey = reverse); splice junctions are thin grey lines. "
             "Each panel shows representative reads from the dominant isoform classes.",
             ha="center", fontsize=9, color=CG, style="italic")

    fig.savefig(FIG_IGV, dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- figure 6: validation on known families
def make_validation():
    rows = []
    with open(os.path.join(HERE, "known_family_showcase.tsv")) as fh:
        for r in csv.DictReader(fh, delimiter="\t"):
            rows.append(dict(
                family=r["family"],
                recall=float(r["recall"]),
                precision_strict=float(r["precision_strict_EpDNA"]),
                precision_bio=float(r["precision_bio_corrected"]),
                n=int(r["n_copies_catalog"]),
            ))

    # sort alphabetically so the plot is not ordered to look cherry-picked
    rows.sort(key=lambda x: x["family"])

    fig, ax = plt.subplots(figsize=(9.2, 6.0))
    ax.set_xlim(0.45, 1.05)
    ax.set_ylim(0.60, 1.05)

    # quadrant thresholds
    thr = 0.85
    ax.axvline(thr, color=CG, lw=1.0, ls="--", alpha=0.5)
    ax.axhline(thr, color=CG, lw=1.0, ls="--", alpha=0.5)

    for r in rows:
        fam, rec, prec, n = r["family"], r["recall"], r["precision_strict"], r["n"]
        clean = rec >= thr and prec >= thr
        color = CT if clean else CO
        size = 120 + n * 35
        ax.scatter(rec, prec, s=size, c=color, ec=CN, lw=1.5, zorder=3, alpha=0.9)
        # label all points; nudge a few to avoid overlap
        dx, dy = 0.015, 0.012
        if fam == "GSTM":
            dx, dy = -0.10, -0.025
        elif fam == "ZNF92":
            dx, dy = 0.015, -0.025
        elif fam == "APOBEC3":
            dx, dy = -0.055, 0.015
        ax.annotate(fam, (rec + dx, prec + dy), fontsize=9, color=CN,
                    fontweight="bold" if clean else "normal")

    # reference diagonal / perfect corner
    ax.plot([0.45, 1.0], [0.45, 1.0], color="#cccccc", lw=1.0, ls="-", zorder=1)

    ax.set_xlabel("recall  (expressed copies grouped / all expressed copies)", fontsize=11, color=CN)
    ax.set_ylabel("strict precision  (whole-protein reciprocal homology)", fontsize=11, color=CN)
    ax.set_title("Strict precision vs recall for all 8 canonical known families",
                 fontsize=14, fontweight="bold", color=CN, pad=12)

    # legend
    ax.scatter([], [], s=150, c=CT, ec=CN, label=f"both ≥ {thr}")
    ax.scatter([], [], s=150, c=CO, ec=CN, label="one metric < threshold")
    ax.legend(loc="lower left", frameon=True, fontsize=9)

    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    fig.text(0.5, 0.01,
             "Each point is one literature-recognized family; point size = number of catalog copies. "
             "Families in the top-right quadrant are clean on both strict precision and recall.",
             ha="center", fontsize=9, color=CG, style="italic")

    fig.savefig(FIG_VALID, dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- figure 7: long-read phasing / copy assignment
def make_phasing():
    fig, ax = plt.subplots(figsize=(13.0, 6.5))
    ax.set_xlim(0, 13.0)
    ax.set_ylim(-1.1, 6.5)
    ax.axis("off")

    # Three copies, four PSV columns.  A SUN is a column where one copy's allele is private.
    copies = [
        ("copy 1\n(RABL2A-like)", CT, ["A", "G", "C", "T"]),
        ("copy 2\n(RABL2B-like)", CO, ["A", "T", "C", "G"]),
        ("copy 3\n(paralog LOC)", CP, ["G", "G", "T", "T"]),
    ]
    labels = ["PSV 1", "PSV 2", "PSV 3", "PSV 4"]
    x_pos = [3.0, 5.4, 7.8, 10.2]
    label_x = 1.15  # right-aligned copy label anchor

    # Determine SUN columns: where exactly one copy has a distinct allele.
    sun_by_col = []
    for col in range(4):
        alleles = [copies[i][2][col] for i in range(3)]
        private = None
        for i, a in enumerate(alleles):
            if alleles.count(a) == 1:
                private = i
                break
        sun_by_col.append(private)

    # Draw copy rows
    row_y = [4.55, 3.15, 1.75]
    for row, (name, color, alleles) in enumerate(copies):
        y = row_y[row]
        # copy label (right-aligned, two lines)
        ax.text(label_x, y, name, ha="right", va="center", fontsize=10,
                color=color, fontweight="bold", linespacing=1.1)
        # chromosome/locus bar
        ax.add_patch(Rectangle((2.6, y - 0.12), 9.4, 0.24,
                     fc="#dfe7f2", ec=CN, lw=1.0, zorder=1))
        # allele boxes
        for col, (x, a) in enumerate(zip(x_pos, alleles)):
            is_sun = sun_by_col[col] == row
            box_fc = color if is_sun else "#eef2f7"
            box_ec = color if is_sun else CN
            text_col = "white" if is_sun else CN
            lw = 2.5 if is_sun else 1.5
            ax.add_patch(FancyBboxPatch((x - 0.36, y - 0.34), 0.72, 0.68,
                         boxstyle="round,pad=0.03,rounding_size=0.08",
                         fc=box_fc, ec=box_ec, lw=lw, zorder=3))
            ax.text(x, y, a, ha="center", va="center", fontsize=16,
                    color=text_col, fontweight="bold", zorder=4)
            if is_sun:
                ax.text(x, y + 0.55, "SUN", ha="center", va="center",
                        fontsize=8, color=color, fontweight="bold")

    # Column labels
    for x, lbl in zip(x_pos, labels):
        ax.text(x, 5.75, lbl, ha="center", va="center",
                fontsize=11, color=CN, fontweight="bold")

    # Read row
    read_y = 0.45
    read_alleles = ["A", "T", "C", "G"]
    ax.text(label_x, read_y, "long\nread", ha="right", va="center",
            fontsize=10, color=CG, fontweight="bold", linespacing=1.1)
    ax.add_patch(FancyBboxPatch((2.6, read_y - 0.22), 9.4, 0.44,
                 boxstyle="round,pad=0.03,rounding_size=0.08",
                 fc="#f4f4f4", ec=CN, lw=1.5, zorder=1))
    for col, (x, a) in enumerate(zip(x_pos, read_alleles)):
        ax.add_patch(FancyBboxPatch((x - 0.32, read_y - 0.28), 0.64, 0.56,
                     boxstyle="round,pad=0.03,rounding_size=0.08",
                     fc="#5d6d7e", ec=CN, lw=1.5, zorder=2))
        ax.text(x, read_y, a, ha="center", va="center", fontsize=15,
                color="white", fontweight="bold", zorder=3)

    # Assignment arrow and label on the right side, clear of all boxes
    arrow_x = 11.45
    ax.annotate("", xy=(arrow_x, read_y + 0.55), xytext=(arrow_x, 3.15),
                arrowprops=dict(arrowstyle="->", color=CO, lw=2.5,
                                connectionstyle="arc3,rad=0.15"))
    ax.text(arrow_x + 0.25, 1.8, "matches\ncopy 2", ha="left", va="center",
            fontsize=10, color=CO, fontweight="bold", linespacing=1.1)

    # Outcome box below the read, no overlap
    ax.add_patch(FancyBboxPatch((1.4, -0.95), 10.2, 0.85,
                 boxstyle="round,pad=0.08,rounding_size=0.10",
                 fc="#fff8f3", ec=CO, lw=2.0, zorder=1))
    ax.text(6.5, -0.45,
            "A read that observes copy 2's private SUN alleles (PSV 2 and PSV 4) "
            "is assigned deterministically to copy 2.",
            ha="center", va="center", fontsize=11, color=CN, fontweight="bold")
    ax.text(6.5, -0.75,
            "No other copy carries that allele pattern, so the per-read gate returns |N(r)| = 1 — "
            "one copy, no guess.",
            ha="center", va="center", fontsize=10, color=CG, linespacing=1.3)

    ax.set_title("Long-read copy assignment: SUNs make assignment deterministic",
                 fontsize=14, fontweight="bold", color=CN, pad=12)

    # Legend in the upper-right corner, clear of the PSV grid
    leg_x = 11.35
    leg_y = 5.35
    ax.add_patch(FancyBboxPatch((leg_x, leg_y), 0.40, 0.32,
                 boxstyle="round,pad=0.02,rounding_size=0.05", fc=CT, ec=CT))
    ax.text(leg_x + 0.55, leg_y + 0.16, "private allele", ha="left", va="center",
            fontsize=9, color=CG)
    ax.add_patch(FancyBboxPatch((leg_x, leg_y - 0.45), 0.40, 0.32,
                 boxstyle="round,pad=0.02,rounding_size=0.05", fc="#eef2f7", ec=CN))
    ax.text(leg_x + 0.55, leg_y - 0.29, "shared allele", ha="left", va="center",
            fontsize=9, color=CG)

    fig.savefig(FIG_PHASE, dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- figure 8: silent copies still matter for copy-number quantification
def make_silent_copy():
    fig, ax = plt.subplots(figsize=(13.0, 6.8))
    ax.set_xlim(0, 13.0)
    ax.set_ylim(-1.0, 6.8)
    ax.axis("off")

    # Three-copy toy family: two expressed, one silent but similar.
    copies = [
        {"name": "copy 1\n(expressed)", "y": 5.0, "color": CT,
         "reads": True, "n_reads": 6, "label_color": CT},
        {"name": "copy 2\n(expressed)", "y": 3.5, "color": CT,
         "reads": True, "n_reads": 6, "label_color": CT},
        {"name": "copy 3\n(silent / not expressed)", "y": 2.0, "color": "#95a5a6",
         "reads": False, "n_reads": 0, "label_color": "#7f8c8d"},
    ]

    # Copy tracks
    strip_x = 2.2
    strip_w = 6.0
    for cp in copies:
        y = cp["y"]
        # label
        ax.text(0.35, y, cp["name"], ha="left", va="center", fontsize=10,
                color=cp["label_color"], fontweight="bold", linespacing=1.1)
        # chromosome strip
        ax.add_patch(Rectangle((strip_x, y - 0.12), strip_w, 0.24,
                     fc="#dfe7f2", ec=CN, lw=1.0, zorder=1))
        # gene body
        ax.add_patch(FancyBboxPatch((strip_x + 0.2, y - 0.22), strip_w - 0.4, 0.44,
                     boxstyle="round,pad=0.02,rounding_size=0.08",
                     fc=cp["color"], ec=CN, lw=1.5, alpha=0.85, zorder=2))
        # read pile
        if cp["reads"]:
            for i in range(cp["n_reads"]):
                rx = strip_x + 0.4 + i * 0.88
                ax.add_patch(Rectangle((rx, y + 0.32), 0.72, 0.16,
                             fc="#5dade2", ec=CN, lw=0.5, zorder=3))
            ax.text(strip_x + 0.4 + cp["n_reads"] * 0.88 + 0.15, y + 0.40, "unique reads",
                    ha="left", va="center", fontsize=8, color=CG)
        else:
            # silent copy gets spillover reads instead
            for i in range(4):
                rx = strip_x + 0.4 + i * 0.88
                ax.add_patch(Rectangle((rx, y + 0.32), 0.72, 0.16,
                             fc="#e74c3c", ec="#c0392b", lw=0.5, zorder=3, alpha=0.8))
            ax.text(strip_x + 0.4 + 4 * 0.88 + 0.15, y + 0.40, "spillover reads",
                    ha="left", va="center", fontsize=8, color="#c0392b")

    # Spillover arrows from expressed copies to silent copy
    ax.annotate("", xy=(strip_x + 3.3, 2.48), xytext=(strip_x + 2.4, 4.65),
                arrowprops=dict(arrowstyle="->", color="#c0392b", lw=1.8,
                                connectionstyle="arc3,rad=0.25", ls="--"),
                zorder=4)
    ax.annotate("", xy=(strip_x + 4.0, 2.48), xytext=(strip_x + 4.0, 3.15),
                arrowprops=dict(arrowstyle="->", color="#c0392b", lw=1.8,
                                connectionstyle="arc3,rad=-0.25", ls="--"),
                zorder=4)
    # annotation in clear empty space on the left
    ax.text(1.0, 4.0, "similar copies\nshare reads by\nmultimapping",
            ha="left", va="center", fontsize=9, color="#c0392b",
            fontweight="bold", linespacing=1.1)

    # Metric cards on the right
    card_x = 9.2
    cards = [
        ("n_loci = 2", "RNA definition sees\nonly expressed copies", CT),
        ("χ(H) = 2", "distinct haplotypes\namong observed copies", CP),
        ("depth_cn ≈ 3", "read depth estimates\nall copies, incl. silent", CO),
    ]
    for i, (val, desc, col) in enumerate(cards):
        y = 5.2 - i * 1.45
        ax.add_patch(FancyBboxPatch((card_x, y - 0.50), 3.2, 1.0,
                     boxstyle="round,pad=0.05,rounding_size=0.10",
                     fc="#f6f8fb", ec=col, lw=2.0, zorder=2))
        ax.text(card_x + 1.6, y + 0.20, val, ha="center", va="center",
                fontsize=12, color=col, fontweight="bold")
        ax.text(card_x + 1.6, y - 0.22, desc, ha="center", va="center",
                fontsize=8, color=CG, linespacing=1.15)

    # Depth intuition bar
    bar_y = 0.45
    ax.text(2.8, bar_y + 0.55, "shared exon read depth",
            ha="center", va="center", fontsize=10, color=CN, fontweight="bold")
    # expected for 2 copies
    ax.add_patch(Rectangle((1.5, bar_y - 0.18), 3.0, 0.28, fc="#d5dbdb", ec=CN, lw=1.0))
    ax.text(3.0, bar_y - 0.32, "expected for 2 copies", ha="center", va="center",
            fontsize=8, color=CG)
    # observed (higher)
    ax.add_patch(Rectangle((1.5, bar_y - 0.18), 5.0, 0.28, fc="#f5b041", ec=CO, lw=1.5, alpha=0.7))
    ax.text(5.6, bar_y - 0.04, "observed depth", ha="center", va="center",
            fontsize=8, color="white", fontweight="bold")
    # marker for the extra portion
    ax.plot([4.5, 4.5], [bar_y + 0.12, bar_y + 0.35], color=CO, lw=1.5)
    ax.text(5.8, bar_y + 0.45, "extra depth from\nsilent copy spillover",
            ha="left", va="center", fontsize=8, color=CO, linespacing=1.1)

    # Bottom explanation box
    ax.add_patch(FancyBboxPatch((0.9, -0.95), 11.2, 0.75,
                 boxstyle="round,pad=0.08,rounding_size=0.10",
                 fc="#fff8f3", ec=CO, lw=2.0, zorder=1))
    ax.text(6.5, -0.55,
            "A silent paralog receives multimapping reads from its expressed relatives. "
            "Those spillover reads raise the family's shared exon depth, so depth_cn can estimate "
            "a larger copy number than the RNA definition alone.",
            ha="center", va="center", fontsize=10, color=CG, linespacing=1.3)
    ax.text(6.5, -0.82,
            "RNA keeps high precision; silent copies are a completeness (recall) axis that depth + DNA "
            "help close — DNA alone sees them but is far less specific.",
            ha="center", va="center", fontsize=9, color=CG, linespacing=1.3)

    ax.set_title("Silent copies still matter: read depth reveals hidden multiplicity",
                 fontsize=14, fontweight="bold", color=CN, pad=12)

    fig.savefig(FIG_SILENT, dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- pptx assembly
def add_figure_slide(prs, title, img, caption):
    from PIL import Image
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.8))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = NAVY

    w, h = Image.open(img).size
    scale = min(Inches(12.4) / w, Inches(5.5) / h)
    iw, ih = int(w * scale), int(h * scale)
    s.shapes.add_picture(img, int((Inches(13.33) - iw) / 2), Inches(1.05), width=iw, height=ih)

    cb = s.shapes.add_textbox(Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.6))
    cp = cb.text_frame.paragraphs[0]
    cp.text = caption
    cp.font.size = Pt(12)
    cp.font.color.rgb = GREY
    cp.alignment = PP_ALIGN.CENTER
    cb.text_frame.word_wrap = True
    return s


def build():
    make_definition()
    make_birdseye()
    make_isoforms()
    make_recombination()
    make_igv()
    make_validation()
    make_phasing()
    make_silent_copy()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_figure_slide(prs,
                     "What is a multi-copy gene family?",
                     FIG_DEF,
                     "A family is ≥2 distinct genomic loci carrying homologous, expressed copies of one gene. "
                     "Copies are different loci (often on different chromosomes); isoforms are different transcripts from the same locus.")

    add_figure_slide(prs,
                     "Birdseye view: a real family as a clique",
                     FIG_BIRD,
                     "Family 69 (RABL2) has five copies on five gorilla chromosomes. The homology graph is a perfect clique: "
                     "every locus is linked to every other locus, so the group is one cohesive family.")

    add_figure_slide(prs,
                     "Zoom: two copies, each with multiple isoforms",
                     FIG_ISO,
                     "Two RABL2 family copies (LOC101142457 on NC_073247.2 and RABL2B on NC_086018.1) each produce two isoforms "
                     "as paths through their own variation graphs. A homology edge links the two loci; isoforms are collapsed to one "
                     "locus before families are called.")

    add_figure_slide(prs,
                     "Why recombinant reads don't merge copies",
                     FIG_RECOMB,
                     "Copy 1 (chromosome A) and copy 2 (chromosome B) each have E1 and E2. A spurious read that splices copy-1-E2 "
                     "into copy-2-E1 forms a non-collinear, unsupported junction; it is rejected, so the two loci stay separate.")

    add_figure_slide(prs,
                     "IGV-style read support for two family copies",
                     FIG_IGV,
                     "Representative RNA-seq reads from GGO.bam over LOC101142457 (NC_073247.2) and RABL2B (NC_086018.1). "
                     "Exons are teal blocks, reads are horizontal bars, and grey lines are splice junctions.")

    add_figure_slide(prs,
                     "Strict precision vs recall on canonical known families",
                     FIG_VALID,
                     "All eight literature-recognized multi-copy families from the gorilla catalog. "
                     "Point size reflects copy number; teal points pass both thresholds.")

    add_figure_slide(prs,
                     "Long-read copy assignment via SUNs",
                     FIG_PHASE,
                     "Each copy carries private Singly Unique Nucleotides (SUNs). A long read spanning "
                     "a SUN observes a copy-specific allele and is assigned deterministically to that copy, "
                     "even when the read maps equally well to several copies overall.")

    add_figure_slide(prs,
                     "Silent copies still matter for copy-number quantification",
                     FIG_SILENT,
                     "A silent paralog is invisible to the RNA family definition, but it still receives "
                     "multimapping spillover reads from expressed relatives. That extra depth raises the "
                     "depth_cn estimate above χ(H) and n_loci, helping recover the true genomic copy number.")

    try:
        prs.save(OUT)
        print(f"[+] wrote {OUT} ({len(prs.slides._sldIdLst)} slides)")
    except PermissionError:
        fallback = OUT.replace(".pptx", "_new.pptx")
        prs.save(fallback)
        print(f"[!] {OUT} is locked (probably open in PowerPoint); wrote {fallback} instead")


if __name__ == "__main__":
    build()
