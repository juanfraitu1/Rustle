#!/usr/bin/env python3
"""
make_definition_slides.py — advisor deck for the multi-copy gene family DEFINITION,
traced end-to-end from reads (no unexplained jumps): every transition reads -> ... ->
families states the operation that produces the next object.

Numbers from bench/family_definition_note.md / family_def_vg_reinforce.json.
Run: /home/juanfra/miniforge3/bin/python bench/make_definition_slides.py
Output: bench/family_definition_v2.pptx + bench/family_def_flow.png
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
FLOW = os.path.join(HERE, "family_def_flow.png")
OUT = os.path.join(HERE, "family_definition_v2.pptx")

NAVY = RGBColor(0x1F, 0x2D, 0x5A)
TEAL = RGBColor(0x12, 0x7A, 0x6E)
GREY = RGBColor(0x44, 0x44, 0x44)
ORANGE = RGBColor(0xD9, 0x6B, 0x27)
mNAVY, mTEAL, mGREY, mORANGE = "#1F2D5A", "#127A6E", "#6B7280", "#D96B27"


# ----------------------------------------------------------------- provenance figure
def make_flow_figure():
    fig, ax = plt.subplots(figsize=(13.5, 9.2))
    ax.set_xlim(0, 13.5); ax.set_ylim(0, 16.5); ax.axis("off")

    # (object box text, color)  and the operation that PRODUCES it (on the arrow above; pre-wrapped)
    stages = [
        ("READS\nHiFi IsoSeq full-length transcripts (one molecule = one read)", mGREY, None),
        ("PLACEMENTS\neach read aligns to 1+ loci, each with divergence de (lower = better fit)",
         mGREY, "align with minimap2\n(-ax splice:hq); de = gap-\ncompressed % mismatch"),
        ("DE-NOVO LOCI  V   — the nodes  (★ from the read-coherence assembly)\ncluster reads sharing a splice-junction chain at overlapping coordinates",
         mNAVY, "★ assembly: group reads by\nshared splice chain → one locus;\nattribute each record by best-\noverlap (most co-aligned exonic bp)"),
        ("~R  EDGES  (read-confusability)\ni — j  whenever the reads cannot tell loci i and j apart",
         mTEAL, "per locus pair, count reads on\nBOTH with de_i,de_j ≤ DE_max=0.05\nAND |de_i−de_j| ≤ Δ=0.005;\n≥ k = 3 ties ⇒ edge"),
        ("READ-COUPLED GROUPS  (candidate families)\nmaximal sets of loci reachable through ~R edges; write  i R* j",
         mTEAL, "connected components of ~R\n(maximal edge-linked node sets,\nvia union-find)"),
        ("~B  EDGES  (backbone-homology)\ntwo copies share a common gene backbone",
         mORANGE, "exon-union copy S(v) per locus;\nalign S(i),S(j); edge iff id ≥ 0.80\nAND min(cov_i, cov_j) ≥ τ = 0.30"),
        ("MULTI-COPY GENE FAMILIES\nbackbone-coherent cores; bridge loci (no backbone) drop out",
         mNAVY, "within each R*-group's vertices,\nunion-find over ~B edges;\neach ~B-component (≥ 2) = family"),
    ]
    n = len(stages)
    ys = [15.2 - i * 2.15 for i in range(n)]
    for (txt, color, op), y in zip(stages, ys):
        box = FancyBboxPatch((0.4, y - 0.62), 7.4, 1.24, boxstyle="round,pad=0.08",
                             linewidth=1.6, edgecolor=color, facecolor=color + "14")
        ax.add_patch(box)
        head = txt.split("\n")[0]; sub = txt.split("\n", 1)[1]
        ax.text(4.1, y + 0.18, head, ha="center", va="center", fontsize=10.5, weight="bold", color=color)
        ax.text(4.1, y - 0.30, sub, ha="center", va="center", fontsize=7.6, color="#222")
    for i in range(1, n):
        y0, y1 = ys[i - 1] - 0.64, ys[i] + 0.64
        ax.add_patch(FancyArrowPatch((4.1, y0), (4.1, y1), arrowstyle="-|>",
                     mutation_scale=18, linewidth=1.6, color="#888"))
        ax.text(8.15, (y0 + y1) / 2, stages[i][2], ha="left", va="center", fontsize=6.9,
                style="italic", color="#333")
    ax.text(4.1, 16.2, "From reads to families — each arrow is a named operation  (★ = the one imported step)",
            ha="center", fontsize=11.5, weight="bold", color=mNAVY)
    fig.tight_layout()
    fig.savefig(FLOW, dpi=150, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------- pptx helpers
def add_title_slide(prs, title, subtitle, foot):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bx = s.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(12.1), Inches(2.2))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(33); p.font.bold = True; p.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(19); p2.font.color.rgb = TEAL
    fb = s.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5))
    fp = fb.text_frame.paragraphs[0]; fp.text = foot; fp.font.size = Pt(12); fp.font.color.rgb = GREY
    return s


def add_content_slide(prs, title, bullets, note=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tp = tb.text_frame.paragraphs[0]; tp.text = title
    tp.font.size = Pt(25); tp.font.bold = True; tp.font.color.rgb = NAVY
    body = s.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.6))
    tf = body.text_frame; tf.word_wrap = True
    for i, (txt, lvl, color) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " if lvl == 0 else "      –  ") + txt
        p.font.size = Pt(17 if lvl == 0 else 14)
        p.font.color.rgb = color or RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(6)
    if note:
        nb = s.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.5))
        np_ = nb.text_frame.paragraphs[0]; np_.text = note
        np_.font.size = Pt(11); np_.font.italic = True; np_.font.color.rgb = GREY
    return s


def add_figure_slide(prs, title, img, caption):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
    tp = tb.text_frame.paragraphs[0]; tp.text = title
    tp.font.size = Pt(23); tp.font.bold = True; tp.font.color.rgb = NAVY
    from PIL import Image
    w, h = Image.open(img).size
    maxw, maxh = Inches(12.0), Inches(5.6)
    scale = min(maxw / w, maxh / h)
    iw, ih = int(w * scale), int(h * scale)
    left = int((Inches(13.33) - iw) / 2)
    s.shapes.add_picture(img, left, Inches(1.0), width=iw, height=ih)
    cb = s.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.5))
    cp = cb.text_frame.paragraphs[0]; cp.text = caption
    cp.font.size = Pt(11.5); cp.font.color.rgb = GREY; cp.alignment = PP_ALIGN.CENTER
    return s


def B(t, lvl=0, color=None):
    return (t, lvl, color)


# ----------------------------------------------------------------- build
def build():
    make_flow_figure()
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "What is a multi-copy gene family — built from RNA reads, step by step",
        "Two relations on expressed loci; the family is where both hold",
        "Each object is produced from the one before by a stated operation; the de-novo loci come from the read-coherence assembly (the one imported step).")

    add_content_slide(prs, "The question, and the answer we will build", [
        B("Question: given only RNA long-reads, which loci form a multi-copy gene family?"),
        B("Answer (assembled over the next slides):", 0, NAVY),
        B("a family = a set of expressed loci that are BOTH read-coupled AND backbone-sharing.", 1, NAVY),
        B("'Read-coupled' = the aligner cannot tell the loci apart (relation ~R).", 0, TEAL),
        B("'Backbone-sharing' = the loci's copies align over a common backbone (relation ~B).", 0, ORANGE),
        B("The rest of this deck derives every word of that sentence from the reads.", 0, GREY),
    ])

    add_figure_slide(
        prs, "The whole chain, end to end (we walk each arrow next)",
        FLOW,
        "Each box is an object; each arrow is the operation that produces the next object from it.")

    add_content_slide(prs, "Step 1 — what we start with: the reads", [
        B("Input = HiFi IsoSeq reads: one molecule → one full-length transcript read."),
        B("Each read is a spliced sequence (exons, with the introns removed)."),
        B("No annotation is used — everything downstream is derived from these reads.", 0, TEAL),
        B("(For a multi-copy gene, reads from different copies look almost identical — that is the whole problem.)", 0, GREY),
    ])

    add_content_slide(prs, "Step 2 — align the reads → placements (where each read can sit)", [
        B("Operation: align every read to the genome with minimap2 (-ax splice:hq)."),
        B("A read that fits two near-identical copies gets MULTIPLE placements (a primary + secondaries).", 0),
        B("Each placement carries a number: de = gap-compressed per-base divergence (how well the read fits there).", 0, NAVY),
        B("de ≈ 0 means the read matches that locus almost perfectly.", 1),
        B("So after this step every read is a set of (locus, de) placements — the raw material for 'confusability'.", 0, TEAL),
    ])

    add_content_slide(prs, "Step 3 — turn placements into LOCI (the nodes of our graph)", [
        B("(a) Build the loci [★ the read-coherence assembly — the one step we import]:", 0, NAVY),
        B("cluster reads that share the same splice-junction chain at overlapping coordinates → one DE-NOVO LOCUS.", 1),
        B("⇒ a de-novo locus = one expressed transcription unit = one NODE (vertex V). No gene name is used.", 1, TEAL),
        B("(b) Attribute each alignment RECORD to its single best-overlap locus:", 0, NAVY),
        B("best-overlap = the locus with the most co-aligned exonic bases (ties broken by lower de).", 1),
        B("⇒ a read on two copies contributes two DISTINCT records → genuine multimapping, not double-counting;", 1),
        B("and a read can never be 'confused' with a locus that merely nests inside its own.", 1),
    ], note="★ This is the only object imported from upstream; everything after Step 3 is derived here.")

    add_content_slide(prs, "Step 4 — the first relation: ~R (read-confusability) = the EDGES", [
        B("For each PAIR of loci (i, j), look at the reads that placed on BOTH.", 0),
        B("A read 'ties' i and j when it fits both well AND equally:", 0, NAVY),
        B("de_i, de_j ≤ DE_max = 0.05  (fits both)   AND   |de_i − de_j| ≤ Δ = 0.005  (equally well).", 1, NAVY),
        B("Δ = single-read divergence resolution at HiFi error (≈ 4σ, not tuned);  DE_max = loose copy-vs-other-gene ceiling.", 1),
        B("Draw an edge  i — j  when at least k = 3 reads tie  (a quorum — one stray read can't make an edge).", 0, NAVY),
        B("Worked example — RABL2A vs RABL2B (illustrative de values):", 0, TEAL),
        B("a read fits RABL2A at de≈0.001 and RABL2B at de≈0.0035 → both ≤ DE_max, |diff|=0.0025 ≤ Δ → it ties them.", 1),
        B("47 reads tie (measured) ⇒ we draw the edge RABL2A — RABL2B.", 1, TEAL),
    ])

    add_content_slide(prs, "Step 5 — edges → graph → CONNECTED COMPONENTS (this is the jump, made explicit)", [
        B("We now have a graph: nodes = loci, edges = 'their reads can't be told apart'."),
        B("A CONNECTED COMPONENT = a maximal set of nodes where you can walk from any node to any other along edges.", 0, NAVY),
        B("Computed by union-find: start each node alone; every edge merges the two nodes' sets; the final sets are the components.", 1),
        B("Why this is the right grouping: if reads confuse A with B, and B with C, then A, B, C form one tangle of", 0, TEAL),
        B("mutually-confusable loci — exactly the unit on which read-to-copy assignment is one coupled problem.", 1, TEAL),
        B("⇒ each connected component is a CANDIDATE family. We say two loci in one component are READ-COUPLED: write  i R* j.", 0, NAVY),
    ], note="This is how 'reads' became 'connected components': align → place → count ties → edges → union-find components.")

    add_content_slide(prs, "Step 6 — why candidates are not yet families: ~R over-couples", [
        B("~R alone makes a mistake: a shared repeat or a retrocopy makes reads of UNRELATED loci cross-map."),
        B("Measured example: OCLN and SEPTIN7 — 3,369 reads tie them, so ~R draws an edge — yet they are not copies.", 0, ORANGE),
        B("the reads are confused via a shared element, not because the genes are the same.", 1),
        B("So a connected component can fuse a real family with a bridged intruder. We need a second, independent test.", 0, NAVY),
    ], note="OCLN~SEPTIN7 (3,369 ties; backbone 0.05 on the next slides) is from the genome-wide read-conflict scan — see family_definition_formal.md.")

    add_content_slide(prs, "Step 7 — build the COPIES: aggregate isoforms into one exon-union per locus", [
        B("A copy is a locus (one node); its reads are different ISOFORMS (alternative transcripts) of that copy."),
        B("Operation: union the exonic intervals of ALL reads at a locus → the locus's exon-union sequence.", 0, NAVY),
        B("'the full gene minus introns' (a retained intron shows up where reads span it).", 1),
        B("⇒ one copy model S(v) per locus. Many isoforms collapse to ONE copy — splicing can't masquerade as extra copies.", 0, TEAL),
        B("(verified: 100s–1000s of distinct isoforms per locus → a single exon-union.)", 1, GREY),
    ])

    add_content_slide(prs, "Step 8 — the second relation: ~B (backbone-homology) = the validation EDGES", [
        B("Operation: align the two copy models S(i), S(j) to each other (minimap2)."),
        B("Measure coverage in EACH direction: cov_i = fraction of S(i)'s bases that align, cov_j likewise.", 0, NAVY),
        B("Draw a backbone edge  i ~B j  iff  identity ≥ 0.80  AND  min(cov_i, cov_j) ≥ τ = 0.30.", 0, NAVY),
        B("the MIN (reciprocal) is the point: a real copy aligns over most of BOTH; a repeat-bridge over only a short insert of one.", 1),
        B("Worked example (measured):", 0, TEAL),
        B("RABL2A ~ RABL2B copies align over 0.94 of each (min-cov 0.94) → backbone edge ✓.", 1, TEAL),
        B("OCLN ~ SEPTIN7 align over only 0.05 of one and ~0 of the other → NO backbone edge ✗.", 1, ORANGE),
    ])

    add_content_slide(prs, "Step 9 — the family: the ~B-connected core inside each read-coupled group", [
        B("Operation: take the vertices of ONE ~R-component; on those vertices only, run union-find over ~B edges;", 0, NAVY),
        B("each resulting ~B-component of size ≥ 2 is a family. A locus with no ~B edge to a family-mate drops out.", 1),
        B("DEFINITION:  a multi-copy gene family = a ~B-connected component (size ≥ 2) inside a read-coupled (R*) group.", 0, NAVY),
        B("Equivalently: connected components of  ~B ∩ R*   (R* = 'reads couple them' from Step 5; ~B = 'they share a backbone').", 1),
        B("RABL2A–RABL2B: read-coupled ✓ and backbone-connected ✓ → a FAMILY.", 0, TEAL),
        B("OCLN–SEPTIN7: read-coupled ✓ but NO backbone → the bridge is cut → NOT a family.", 0, ORANGE),
    ])

    add_content_slide(prs, "The definition, fully assembled (every term now derived)", [
        B("A multi-copy gene family is a set of de-novo loci that is:", 0, NAVY),
        B("(1) read-coupled — its loci are mutually confusable: a connected component of ~R", 0, TEAL),
        B("   (≥ 3 reads tie each link at |de_i−de_j| ≤ Δ), AND", 1, TEAL),
        B("(2) backbone-connected — its copies (all-isoform exon-unions) share a common backbone: ~B", 0, ORANGE),
        B("   (identity ≥ 0.80, reciprocal coverage ≥ τ).", 1, ORANGE),
        B("~R alone over-couples (repeat bridges); ~B alone is plain homology (ignores expression/confusability).", 0, GREY),
        B("The family is where BOTH hold — exactly the unit where copy-assignment is both needed and well-posed.", 0, NAVY),
    ])

    add_content_slide(prs, "What this correctly handles — and why (by the two relations)", [
        B("Domain-sharers (adjacent genes sharing one exon) → 0 conflicting reads (best-overlap) → not ~R → excluded.", 0),
        B("Repeat / retro bridges between non-copies → no shared backbone → not ~B → excluded (OCLN~SEPTIN7).", 0),
        B("Isoforms of one copy → collapse into one exon-union → never counted as extra copies.", 0),
        B("Single pairs (size-2, the largest class — 73/196 here) → admitted: one ~B edge in a coupled group, no density needed.", 0),
        B("Read-resolvable diverged paralogs → fewer than 3 reads tie → not ~R → out of scope by design.", 0),
    ])

    add_content_slide(prs, "Evidence", [
        B("Hand-labelled panel (17 IsoSeq candidates):  TP = 7, TN = 10, FP = 0, FN = 0.", 0, NAVY),
        B("~B is the decisive lever: per candidate, repeat/retro bridges score ≤ 0.1 vs real copies ≥ 0.8.", 0, TEAL),
        B("Genome-wide, OOM-safe (de-novo loci): 212 candidate components → 196 validated families.", 0),
        B("cutting 14 backbone-less member loci dissolves/refines 16 components; cross-chrom bridges 20 → 12; size-2 80 → 73.", 1),
        B("Honest cost: of the cut edges, 9 had no DNA homology vs 12 DNA-paralog/sub-bar → conservative (cuts confident bridges).", 1, GREY),
        B("Six alternative refinements (clustering, read coverage/intron, junction concordance) tested and rejected;", 0, GREY),
        B("only ~B (backbone) and external repeat-masking survive — ~B is the one that closes the repeat-bridge gap.", 1, GREY),
    ])

    add_content_slide(prs, "Summary — one sentence, fully grounded", [
        B("A multi-copy gene family = de-novo loci whose reads cannot be told apart (~R)", 0, NAVY),
        B("and whose all-isoform copies share a common backbone (~B).", 0, NAVY),
        B("Built only from reads: align → attribute per record → loci → count ties → ~R edges → components", 0, TEAL),
        B("→ exon-union copies → align → ~B edges → backbone-connected core = family.", 0, TEAL),
        B("Scope (honest): expressed + confusable copies only; copy NUMBER (PSV haplotypes) is the next step, deferred.", 0, GREY),
    ])

    prs.save(OUT)
    print(f"[+] wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")
    print(f"[+] wrote {FLOW}")


if __name__ == "__main__":
    build()
