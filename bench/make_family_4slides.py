#!/usr/bin/env python
"""A focused family-definition deck (5 slides): definition, BAM->families pipeline, single-core limit,
methods table, runaway-read defense.

Output: bench/slides/three_slides.pptx (+ PNGs). Run: /home/juanfra/miniforge3/bin/python bench/make_3slide_deck.py
"""
import os
import textwrap

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import networkx as nx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = os.path.join(os.path.dirname(__file__), "slides")
os.makedirs(OUT, exist_ok=True)
plt.rcParams.update({"font.size": 15, "font.family": "DejaVu Sans", "axes.spines.top": False,
                     "axes.spines.right": False, "figure.dpi": 150})
BLUE, RED, GREEN, GREY, NAVY = "#377eb8", "#e41a1c", "#4daf4a", "#999999", (20, 40, 90)


def savefig(name):
    p = os.path.join(OUT, name)
    plt.savefig(p, bbox_inches="tight", facecolor="white"); plt.close(); return p


def fig_pipeline():
    from matplotlib.patches import FancyBboxPatch
    fig, ax = plt.subplots(figsize=(12.6, 5.2)); ax.axis("off")
    NH = "#14285a"
    boxes = [
        (0.2, "BAM alignments", "each read: intron chain\n(CIGAR) · de/NM · MAPQ\nprimary + secondary"),
        (3.45, "de-novo loci", "group reads by intron\nchain → gate → collapse\nisoforms into gene loci"),
        (6.7, "read-conflict graph", "loci = nodes\na de-tie = an edge\n(a read fits both loci)"),
        (9.95, "multi-copy\nFAMILIES", "connected\ncomponents\n(|C| ≥ 2)"),
    ]
    bw, bh, by = 2.6, 1.7, 2.5
    for (x, title, sub) in boxes:
        ax.add_patch(FancyBboxPatch((x, by), bw, bh, boxstyle="round,pad=0.04,rounding_size=0.18",
                                    fc="#eef3fb", ec=NH, lw=2.2))
        ax.text(x + bw / 2, by + bh - 0.34, title, ha="center", va="center", fontsize=13.5, weight="bold", color=NH)
        ax.text(x + bw / 2, by + 0.55, sub, ha="center", va="center", fontsize=10.3, color="#282828")
    # mini read-stack on box 1
    for k, yy in enumerate([4.55, 4.78, 5.01]):
        ax.add_patch(plt.Rectangle((0.45, yy), 0.7, 0.07, fc=GREY, ec="black"))
        ax.add_patch(plt.Rectangle((1.4, yy), 0.7, 0.07, fc=GREY, ec="black"))
    ax.plot([1.15, 1.4], [4.62, 4.62], color=GREY, lw=0.7)  # splice gap hint
    # mini conflict graph on box 4 (2 components)
    gx, gy = 11.25, 4.9
    ax.plot([gx, gx + .5], [gy, gy + .25], color=BLUE, lw=2); ax.plot([gx, gx + .5], [gy, gy - .25], color=BLUE, lw=2)
    ax.plot([gx + .5, gx + .5], [gy + .25, gy - .25], color=BLUE, lw=2)
    for px, py in [(gx, gy), (gx + .5, gy + .25), (gx + .5, gy - .25)]:
        ax.add_patch(plt.Circle((px, py), .08, fc=BLUE, ec="black", zorder=3))
    # arrows + numbered mechanism labels
    arrow_x = [(2.8, 3.45), (6.05, 6.7), (9.3, 9.95)]
    for (x0, x1) in arrow_x:
        ax.annotate("", (x1, by + bh / 2), (x0, by + bh / 2), arrowprops=dict(arrowstyle="-|>", color=NH, lw=2.5))
    for i, (x0, x1) in enumerate(arrow_x):
        ax.text((x0 + x1) / 2, by + bh / 2 + 0.25, f"{['①','②','③'][i]}", ha="center", fontsize=15, weight="bold", color=RED)
    # legend (the mechanism of each arrow, code functions named) — one auto-sizing bbox that WRAPS the text
    items = [
        "①  group primary reads by EXACT intron chain → de-novo skeletons; keep ≥3 reads + all-canonical GT–AG "
        "junctions; collapse isoforms sharing a junction into one gene LOCUS.   (pass1_skeletons → assemble_gate → collapse_loci)",
        "②  for each read, list its placements across the loci; draw an EDGE between two loci when a read DE-TIES — "
        "|Δde| ≤ δ AND both alignments fit (≤ de_max), ≥3 such reads.   (conflict_edges)",
        "③  a family = a connected component of that graph (no global core, no similarity bar).   (conflict_families)",
    ]
    legend = "\n".join(textwrap.fill(it, 112, subsequent_indent="      ") for it in items)
    ax.text(0.2, -0.05, legend, fontsize=10.3, va="top", ha="left", color="#141414",
            bbox=dict(boxstyle="round,pad=0.6", fc="#f6f6f6", ec=GREY))
    ax.set_xlim(-0.1, 12.9); ax.set_ylim(-1.6, 5.25)
    return savefig("s_pipeline.png")


def fig_conflict():
    fig, ax = plt.subplots(figsize=(9.5, 4.4))
    G = nx.Graph()
    fam1, fam2, single, dom = ["A0", "A1", "A2"], ["B0", "B1"], ["S"], ["D0", "D1"]
    G.add_nodes_from(fam1 + fam2 + single + dom)
    G.add_edges_from([("A0", "A1"), ("A1", "A2"), ("A0", "A2"), ("B0", "B1")])
    pos = {"A0": (0, 1), "A1": (1, 1.6), "A2": (1, .4), "B0": (3, 1.3), "B1": (4, .8),
           "S": (5.7, 1.1), "D0": (3.3, -.7), "D1": (4.3, -.7)}
    nx.draw_networkx_nodes(G, pos, nodelist=fam1, node_color=BLUE, node_size=1500, ax=ax)
    nx.draw_networkx_nodes(G, pos, nodelist=fam2, node_color=GREEN, node_size=1500, ax=ax)
    nx.draw_networkx_nodes(G, pos, nodelist=single + dom, node_color=GREY, node_size=1500, ax=ax)
    nx.draw_networkx_edges(G, pos, width=3, ax=ax)
    nx.draw_networkx_labels(G, pos, font_color="white", font_weight="bold", font_size=12, ax=ax)
    ax.text(0.5, 2.15, "family 1\n(tandem array)", ha="center", color=BLUE, fontsize=12, weight="bold")
    ax.text(3.5, 1.85, "family 2\n(cross-chrom pair)", ha="center", color=GREEN, fontsize=12, weight="bold")
    ax.text(5.7, 1.7, "single-copy\n(no tie → not a family)", ha="center", fontsize=11)
    ax.text(3.8, -1.25, "domain-sharers\n(share 1 exon, NO tie → not merged)", ha="center", fontsize=11)
    ax.text(2.0, -2.05, "edge = a read fits BOTH loci comparably (a 'tie')   •   family = connected component",
            ha="center", fontsize=12.5, style="italic")
    ax.axis("off"); ax.set_ylim(-2.35, 2.6); ax.set_xlim(-0.8, 6.6)
    return savefig("s_conflict.png")


def fig_singlecore():
    fig, ax = plt.subplots(figsize=(11, 5.4)); ax.axis("off")
    # 4 copies as sequence bars; consecutive copies share a homology block in SHIFTING positions.
    GR, BL, OR = "#4daf4a", "#377eb8", "#ff7f00"
    bars = {  # copy -> list of (x0, x1, colour, label)
        "A": [(1.0, 4.0, GR, ""), (4.0, 8.0, GREY, "")],
        "B": [(1.0, 4.0, GR, ""), (4.0, 6.0, BL, ""), (6.0, 8.0, GREY, "")],
        "C": [(1.0, 4.0, GREY, ""), (4.0, 6.0, BL, ""), (6.0, 8.0, OR, "")],
        "D": [(1.0, 6.0, GREY, ""), (6.0, 8.0, OR, "")],
    }
    ys = {"A": 4.0, "B": 3.2, "C": 2.4, "D": 1.6}
    for cp, segs in bars.items():
        y = ys[cp]
        for (x0, x1, c, _) in segs:
            ax.add_patch(plt.Rectangle((x0, y), x1 - x0, 0.5, fc=c, ec="black"))
        ax.text(0.8, y + 0.25, f"copy {cp}", ha="right", va="center", fontsize=12, weight="bold")
    ax.text(2.5, 4.78, "homology A↔B", ha="center", color=GR, fontsize=11, weight="bold")
    ax.text(5.0, 1.35, "homology C↔D", ha="center", color=OR, fontsize=11, weight="bold")
    ax.text(5.0, 3.0, "B↔C", ha="center", color=BL, fontsize=10.5, weight="bold")
    ax.text(4.0, 5.15, "no block spans all 4 copies  →  common core = ∅", ha="center", color=RED, fontsize=11.5, weight="bold")

    # INSET: the family as a union-find CHAIN — consecutive ties solid, A⋯D NOT directly tied (transitive closure).
    gx = 9.6
    ny = {n: ys[n] + 0.25 for n in "ABCD"}
    for (u, v), c in {("A", "B"): GR, ("B", "C"): BL, ("C", "D"): OR}.items():
        ax.plot([gx, gx], [ny[u], ny[v]], color=c, lw=3.5, zorder=1)
        ax.text(gx + 0.2, (ny[u] + ny[v]) / 2, "tie", fontsize=8.5, color=c, va="center")
    for n in "ABCD":  # dotted leader from each bar to its node
        ax.plot([8.05, gx - 0.17], [ny[n], ny[n]], color=GREY, lw=0.8, ls=":", zorder=0)
        ax.add_patch(plt.Circle((gx, ny[n]), 0.17, fc="#14285a", ec="black", zorder=2))
        ax.text(gx, ny[n], n, ha="center", va="center", color="white", fontsize=10, weight="bold", zorder=3)
    # A⋯D dashed NON-edge, arcing out to the right
    ax.annotate("", (gx, ny["D"]), (gx, ny["A"]), zorder=1,
                arrowprops=dict(arrowstyle="-", color=GREY, ls="--", lw=1.6, connectionstyle="arc3,rad=-0.55"))
    ax.text(gx + 1.55, (ny["A"] + ny["D"]) / 2, "A⋯D:\nNO direct tie", ha="left", va="center", fontsize=9.5, color=GREY)
    ax.text(gx, ny["D"] - 0.62, "1 connected component\n= ONE family", ha="center", va="top",
            fontsize=10, color="#146e14", weight="bold")
    # two interpretations — each an AUTO-SIZING bbox that wraps its own text (red = single-core, green = conflict)
    sc = "Single core (longest common substring / highest-Jaccard block, shared by ALL):\n" + textwrap.fill(
        "the all-copies intersection SHRINKS as members diverge → high θ drops the divergent tail (UNDER-merge); "
        "low θ rewards a shared domain/repeat present in unrelated genes (OVER-merge). Either way it fails.", 90)
    ax.text(0.5, 0.85, sc, fontsize=11, va="top", ha="left", color="#7a1010",
            bbox=dict(boxstyle="round,pad=0.4", fc="#fbeeee", ec="#b01515"))
    rc = "Read-conflict / pairwise POA-core → connected component (union-find):\n" + textwrap.fill(
        "membership = a PAIRWISE tie to ANY existing member → the chain A–B–C–D is ONE family, even though A and D "
        "share NO common core and are not directly similar. Transitivity, not a global block.", 90)
    ax.text(0.5, -1.0, rc, fontsize=11, va="top", ha="left", color="#0d5a0d",
            bbox=dict(boxstyle="round,pad=0.4", fc="#eef7ee", ec="#146e14"))
    ax.set_xlim(-1.0, 12.6); ax.set_ylim(-2.1, 5.45)
    return savefig("s_singlecore.png")


def fig_runaway():
    fig, ax = plt.subplots(figsize=(11, 5.2)); ax.axis("off")
    # reference transcript: 3 exons (boxes) joined by introns (lines). terminal = exon 3.
    exons = [(1.0, 2.2), (3.4, 4.4), (5.6, 6.6)]  # (x0,x1) of each exon
    ey = 4.0
    for i, (x0, x1) in enumerate(exons):
        ax.add_patch(plt.Rectangle((x0, ey), x1 - x0, 0.45, fc=BLUE, ec="black"))
        ax.text((x0 + x1) / 2, ey + 0.22, f"e{i+1}", ha="center", va="center", color="white", weight="bold")
    for (a, b) in [(exons[0][1], exons[1][0]), (exons[1][1], exons[2][0])]:
        ax.plot([a, b], [ey + 0.22, ey + 0.22], color=BLUE, lw=1.5)
    ax.text(3.8, ey + 0.75, "real transcript:  intron chain = (e1|e2, e2|e3),  terminal exon = e3", ha="center", fontsize=12, weight="bold")
    # 4 good reads agreeing
    for k, y in enumerate([3.3, 3.0, 2.7, 2.4]):
        ax.add_patch(plt.Rectangle((1.0, y), 5.6, 0.13, fc=GREY, ec="black"))
    ax.text(0.9, 2.85, "4 reads agree", ha="right", fontsize=11, color="black")
    # runaway read: overshoots e3 AND splices into a phantom downstream exon
    yr = 1.7
    ax.add_patch(plt.Rectangle((1.0, yr), 5.6, 0.16, fc=RED, ec="black"))
    ax.plot([6.6, 7.8], [yr + 0.08, yr + 0.08], color=RED, lw=1.5, ls="--")  # phantom intron
    ax.add_patch(plt.Rectangle((7.8, yr), 1.4, 0.16, fc=RED, ec="black"))    # phantom downstream exon
    ax.text(0.9, yr + 0.08, "runaway read", ha="right", fontsize=11, color=RED, weight="bold")
    ax.annotate("read-through / chimera / mis-clip:\novershoots e3 + adds a PHANTOM junction", (8.5, yr + 0.4),
                (6.7, 0.9), fontsize=11, color=RED, ha="left", arrowprops=dict(arrowstyle="->", color=RED))
    # the two defenses
    # two defenses — one auto-sizing bbox that wraps its text
    defenses = ("Two defenses (de-novo assembly, denovo_assemble.rs):\n"
                + textwrap.fill("①  EXACT intron-chain key — the runaway's chain (extra phantom junction) ≠ the real chain "
                                "→ it forms its OWN skeleton; a lone runaway fails the ≥3-read quorum → DROPPED. It never "
                                "extends the real chain.", 92, subsequent_indent="     ")
                + "\n" + textwrap.fill("②  Terminal boundary = the k-th-SUPPORTED end (k=2), not the union (max-end) → a "
                                       "single read's overshoot of e3 is trimmed to the 2nd-furthest position. Verified: "
                                       "pass1_robust_trims_a_runaway_terminal_read.", 92, subsequent_indent="     "))
    ax.text(0.4, 0.35, defenses, fontsize=11, va="top", ha="left", color="#141414",
            bbox=dict(boxstyle="round,pad=0.5", fc="#f4f4f4", ec=GREY))
    ax.set_xlim(-1.2, 10.6); ax.set_ylim(-1.7, 5.0)
    return savefig("s_runaway.png")


METHODS = [
    ("Method", "Signal / level", "Approach", "Limit on RECENT near-identical paralogs"),
    ("OrthoFinder", "protein homology", "DIAMOND all-vs-all + MCL + gene trees", "near-identical copies collapse / over-merge"),
    ("Pfam / InterPro", "domain HMM", "group by shared domain architecture", "groups domain-sharers (FP); domain ≠ copy"),
    ("Ensembl Compara", "phylogenetic", "gene-tree ↔ species-tree reconciliation", "fails on recent near-identical paralogs"),
    ("Mash / minimizer-Jaccard", "alignment-free", "MinHash / minimizer distance", "cannot separate domain-sharers"),
    ("SEDEF / BISER", "DNA self-align", "Jaccard + chaining → segdup pairs", "strong on recent SDs; DNA only, no expression"),
    ("Soto 2025 (HSD)", "DNA + copy-number", "SD98 + shared-exon + famCN grouping", "reference human catalog; DNA, no RNA/isoform"),
    ("Eichler 2024 (TBC1D3)", "long-read RNA", "map FLNC; assign read iff AS ≥ 10", "ONE family; not de-novo / genome-wide"),
    ("longcallR", "long-read RNA", "CNN SNP + MEC phasing + ASE/ASJ", "uniquely-mappable only; never assigns to COPIES"),
    ("RSEM / Salmon / kallisto", "RNA quant (EM)", "EM apportions multireads", "needs a PRE-DEFINED transcript set; no discovery"),
    ("IsoCon", "long-read RNA", "NN cluster/correct + real-vs-error test", "closest prior art; WITHIN a known family, targeted"),
    ("★ OURS (read-conflict)", "long-read RNA, de-novo", "loci a read CAN'T tell apart (conflict graph) + copy assign", "fills the empty niche: de-novo · genome-wide · RNA · copy-level"),
]


def build():
    prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def tb(slide, text, l, t, w, h, sz, bold=False, color=(20, 20, 20), align=PP_ALIGN.LEFT, italic=False):
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = box.text_frame; tf.word_wrap = True
        for i, ln in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ln; p.alignment = align
            r = p.runs[0]; r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = RGBColor(*color)

    def title(slide, s): tb(slide, s, 0.6, 0.28, 12.1, 1.0, 28, bold=True, color=NAVY)

    def fig(slide, png, top=1.5, maxw=11.8, maxh=5.0):
        from PIL import Image
        iw, ih = Image.open(png).size
        sc = min(maxw / (iw / 150), maxh / (ih / 150)); w = (iw / 150) * sc
        slide.shapes.add_picture(png, Inches((13.333 - w) / 2), Inches(top), width=Inches(w))

    # Slide 1
    s = prs.slides.add_slide(blank); title(s, "1 · What is a multi-copy gene family (our definition)?")
    tb(s, "A set of loci the READS cannot tell apart — a family = a connected component of the read-conflict graph "
          "(an edge = some read fits both loci comparably). Relational, annotation-free; no absolute-similarity bar.",
       0.7, 1.15, 11.9, 0.9, 18, italic=True, color=(60, 60, 60))
    fig(s, fig_conflict(), top=2.05, maxh=4.6)

    # Slide 2 — the pipeline: from BAM alignments to families
    s = prs.slides.add_slide(blank); title(s, "2 · From reads (BAM alignments) to families — the pipeline")
    fig(s, fig_pipeline(), top=1.35, maxh=5.0)
    tb(s, "isoform = a distinct splice variant of a gene — a particular intron chain (which exons are joined). "
          "One gene can express several isoforms; step ① collapses isoforms that share a junction into ONE gene LOCUS, "
          "so a family is a set of loci (genes/copies), not of isoforms.",
       0.55, 6.75, 12.25, 0.7, 13, italic=True, color=(90, 90, 90))

    # Slide 3 — why a single common core cannot find all members
    s = prs.slides.add_slide(blank); title(s, "3 · Why a single common core can't discover all members")
    fig(s, fig_singlecore(), top=1.4, maxh=5.6)

    # Slide 4 — methods table
    s = prs.slides.add_slide(blank); title(s, "4 · Methods to detect multi-copy gene families (gathered)")
    rows, cols = len(METHODS), 4
    gt = s.shapes.add_table(rows, cols, Inches(0.35), Inches(1.35), Inches(12.6), Inches(5.7)).table
    gt.columns[0].width = Inches(2.5); gt.columns[1].width = Inches(2.0); gt.columns[2].width = Inches(4.2); gt.columns[3].width = Inches(3.9)
    for ci in range(cols):
        for ri in range(rows):
            cell = gt.cell(ri, ci); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06); cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            p = cell.text_frame.paragraphs[0]; p.text = METHODS[ri][ci]; p.runs[0].font.size = Pt(11.5 if ri else 12.5)
            header = ri == 0; ours = ri == rows - 1
            p.runs[0].font.bold = header or ours or ci == 0
            if header:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(*NAVY); p.runs[0].font.color.rgb = RGBColor(255, 255, 255)
            elif ours:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(225, 240, 225); p.runs[0].font.color.rgb = RGBColor(20, 110, 20)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if ri % 2 else RGBColor(244, 247, 251)
    tb(s, "The RNA / long-read / de-novo / copy-level niche was empty — IsoCon is the closest, but only WITHIN a known family.",
       0.4, 7.05, 12.5, 0.4, 11.5, italic=True, color=(110, 110, 110))

    # Slide 5 — runaway read
    s = prs.slides.add_slide(blank); title(s, "5 · A runaway read can't wrongly extend the intron chain")
    fig(s, fig_runaway(), top=1.35, maxh=5.5)

    path = os.path.join(OUT, "family_slides.pptx"); prs.save(path); return path


if __name__ == "__main__":
    print("wrote", build())
